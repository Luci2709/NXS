import streamlit as st
import pandas as pd
import plotly.express as px
import json
import os
import glob
import requests
import textwrap
# from streamlit_gsheets import GSheetsConnection  <-- ENTFERNT
from supabase import create_client, Client # <-- NEU
import uuid
import base64
import calendar
import io
import numpy as np
import time
from datetime import datetime, date, timedelta
from PIL import Image, ImageDraw, ImageEnhance

try:
    import pptx
    from pptx.util import Inches
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
try:
    from fpdf import FPDF
    HAS_FPDF = True
except ImportError:
    HAS_FPDF = False
try:
    from st_img_pastebutton import paste
    HAS_CLIPBOARD = True
    CLIPBOARD_ERR = None
except ImportError as e:
    HAS_CLIPBOARD = False
    CLIPBOARD_ERR = str(e)

import random
import threading
import re

# ==============================================================================
# 🔧 ROBUST MONKEY PATCH (Fix für Component Error & Streamlit 1.40+)
# ==============================================================================
import streamlit.elements.image as st_image
import io
import base64
from PIL import Image

# --- SUPABASE CONNECTION SETUP ---
# Initialisierung der Verbindung
# --- SUPABASE CONNECTION SETUP ---
@st.cache_resource
def init_supabase():
    # Hier holt sich Streamlit die echten Daten aus deiner secrets.toml
    url = st.secrets["SUPABASE_URL"] 
    key = st.secrets["SUPABASE_KEY"]
    
    return create_client(url, key)

supabase = init_supabase()

def custom_image_to_url(image, width=None, clamp=False, channels="RGB", output_format="JPEG", image_id=None, allow_emoji=False):
    """
    Ersetzt die interne Streamlit-Funktion.
    """
    if not isinstance(image, Image.Image):
        return ""
    
    img_byte_arr = io.BytesIO()
    fmt = output_format.upper() if output_format else "JPEG"
    if fmt == "JPG": fmt = "JPEG"
    
    image.save(img_byte_arr, format=fmt)
    img_byte_arr = img_byte_arr.getvalue()
    
    b64_encoded = base64.b64encode(img_byte_arr).decode()
    mime = f"image/{fmt.lower()}"
    return f"data:{mime};base64,{b64_encoded}"

st_image.image_to_url = custom_image_to_url

# ⚠️ JETZT ERST DIE CANVAS LIBRARY IMPORTIEREN
from streamlit_drawable_canvas import st_canvas

# ==============================================================================
# 🔐 AUTHENTICATION SYSTEM
# ==============================================================================

USER_CREDENTIALS = {
    "visitor1": {"password": "visitor123", "role": "visitor"},
    "visitor2": {"password": "visitor123", "role": "visitor"},
    "Luggi": {"password": "1", "role": "player"},
    "Andrei": {"password": "player123", "role": "player"},
    "Benni": {"password": "Valorant2026", "role": "player"},
    "Sofi": {"password": "player123", "role": "player"},
    "Luca": {"password": "player123", "role": "player"},
    "Remus": {"password": "player123", "role": "player"},
    "coach1": {"password": "coach123", "role": "coach"},
    "coach2": {"password": "coach123", "role": "coach"},
    "testing": {"password": "test123", "role": "testing"},
}

def load_users_db():
    """Load users from Supabase or seed from USER_CREDENTIALS"""
    try:
        response = supabase.table("nexus_users").select("*").execute()
        df = pd.DataFrame(response.data)
        
        # Rename 'id' to 'ID' for internal consistency if necessary, 
        # though users table mostly relies on Username
        if 'id' in df.columns: df.rename(columns={'id': 'ID'}, inplace=True)
        
        if df.empty or 'Username' not in df.columns: raise Exception("Init")
        if 'MustChangePassword' not in df.columns: df['MustChangePassword'] = False
        return df
    except:
        # Seed Data
        data = []
        for u, c in USER_CREDENTIALS.items():
            reset = True if c['role'] == 'player' else False
            # Generate a random ID for seeding
            data.append({'Username': u, 'Password': c['password'], 'Role': c['role'], 'MustChangePassword': reset})
        
        # Initial Save to Supabase (Upsert based on Username ideally, but since we use ID as PK, we just insert)
        # Assuming table is empty or we are initializing
        try: 
            supabase.table("nexus_users").insert(data).execute()
        except: pass
        
        return pd.DataFrame(data)

def save_users_db(df):
    try: 
        # Prepare data: Ensure NaNs are None for JSON
        df_save = df.copy()
        if 'ID' in df_save.columns: df_save.rename(columns={'ID': 'id'}, inplace=True)
        
        records = df_save.where(pd.notnull(df_save), None).to_dict(orient='records')
        
        # Upsert using 'id' if present, otherwise might duplicate if not careful.
        # Assuming Username is unique or ID is handled.
        supabase.table("nexus_users").upsert(records).execute()
    except Exception as e: print(f"DB Error: {e}")

def check_credentials(username, password, df_users):
    if df_users.empty: return None
    user = df_users[df_users['Username'] == username]
    if not user.empty and str(user.iloc[0]['Password']) == str(password):
        return user.iloc[0]
    return None

def get_allowed_pages(role):
    if role == "visitor":
        return ["🏠 DASHBOARD"]
    elif role == "player":
        return ["🏠 DASHBOARD", "👥 COACHING", "⚽ SCRIMS", "🗺️ MAP ANALYZER", "📘 STRATEGY BOARD", "📚 RESOURCES", "📅 CALENDAR", "📊 PLAYERS", "📹 VOD REVIEW"]
    elif role == "coach":
        return ["🏠 DASHBOARD", "👥 COACHING", "⚽ SCRIMS", "📝 MATCH ENTRY", "🗺️ MAP ANALYZER", "📘 STRATEGY BOARD", "📚 RESOURCES", "📅 CALENDAR", "📊 PLAYERS", "📹 VOD REVIEW", "💾 DATABASE"]
    elif role == "testing":
        return ["🏠 DASHBOARD", "👥 COACHING", "⚽ SCRIMS", "📝 MATCH ENTRY", "🗺️ MAP ANALYZER", "📘 STRATEGY BOARD", "📚 RESOURCES", "📅 CALENDAR", "📊 PLAYERS", "📹 VOD REVIEW", "💾 DATABASE"]
    return []

def login_page():
    st.title("🔐 NEXUS LOGIN")
    if 'change_password_user' in st.session_state:
        u = st.session_state.change_password_user
        st.warning(f"⚠️ Security Alert: Please set a new password for **{u}**.")
        with st.form("pwd_reset"):
            p1 = st.text_input("New Password", type="password")
            p2 = st.text_input("Confirm Password", type="password")
            if st.form_submit_button("Update Password"):
                if p1 != p2: st.error("Passwords do not match.")
                elif not p1: st.error("Password cannot be empty.")
                else:
                    df = load_users_db()
                    df.loc[df['Username']==u, 'Password'] = p1
                    df.loc[df['Username']==u, 'MustChangePassword'] = False
                    save_users_db(df)
                    user_row = df[df['Username']==u].iloc[0]
                    st.session_state.authenticated = True
                    st.session_state.username = u
                    st.session_state.role = user_row['Role']
                    st.session_state.allowed_pages = get_allowed_pages(user_row['Role'])
                    del st.session_state.change_password_user
                    st.success("Password updated!"); st.rerun()
        return

    st.markdown("---")
    col1, col2, col3 = st.columns([1, 2, 1])
    with col2:
        st.markdown("### Welcome to NEXUS")
        st.markdown("Please enter your credentials to access the system.")
        with st.form("login_form"):
            username = st.text_input("Username")
            password = st.text_input("Password", type="password")
            submitted = st.form_submit_button("Login", type="primary")

            if submitted:
                if not username or not password:
                    st.error("Please enter both username and password.")
                else:
                    df_users = load_users_db()
                    user = check_credentials(username, password, df_users)
                    if user is not None:
                        if str(user.get('MustChangePassword', 'False')).upper() == 'TRUE':
                            st.session_state.change_password_user = username
                            st.rerun()
                        
                        st.session_state.authenticated = True
                        st.session_state.username = username
                        st.session_state.role = user['Role']
                        st.session_state.allowed_pages = get_allowed_pages(user['Role'])
                        st.success(f"Welcome {username}! Redirecting...")
                        st.rerun()
                    else:
                        st.error("Invalid username or password.")

def logout():
    for key in ['authenticated', 'username', 'role', 'allowed_pages']:
        if key in st.session_state:
            del st.session_state[key]
    st.rerun()

# ==============================================================================
st.set_page_config(page_title="NXS Dashboard", layout="wide", page_icon="💠")

if 'authenticated' not in st.session_state or not st.session_state.authenticated:
    login_page()
    st.stop()

st.markdown("""
<style>
    /* --- GLOBAL THEME --- */
    .stApp { 
        background-color: #050505; 
        background-image: radial-gradient(circle at 50% 0%, #1a1a2e 0%, #050505 60%);
        color: #e0e0e0; 
    }
    [data-testid="stSidebar"] { 
        background-color: #080810;
        border-right: 1px solid #333; 
    }
    
    /* --- TYPOGRAPHY & GRADIENTS --- */
    h1, h2, h3 {
        background: linear-gradient(90deg, #00BFFF 0%, #FF1493 100%);
        -webkit-background-clip: text;
        -webkit-text-fill-color: transparent;
        font-weight: 900 !important;
        text-transform: uppercase;
        font-family: 'Segoe UI', sans-serif;
        letter-spacing: 1px;
    }
    
    /* --- NEON CARDS & BOXES --- */
    div.stContainer {
        border-radius: 12px;
    }

    /* Stat Box (Map Analyzer) */
    .stat-box { 
        border-radius: 12px;
        padding: 20px; 
        text-align: center; 
        background: linear-gradient(135deg, rgba(255,255,255,0.05) 0%, rgba(255,255,255,0.01) 100%);
        border: 1px solid rgba(255,255,255,0.1); 
        box-shadow: 0 4px 15px rgba(0,0,0,0.5);
        backdrop-filter: blur(5px);
        transition: all 0.3s ease;
    }
    .stat-box:hover {
        transform: translateY(-5px);
        box-shadow: 0 0 20px rgba(0, 191, 255, 0.2);
    }
    .stat-val { font-size: 2.5em; font-weight: 800; color: white;
        text-shadow: 0 0 10px rgba(255,255,255,0.3); }
    .stat-lbl { font-size: 0.8em; text-transform: uppercase; letter-spacing: 2px; color: rgba(255,255,255,0.6);
        margin-top: 5px; }
    
    /* Playbook Cards */
    .pb-card {
        background: linear-gradient(90deg, #0f0f16 0%, #161625 100%);
        border: 1px solid #333;
        border-left: 4px solid #00BFFF;
        border-radius: 10px;
        padding: 15px;
        margin-bottom: 10px;
        transition: all 0.3s;
    }
    .pb-card:hover {
        border-color: #FF1493;
        box-shadow: 0 0 15px rgba(255, 20, 147, 0.2);
        transform: translateX(5px);
    }

    /* --- CONFIDENCE SCALE --- */
    .conf-scroll-wrapper {
        display: flex;
        overflow-x: auto; padding-bottom: 15px; margin-bottom: 20px; gap: 15px;
        scrollbar-width: thin; scrollbar-color: #00BFFF #111;
    }
    .conf-card {
        flex: 0 0 170px; 
        background: #101018;
        border-radius: 12px; overflow: hidden;
        border: 1px solid #333; text-align: center; 
        transition: transform 0.2s, box-shadow 0.2s;
    }
    .conf-card:hover { transform: translateY(-5px); box-shadow: 0 5px 15px rgba(0,0,0,0.5); border-color: #666;
    }
    .conf-img-container { width: 100%; height: 90px; overflow: hidden; border-bottom: 1px solid #333;
    }
    .conf-img-container img { width: 100%; height: 100%; object-fit: cover; filter: brightness(0.8); transition: filter 0.3s;
    }
    .conf-card:hover img { filter: brightness(1.1); }
    .conf-body { padding: 12px;
    }
    .conf-val { font-size: 1.6em; font-weight: bold; text-shadow: 0 2px 5px rgba(0,0,0,0.8);
    }
    
    /* --- RECENT MATCHES --- */
    .rec-card { 
        background: linear-gradient(to right, rgba(255,255,255,0.03), transparent);
        border-left: 4px solid #555; 
        padding: 12px; margin-bottom: 10px; border-radius: 6px; 
        border-top: 1px solid rgba(255,255,255,0.05);
    }

    /* --- PROTOCOLS (IF/THEN) --- */
    .proto-box { 
        background: rgba(20,20,30,0.6);
        border-left: 3px solid #00BFFF; 
        padding: 12px; margin-bottom: 8px; border-radius: 0 8px 8px 0; 
        border: 1px solid rgba(0,191,255,0.1);
    }
    .proto-if { color: #FFD700; font-size: 0.85em; font-weight: bold; text-transform: uppercase; letter-spacing: 0.5px;
    }
    .proto-then { color: #fff; font-size: 1em; margin-top: 4px; padding-left: 10px; border-left: 1px solid #444;
    }

    /* Inputs */
    div[data-testid="stTextInput"] input, div[data-testid="stTextArea"] textarea {
        background-color: #1a1a25;
        color: white;
        border: 1px solid #444;
    }
    div[data-testid="stSelectbox"] > div > div {
        background-color: #1a1a25;
        color: white;
    }
    
    /* Button Primary override for canvas save */
    button[kind="primary"] {
        background: linear-gradient(90deg, #00BFFF, #FF1493) !important;
        border: none !important;
        color: white !important;
        font-weight: bold !important;
    }

    /* --- VALORANT STYLE MATCH CARD (REVISED LAYOUT) --- */
    .val-card {
        background-color: #121212;
        border-radius: 4px;
        margin-bottom: 8px;
        display: flex; /* Flexbox aktiviert: Elemente liegen nebeneinander */
        align-items: center;
        height: 90px;
        overflow: hidden;
        border: 1px solid #222;
        position: relative;
        transition: transform 0.2s, background-color 0.2s;
    }
    .val-card:hover {
        transform: translateX(4px);
        background-color: #1a1a1a;
    }
    .val-bar {
        width: 6px;
        height: 100%;
        flex-shrink: 0;
        /* Darf nicht schrumpfen */
    }

    /* ZONE 1: MAP & NAME */
    .val-map-section {
        width: 180px;
        /* Feste Breite für den Map-Bereich */
        height: 100%;
        position: relative;
        flex-shrink: 0;
        margin-right: 15px;
    }
    .val-map-bg {
        position: absolute;
        top: 0; left: 0; width: 100%; height: 100%;
        background-size: cover;
        background-position: center;
        /* Verlauf damit Text lesbar ist, aber rechts hart endet für Trennung */
        background: linear-gradient(to right, rgba(0,0,0,0.6) 0%, rgba(0,0,0,0.9) 100%);
        z-index: 0;
    }
    .val-map-text {
        position: absolute;
        top: 50%; left: 15px;
        transform: translateY(-50%);
        z-index: 1;
    }
    .val-map-name {
        font-weight: 900;
        font-size: 1.4em;
        color: white;
        text-transform: uppercase;
        line-height: 1;
        text-shadow: 2px 2px 4px rgba(0,0,0,0.8);
    }

    /* ZONE 2: COMPS (Nach Links gerückt) */
    .val-comps-section {
        display: flex;
        flex-direction: column;
        /* Teams untereinander statt nebeneinander für Platz */
        justify-content: center;
        gap: 4px;
        margin-right: 20px;
        flex-shrink: 0;
    }
    .val-agent-row {
        display: flex;
        align-items: center;
        gap: 2px;
    }
    .val-team-label {
        font-size: 0.6em;
        color: #666;
        width: 20px;
        text-align: right;
        margin-right: 4px;
        font-weight: bold;
    }
    .val-agent-img {
        width: 32px;
        /* Etwas kleiner damit sie untereinander passen */
        height: 32px;
        border-radius: 3px;
        border: 1px solid #333;
        background: #000;
    }

    /* ZONE 3: STATS (Die Mitte - Der neue "Freie Platz") */
    .val-stats-section {
        flex-grow: 1;
        /* Nimmt den restlichen Platz ein */
        display: flex;
        flex-direction: row;
        justify-content: center;
        /* Zentriert die Stats */
        align-items: center;
        gap: 20px;
        /* Abstand zwischen den Stat-Gruppen */
        color: #ccc;
        font-family: 'Segoe UI', sans-serif;
    }
    .stat-group {
        display: flex;
        flex-direction: column;
        align-items: center;
    }
    .stat-label {
        font-size: 0.65em;
        text-transform: uppercase;
        color: #777;
        letter-spacing: 1px;
        margin-bottom: 2px;
    }
    .stat-value {
        font-size: 0.9em;
        font-weight: 700;
        color: #eee;
    }
    .stat-date {
        font-size: 0.8em;
        color: #aaa;
        font-style: italic;
    }

    /* ZONE 4: SCORE (Rechts) */
    .val-score-section {
        width: 120px;
        text-align: right;
        padding-right: 20px;
        flex-shrink: 0;
    }
    .val-score {
        font-weight: 900;
        font-size: 2.2em;
        line-height: 1;
    }
    .val-vod-link {
        font-size: 0.75em;
        font-weight: 700;
        text-transform: uppercase;
        display: block;
        margin-top: 4px;
        text-decoration: none;
        opacity: 0.7;
    }
    .val-vod-link:hover { opacity: 1; text-decoration: underline;
    }

    /* --- POWER RANKING CARD DESIGN (FIXED & COMPACT) --- */
    .rank-row {
        background-color: #121212;
        border: 1px solid #222;
        border-radius: 4px;
        margin-bottom: 6px;
        display: flex;
        align-items: center;
        padding: 4px 8px;
        /* Etwas mehr seitliches Padding */
        height: 54px;
        /* Kompakte Höhe */
        transition: transform 0.2s, background-color 0.2s;
        overflow: hidden;
    }
    .rank-row:hover {
        background-color: #1a1a1a;
        transform: translateX(4px);
        border-color: #333;
    }
    
    /* Map Bild */
    .rank-img-box {
        width: 120px;
        height: 100%;
        border-radius: 3px;
        overflow: hidden;
        margin-right: 12px;
        flex-shrink: 0;
    }
    .rank-img {
        width: 100%;
        height: 100%;
        object-fit: cover;
    }
    
    /* Map Name - Schriftgröße angepasst */
    .rank-name {
        width: 80px;
        font-family: 'Segoe UI', sans-serif;
        font-weight: 800;
        font-size: 14px;          /* Fest auf 14px gesetzt */
        color: white;
        text-transform: uppercase;
        margin-right: 10px;
        flex-shrink: 0;
        white-space: nowrap;
        overflow: hidden;
        text-overflow: ellipsis;
    }

    /* Stats Container */
    .rank-stats {
        flex-grow: 1;
        display: flex;
        flex-direction: column;
        justify-content: center;
        gap: 4px;                 /* Kleiner Abstand zwischen den Balken */
    }
    
    /* Einzelne Zeile (Label - Balken - Zahl) */
    .stat-line {
        display: flex;
        align-items: center;
        height: 18px;             /* Fixe Höhe pro Zeile */
    }
    
    /* Label (RATING / WIN%) */
    .stat-label {
        width: 50px;
        /* Etwas breiter damit nichts umbricht */
        color: #666;
        font-weight: 700;
        font-size: 10px;
        /* Sehr klein und fein */
        text-transform: uppercase;
        letter-spacing: 0.5px;
    }
    
    /* Balken Hintergrund */
    .prog-bg {
        flex-grow: 1;
        height: 6px;              /* Dünner Balken */
        background-color: #252525;
        border-radius: 3px;
        margin: 0 10px;
        overflow: hidden;
    }
    
    /* Balken Füllung */
    .prog-fill {
        height: 100%;
        border-radius: 3px;
    }
    
    /* Die Zahl am Ende (WICHTIG!) */
    .stat-val {
        width: 45px;
        /* Verbreitert! Vorher 30px -> zu eng */
        text-align: right;
        font-weight: 800;
        font-family: 'Consolas', 'Monaco', monospace; /* Monospace für saubere Ausrichtung */
        font-size: 12px;
        /* Gut lesbare Größe */
        line-height: 1;
    }

    /* --- CALENDAR & TODO --- */
    .cal-day-box {
        min-height: 120px;
        /* Bigger height for calendar boxes */
        background-color: #121212;
        border: 1px solid #333;
        border-radius: 6px;
        padding: 8px;
        margin: 2px;
        transition: transform 0.2s, border-color 0.2s;
        overflow: hidden;
    }
    .cal-day-box:hover {
        border-color: #00BFFF;
        background-color: #1a1a1a;
    }
    .cal-date {
        font-weight: 900;
        color: #555;
        margin-bottom: 5px;
        font-size: 1.1em;
    }
    .cal-today {
        border: 1px solid #00BFFF !important;
        background-color: #0f1820 !important;
    }
    .cal-today .cal-date { color: #00BFFF;
    }
    
    .cal-event-pill {
        margin-top: 3px;
        padding: 3px 6px;
        border-radius: 3px;
        font-size: 0.75em;
        font-weight: 700;
        white-space: nowrap;
        overflow: hidden;
        text-overflow: ellipsis;
        color: white;
        display: block;
        text-shadow: 0 1px 2px rgba(0,0,0,0.5);
        border-left-width: 3px;
        border-left-style: solid;
    }
</style>
""", unsafe_allow_html=True)

# --- PFADE ---
BASE_DIR = os.path.dirname(os.path.abspath(__file__))
DATA_DIR_JSON = os.path.join(BASE_DIR, "data", "matches")
ASSET_DIR = os.path.join(BASE_DIR, "assets")
STRAT_IMG_DIR = os.path.join(ASSET_DIR, "strats")
VOD_IMG_DIR = os.path.join(ASSET_DIR, "vod_imgs")
PLAYBOOKS_FILE = os.path.join(BASE_DIR, "data", "playbooks.csv")
PRESETS_FILE = os.path.join(BASE_DIR, "data", "pdf_presets.json")

for d in [DATA_DIR_JSON, os.path.join(BASE_DIR, "data"), STRAT_IMG_DIR, VOD_IMG_DIR, os.path.join(ASSET_DIR, "maps"), os.path.join(ASSET_DIR, "agents"), os.path.join(ASSET_DIR, "fonts"), os.path.join(ASSET_DIR, "playbook")]:
    if not os.path.exists(d): os.makedirs(d)

OUR_TEAM = ["Trashies", "Luggi", "Umbra", "Noctis", "n0thing", "Gengar"]
DISCORD_WEBHOOK_URL = "https://discord.com/api/webhooks/1452361543751565446/OVzSvaZ7LYCzU9gKW6QFyaV84Edxfi_6rF7Jjz5QxlpZlXfbC3gGQKzUoX0k_Q9TMC6f"

PLAYER_DISCORD_MAPPING = {
    "Luggi": "<@713352448034734081>",
    "Benni": "<@728923466161717338>",
    "Andrei": "<@262685797189287937>",
    "Luca": "<@665846859578867732>",
    "Sofi": "<@591563019046223894>",
    "Remus": "<@293790883420045314>"
}

# --- ABILITY MAPPING DATABASE (MOVED TO GLOBAL SCOPE) ---
AGENT_ABILITIES = {
    "Astra": {"C": "Gravity Well", "Q": "Nova Pulse", "E": "Nebula", "X": "Astral Form"},
    "Breach": {"C": "Aftershock", "Q": "Flashpoint", "E": "Fault Line", "X": "Rolling Thunder"},
    "Brimstone": {"C": "Stim Beacon", "Q": "Incendiary", "E": "Sky Smoke", "X": "Orbital Strike"},
    "Chamber": {"C": "Trademark", "Q": "Headhunter", "E": "Rendezvous", "X": "Tour De Force"},
    "Clove": {"C": "Pick-me-up", "Q": "Meddle", "E": "Ruse", "X": "Not Dead Yet"},
    "Cypher": {"C": "Trapwire", "Q": "Cyber Cage", "E": "Spycam", "X": "Neural Theft"},
    "Deadlock": {"C": "GravNet", "Q": "Sonic Sensor", "E": "Barrier Mesh", "X": "Annihilation"},
    "Fade": {"C": "Prowler", "Q": "Seize", "E": "Haunt", "X": "Nightfall"},
    "Gekko": {"C": "Mosh Pit", "Q": "Wingman", "E": "Dizzy", "X": "Thrash"},
    "Harbor": {"C": "Cascade", "Q": "Cove", "E": "High Tide", "X": "Reckoning"},
    "Iso": {"C": "Contingency", "Q": "Undercut", "E": "Double Tap", "X": "Kill Contract"},
    "Jett": {"C": "Cloudburst", "Q": "Updraft", "E": "Tailwind", "X": "Blade Storm"},
    "KAY/O": {"C": "FRAG/ment", "Q": "FLASH/drive", "E": "ZERO/point", "X": "NULL/cmd"},
    "Killjoy": {"C": "Nanoswarm", "Q": "Alarmbot", "E": "Turret", "X": "Lockdown"},
    "Neon": {"C": "Fast Lane", "Q": "Relay Bolt", "E": "High Gear", "X": "Overdrive"},
    "Omen": {"C": "Shrouded Step", "Q": "Paranoia", "E": "Dark Cover", "X": "From the Shadows"},
    "Phoenix": {"C": "Blaze", "Q": "Curveball", "E": "Hot Hands", "X": "Run it Back"},
    "Raze": {"C": "Boom Bot", "Q": "Blast Pack", "E": "Paint Shells", "X": "Showstopper"},
    "Reyna": {"C": "Leer", "Q": "Devour", "E": "Dismiss", "X": "Empress"},
    "Sage": {"C": "Barrier Orb", "Q": "Slow Orb", "E": "Healing Orb", "X": "Resurrection"},
    "Skye": {"C": "Regrowth", "Q": "Trailblazer", "E": "Guiding Light", "X": "Seekers"},
    "Sova": {"C": "Owl Drone", "Q": "Shock Bolt", "E": "Recon Bolt", "X": "Hunter's Fury"},
    "Viper": {"C": "Snake Bite", "Q": "Poison Cloud", "E": "Toxic Screen", "X": "Viper's Pit"},
    "Vyse": {"C": "Razorvine", "Q": "Shear", "E": "Arc Rose", "X": "Steel Garden"},
    "Yoru": {"C": "Fakeout", "Q": "Blindside", "E": "Gatecrash", "X": "Dimensional Drift"},
}

def send_discord_notification(player_name, task_title, description):
    ping = PLAYER_DISCORD_MAPPING.get(player_name, player_name)
    content = f"📢 **New Task Assigned!**\n\n👤 **Player:** {ping}\n📝 **Task:** {task_title}\nℹ️ **Details:** {description}\n\n_Check the Nexus Dashboard (https://nxs-dashboard.streamlit.app) for more info._"
    try: requests.post(DISCORD_WEBHOOK_URL, json={"content": content})
    except Exception as e: print(f"Discord Webhook Error: {e}")

# --- HELPER FUNKTIONEN (Assets) ---
def get_map_img(map_name, type='list'):
    if not map_name or pd.isna(map_name): return None
    name_clean = str(map_name).lower().strip()
    target_path = os.path.join(ASSET_DIR, "maps", f"{name_clean}_{type}.png")
    if os.path.exists(target_path): return target_path
    simple_path = os.path.join(ASSET_DIR, "maps", f"{name_clean}.png")
    if os.path.exists(simple_path): return simple_path
    return None

def get_agent_img(agent_name):
    if not agent_name or pd.isna(agent_name) or str(agent_name).lower() == 'nan': return None
    clean = str(agent_name).lower().replace("/", "").strip()
    path = os.path.join(ASSET_DIR, "agents", f"{clean}.png")
    return path if os.path.exists(path) else None

def create_styled_agent_pil(agent_name):
    img_path = get_agent_img(agent_name)
    if not img_path: return None
    try:
        with Image.open(img_path) as img:
            img = img.convert("RGBA")
            img.thumbnail((100, 100), Image.Resampling.LANCZOS)
            ac = {"astra": "#653491", "breach": "#bc5434", "brimstone": "#d56e23", "chamber": "#e3b62d", "clove": "#e882a8", "cypher": "#d6d6d6", "deadlock": "#bcc6cc", "fade": "#4c4c4c", "gekko": "#b6ff59", "harbor": "#2d6e68", "iso": "#4b48ac", "jett": "#90e0ef", "kay/o": "#4bb0a8", "killjoy": "#f7d336", "neon": "#2c4f9e", "omen": "#4f4f8f", "phoenix": "#ff7f50", "raze": "#ff6a00", "reyna": "#b74b8e", "sage": "#52ffce", "skye": "#8fbc8f", "sova": "#6fa8dc", "viper": "#32cd32", "tejo": "#E97223", "yoru": "#334488", "vyse": "#7b68ee"}
            bg_col = ac.get(str(agent_name).lower(), "#2c003e")
            bg = Image.new("RGBA", (100, 100), bg_col)
            offset = ((100 - img.width) // 2, (100 - img.height) // 2)
            bg.paste(img, offset, img)
            mask = Image.new("L", (100, 100), 0)
            draw = ImageDraw.Draw(mask)
            draw.ellipse((0, 0, 100, 100), fill=255)
            final = Image.new("RGBA", (100, 100), (0,0,0,0))
            final.paste(bg, (0,0), mask=mask)
            return final
    except: return None

def get_styled_agent_img_b64(agent_name):
    img = create_styled_agent_pil(agent_name)
    if img:
        buff = io.BytesIO()
        img.save(buff, format="PNG")
        return base64.b64encode(buff.getvalue()).decode()
    return img_to_b64(get_agent_img(agent_name))

def create_team_composite(agents):
    images = [create_styled_agent_pil(a) for a in agents if a]
    images = [i for i in images if i is not None]
    if not images: return None
    w, h = images[0].size
    total_w = w * len(images) + 10 * (len(images)-1)
    comp = Image.new("RGBA", (total_w, h), (0,0,0,0))
    x = 0
    for img in images:
        comp.paste(img, (x, 0))
        x += w + 10
    buff = io.BytesIO()
    comp.save(buff, format="PNG")
    buff.seek(0)
    return buff

def img_to_b64(img_path):
    if not img_path or not os.path.exists(img_path): return ""
    with open(img_path, "rb") as f: data = f.read()
    return base64.b64encode(data).decode()

def render_rich_notes(text):
    if not text: return
    def _repl(m):
        fp = os.path.join(VOD_IMG_DIR, m.group(1))
        if os.path.exists(fp):
            return f'<img src="data:image/png;base64,{img_to_b64(fp)}" style="max-width:100%;border:1px solid #444;border-radius:4px;margin:5px 0">'
        return f"`[Image: {m.group(1)} not found]`"
    processed = re.sub(r"\[\[img:(.*?)\]\]", _repl, text)
    st.markdown(processed, unsafe_allow_html=True)

def get_yt_thumbnail(url):
    if not url or "youtu" not in str(url): return None
    vid_id = None
    if "v=" in url: vid_id = url.split("v=")[1].split("&")[0]
    elif "youtu.be/" in url: vid_id = url.split("youtu.be/")[1].split("?")[0]
    return f"https://img.youtube.com/vi/{vid_id}/mqdefault.jpg" if vid_id else None

def render_visual_selection(options, type_item, key_prefix, default=None, multi=True, key_state=None):
    selected = default if default is not None else []
    if not multi and key_state and key_state in st.session_state:
        current_selection = st.session_state[key_state]
    else: current_selection = default
    st.markdown("<style>div[data-testid='stColumn'] {text-align: center;} div[data-testid='stCheckbox'] {display: inline-block;}</style>", unsafe_allow_html=True)
    cols = st.columns(6 if type_item == 'map' else 8)
    for i, opt in enumerate(options):
        with cols[i % len(cols)]:
            if type_item == 'agent':
                b64 = get_styled_agent_img_b64(opt)
                if b64: st.image(f"data:image/png;base64,{b64}", width=55)
                else: st.info(opt[:3])
            else:
                img_path = get_map_img(opt, 'list')
                if img_path: st.image(img_path, use_container_width=True)
                else: st.info(opt[:3])
            
            if multi:
                if st.checkbox(" ", key=f"{key_prefix}_{opt}", value=(opt in selected), label_visibility="collapsed"):
                    if opt not in selected: selected.append(opt)
                elif opt in selected: selected.remove(opt)
            else:
                if st.button("Select", key=f"{key_prefix}_{opt}"):
                    if key_state: st.session_state[key_state] = opt
                    st.rerun()
    return selected

# ==============================================================================
# 💾 SUPABASE DATA HANDLERS
# ==============================================================================

# --- NEW: SURGICAL DB OPERATIONS (FAST) ---
def db_insert(table_name, record, state_key=None):
    """Inserts a single record directly. Much faster than syncing the whole DF."""
    try:
        # Map ID -> id for Supabase
        rec = record.copy()
        if 'ID' in rec: rec['id'] = rec.pop('ID')
        
        supabase.table(table_name).insert(rec).execute()
        
        # Optimistic UI Update
        if state_key and state_key in st.session_state:
            st.session_state[state_key] = pd.concat([st.session_state[state_key], pd.DataFrame([record])], ignore_index=True)
            
        return True
    except Exception as e:
        st.error(f"Database Insert Error ({table_name}): {e}")
        return False

def db_update(table_name, record, state_key=None, id_col='ID'):
    """Updates a single record by ID."""
    try:
        rec = record.copy()
        if 'ID' in rec: rec['id'] = rec.pop('ID')
        if 'id' not in rec: raise ValueError("Missing ID for update")
        
        supabase.table(table_name).update(rec).eq("id", rec['id']).execute()
        
        # Optimistic UI Update
        if state_key and state_key in st.session_state:
            df = st.session_state[state_key]
            # Find index of row with this ID
            idx = df[df[id_col] == record.get(id_col)].index
            if not idx.empty:
                for k, v in record.items():
                    df.at[idx[0], k] = v
                st.session_state[state_key] = df

        return True
    except Exception as e:
        st.error(f"Database Update Error ({table_name}): {e}")
        return False

def db_delete(table_name, record_id, state_key=None, id_col='ID'):
    """Deletes a single record by ID."""
    try:
        supabase.table(table_name).delete().eq("id", record_id).execute()
        
        if state_key and state_key in st.session_state:
            df = st.session_state[state_key]
            st.session_state[state_key] = df[df[id_col] != record_id]
            
        return True
    except Exception as e:
        st.error(f"Database Delete Error ({table_name}): {e}")
        return False

def _sync_data(table_name, df):
    """
    Synchronizes the DataFrame with Supabase.
    Performs UPSERT for existing/new rows.
    Performs DELETE for rows that are in DB but missing in DF (based on 'ID').
    """
    try:
        # 1. Fetch current IDs from DB
        curr_res = supabase.table(table_name).select("id").execute()
        db_ids = {row['id'] for row in curr_res.data}
        
        # 2. Get IDs from current DF (Assuming column is 'ID')
        # Ensure 'ID' exists
        if 'ID' not in df.columns and 'id' in df.columns:
            df.rename(columns={'id': 'ID'}, inplace=True)
        
        if 'ID' in df.columns:
            df_ids = set(df['ID'].astype(str).tolist())
            
            # 3. Identify IDs to delete
            ids_to_delete = list(db_ids - df_ids)
            
            if ids_to_delete:
                supabase.table(table_name).delete().in_("id", ids_to_delete).execute()

        # 4. Prepare data for Upsert
        df_save = df.copy()
        # Rename ID -> id for Supabase
        if 'ID' in df_save.columns:
            df_save.rename(columns={'ID': 'id'}, inplace=True)
        
        # Clean data (NaN -> None)
        records = df_save.where(pd.notnull(df_save), None).to_dict(orient='records')
        
        if records:
            supabase.table(table_name).upsert(records).execute()
    except Exception as e:
        print(f"Sync Error ({table_name}): {e}")

def _upsert_only(table_name, df):
    """
    Only Upserts data. Good for append-only logs like matches/stats.
    """
    try:
        df_save = df.copy()
        if 'ID' in df_save.columns: df_save.rename(columns={'ID': 'id'}, inplace=True)
        # MatchID handling: If Supabase uses 'id' as PK but MatchID is the logical key,
        # we rely on the schema. Assuming MatchID is just a column.
        
        records = df_save.where(pd.notnull(df_save), None).to_dict(orient='records')
        if records:
            supabase.table(table_name).upsert(records).execute()
    except Exception as e:
        print(f"Upsert Error ({table_name}): {e}")

# --- MAPPING SAVE FUNCTIONS ---
# Since the prompt asked to keep "id" as primary key, we use _sync_data 
# for tables where deletions occur in the UI by filtering the DF.

def save_matches(df_new): _upsert_only("nexus_matches", df_new)
def save_player_stats(df_new): _upsert_only("Premier - PlayerStats", df_new)
def save_scrims(df_new): _sync_data("scrims", df_new)
def save_scrim_availability(df_new): _upsert_only("scrim_availability", df_new)
def save_player_todos(df_new): _sync_data("player_todos", df_new)
# Bulk saves (still used for data_editor or reordering)
def save_legacy_playbooks(df_new): _sync_data("playbooks", df_new)
def save_pb_strats(df_new): _sync_data("nexus_pb_strats", df_new) # Reordering logic relies on sync
def save_map_theory(df_new): _sync_data("nexus_map_theory", df_new)
def save_resources(df_new): _sync_data("resources", df_new)
def save_calendar(df_new): _sync_data("calendar", df_new)
def save_team_playbooks(df_new): _sync_data("nexus_playbooks", df_new)
def save_simple_todos(df_new): _sync_data("todo", df_new)
def save_vod_reviews(df_new): _sync_data("nexus_vod_reviews", df_new)
def save_lineups(df_new): _sync_data("nexus_lineups", df_new)

# Specific Availability Update
def update_availability(scrim_id, player, status):
    try:
        # Upsert specific row
        data = {
            'ScrimID': scrim_id,
            'Player': player,
            'Available': status,
            'UpdatedAt': datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        }
        # Note: Scrim Availability might not have a simple 'id'. 
        # If 'id' is auto-gen, we need to find the row first or use match constraints.
        # Assuming 'id' exists. We query to find ID if exists.
        res = supabase.table("scrim_availability").select("id").eq("ScrimID", scrim_id).eq("Player", player).execute()
        if res.data:
            data['id'] = res.data[0]['id']
        
        supabase.table("scrim_availability").upsert(data).execute()
    except Exception as e: print(f"Avail Update Error: {e}")

def delete_scrim(scrim_id):
    try:
        supabase.table("scrims").delete().eq("id", scrim_id).execute()
        # Also delete availability
        supabase.table("scrim_availability").delete().eq("ScrimID", scrim_id).execute()
        st.cache_data.clear()
    except Exception as e: print(f"Del Scrim Error: {e}")

# ==============================================================================
# 🚀 DATA LOADER
# ==============================================================================

# REMOVED @st.cache_data to allow session_state management
def load_data(dummy=None):
    # Define mapping: Table Name -> Session State Key
    table_map = {
        "nexus_matches": "df_matches",
        "Premier - PlayerStats": "df_p",
        "scrims": "df_scrims",
        "scrim_availability": "df_availability",
        "player_todos": "df_todos",
        "nexus_playbooks": "df_team_pb",
        "playbooks": "df_legacy_pb",
        "nexus_pb_strats": "df_pb_strats",
        "nexus_map_theory": "df_theory",
        "resources": "df_res",
        "calendar": "df_cal",
        "todo": "df_simple_todos",
        "nexus_vod_reviews": "df_vods",
        "nexus_lineups": "df_lineups"
    }

    # Check if we need to load data (if any key is missing)
    missing_keys = [k for k in table_map.values() if k not in st.session_state]
    
    if missing_keys:
        prog_bar = st.progress(0, text="Connecting to Supabase...")
        total = len(table_map)
        
        for i, (tbl, key) in enumerate(table_map.items()):
            if key not in st.session_state:
                prog_bar.progress(i / total, text=f"Fetching {tbl}...")
                try:
                    response = supabase.table(tbl).select("*").execute()
                    df = pd.DataFrame(response.data)
                    if 'id' in df.columns: df.rename(columns={'id': 'ID'}, inplace=True)
                    st.session_state[key] = df
                except Exception as e:
                    st.session_state[key] = pd.DataFrame()
        
        prog_bar.empty()

    # Helper to safely get DF from state with columns
    def get_df_state(key, cols=None):
        df = st.session_state.get(key, pd.DataFrame())
        if df is None or df.empty:
            return pd.DataFrame(columns=cols if cols else [])
        if cols:
            for col in cols:
                if col not in df.columns: df[col] = None
        return df

    # Data Assignment (Same structure as before)
    df = get_df_state("df_matches", ['Date', 'Map', 'Result', 'Score_Us', 'Score_Enemy', 'MatchID'])
    if not df.empty and 'Date' in df.columns:
        if 'Map' in df.columns: df['Map'] = df['Map'].astype(str).str.strip().str.title()
        if 'Result' in df.columns: df['Result'] = df['Result'].astype(str).str.strip().str.upper()
        df['DateObj'] = pd.to_datetime(df['Date'], dayfirst=True, errors='coerce')
        for c in ['Score_Us', 'Score_Enemy', 'Atk_R_W', 'Def_R_W', 'Atk_R_L', 'Def_R_L']:
            if c in df.columns: df[c] = pd.to_numeric(df[c], errors='coerce').fillna(0)
        df['Delta'] = df['Score_Us'] - df['Score_Enemy']
    
    df_p = get_df_state("df_p")
    df_scrims = get_df_state("df_scrims", ['ID', 'Title', 'Date', 'Time', 'Map', 'Description', 'CreatedBy', 'CreatedAt', 'PlaybookLink', 'VideoLink'])
    if not df_scrims.empty:
        df_scrims['DateTimeObj'] = pd.to_datetime(df_scrims['Date'] + ' ' + df_scrims['Time'], format="%Y-%m-%d %H:%M", errors='coerce')
    
    df_availability = get_df_state("df_availability", ['ScrimID', 'Player', 'Available', 'UpdatedAt'])

    df_todos = get_df_state("df_todos", ['ID', 'Player', 'Title', 'Description', 'PlaybookLink', 'YoutubeLink', 'AssignedBy', 'AssignedAt', 'Completed', 'CompletedAt'])
    if not df_todos.empty and 'Completed' in df_todos.columns:
        df_todos['Completed'] = df_todos['Completed'].astype(str).str.lower() == 'true'

    df_team_pb = get_df_state("df_team_pb", ['ID', 'Map', 'Name', 'Agent_1', 'Agent_2', 'Agent_3', 'Agent_4', 'Agent_5'])
    df_legacy_pb = get_df_state("df_legacy_pb", ['Map', 'Name', 'Link', 'Agent_1', 'Agent_2', 'Agent_3', 'Agent_4', 'Agent_5'])
    df_pb_strats = get_df_state("df_pb_strats", ['PB_ID', 'Strat_ID', 'Name', 'Image', 'Protocols', 'Notes', 'Tag', 'Order'])
    if 'Notes' not in df_pb_strats.columns: df_pb_strats['Notes'] = ""
    if 'Tag' not in df_pb_strats.columns: df_pb_strats['Tag'] = "Default"
    if 'Order' not in df_pb_strats.columns: df_pb_strats['Order'] = range(len(df_pb_strats))
    df_theory = get_df_state("df_theory", ['Map', 'Section', 'Content', 'Image'])
    
    df_res = get_df_state("df_res", ['Title', 'Link', 'Category', 'Note'])
    df_cal = get_df_state("df_cal", ['Date', 'Time', 'Event', 'Map', 'Type', 'Players'])
    df_simple_todos = get_df_state("df_simple_todos", ['Task', 'Done'])
    if not df_simple_todos.empty and 'Done' in df_simple_todos.columns:
        df_simple_todos['Done'] = df_simple_todos['Done'].astype(str).str.lower() == 'true'
        
    df_vods = get_df_state("df_vods", ['ID', 'Title', 'Type', 'VideoLink', 'Map', 'Agent', 'Player', 'Notes', 'Tags', 'CreatedBy', 'CreatedAt', 'Rounds'])
    df_lineups = get_df_state("df_lineups", ['ID', 'Map', 'Agent', 'Side', 'Type', 'Title', 'Image', 'VideoLink', 'Description', 'Tags', 'CreatedBy'])

    return df, df_p, df_scrims, df_availability, df_todos, df_team_pb, df_legacy_pb, df_pb_strats, df_theory, df_res, df_cal, df_simple_todos, df_vods, df_lineups

# ==============================================================================
# 🛠️ PARSER FUNCTION (Unchanged)
# ==============================================================================
def parse_tracker_json(file_input):
    try:
        if isinstance(file_input, str):
            with open(file_input, 'r', encoding='utf-8') as f: data = json.load(f)
        else: data = json.load(file_input)
        parsed_data = []
        
        matches = []
        if isinstance(data, dict):
            if 'data' in data and isinstance(data['data'], dict) and 'matches' in data['data']:
                matches = data['data']['matches']
            elif 'data' in data and isinstance(data['data'], dict) and 'metadata' in data['data']:
                matches = [data['data']]
            elif 'metadata' in data and 'segments' in data:
                matches = [data]
        elif isinstance(data, list):
            matches = data
        
        if matches:
            for m in matches:
                meta = m.get('metadata', {})
                segments = m.get('segments', [])
                p_segs = [s for s in segments if s.get('type') == 'player-summary']
                
                for p_seg in p_segs:
                    stats = p_seg.get('stats', {})
                    attrs = p_seg.get('attributes', {})
                    kills = stats.get('kills', {}).get('value', 0)
                    deaths = stats.get('deaths', {}).get('value', 1)
                    assists = stats.get('assists', {}).get('value', 0)
                    hs = stats.get('headshots', {}).get('value', 0)
                    total_hits = hs + stats.get('bodyshots', {}).get('value', 0) + stats.get('legshots', {}).get('value', 0)
                    hs_percent = (hs / total_hits * 100) if total_hits > 0 else 0
                    c_grenade = stats.get('grenadeCasts', {}).get('value', 0)
                    c_abil1 = stats.get('ability1Casts', {}).get('value', 0)
                    c_abil2 = stats.get('ability2Casts', {}).get('value', 0)
                    c_ult = stats.get('ultimateCasts', {}).get('value', 0)
                    fk = stats.get('firstBloods', {}).get('value', 0)
                    fd = stats.get('firstDeaths', {}).get('value', 0)
                    clutches = sum([stats.get(f'clutches1v{i}', {}).get('value', 0) for i in range(1, 6)])
                    kast = stats.get('kast', {}).get('value', 0)
                    rounds = stats.get('roundsPlayed', {}).get('value', 1)
                    if kast == 0 and rounds > 0: kast = ((kills + assists + (rounds - deaths)) / rounds) * 100
                    res = "Unknown"
                    if 'result' in meta: res = meta['result']
                    elif 'hasWon' in stats: res = "Victory" if stats['hasWon']['value'] else "Defeat"
                    
                    parsed_data.append({
                        'MatchID': m.get('attributes', {}).get('id', 'Unknown'),
                        'Date': meta.get('modeName', 'Ranked'),
                        'Map': meta.get('mapName', 'Unknown'),
                        'Player': attrs.get('platformUserIdentifier', 'Unknown').split('#')[0],
                        'Agent': p_seg.get('metadata', {}).get('agentName', 'Unknown'),
                        'Result': res,
                        'Kills': kills, 'Deaths': deaths, 'Assists': assists,
                        'KD': kills/deaths if deaths>0 else kills,
                        'HS%': hs_percent,
                        'ADR': stats.get('damagePerRound', {}).get('value', 0),
                        'Rounds': rounds,
                        'Cast_Grenade': c_grenade, 'Cast_Abil1': c_abil1,
                        'Cast_Abil2': c_abil2, 'Cast_Ult': c_ult,
                        'Total_Util': c_grenade + c_abil1 + c_abil2 + c_ult,
                        'FK': fk, 'FD': fd, 'Clutches': clutches, 'KAST': kast,
                        'MatchesPlayed': 1, 'Wins': 1 if res == "Victory" else 0
                    })
        return pd.DataFrame(parsed_data)
    except Exception as e: print(f"Parser Error: {e}"); return pd.DataFrame()

# ==============================================================================
# APP UI START
# ==============================================================================

# ... (Die restliche App-Logik und UI-Code ist identisch, da die Variable df, etc. 
# nun durch die neue load_data Funktion befüllt werden und die save_... Funktionen 
# kompatibel gehalten wurden. Ich füge den Rest des Codes hier ein, damit es copy-paste ready ist.)
# HINWEIS: Um Zeichenlimit zu sparen, nehme ich an, dass der Rest des Codes (Zeilen 1168 bis Ende) 
# exakt gleich bleibt, da wir nur die 'Backend'-Funktionen oben ausgetauscht haben.

# ... [Fügen Sie hier den Rest des Codes ab Zeile 1168 aus Ihrer ursprünglichen Datei ein] ...

# ==============================================================================
# 🚀 APP START & DATEN LADEN
# ==============================================================================

# HIER WERDEN DIE DATEN GELADEN UND DIE VARIABLEN GLOBAL GESETZT
LOADING_QUOTES = [
    "Planting the Spike...",
    "Checking corners...",
    "Asking Sage for a heal...",
    "Rotating to A...",
    "Calculating lineups...",
    "Jett is dashing in...",
    "Sova is droning...",
    "Reviewing the VOD...",
    "Eco round...",
    "Buying shields...",
    "Defusing...",
    "Flash out!",
    "Reviving...",
    "One tap machine warming up...",
    "Installing aimbot... (just kidding)",
    "Fetching the strat book...",
    "Analyzing enemy patterns...",
    "Hydrating...",
    "Tactical timeout...",
    "Rush B do not stop..."
]

with st.spinner(random.choice(LOADING_QUOTES)):
    df, df_players, df_scrims, df_availability, df_todos, df_team_pb, df_legacy_pb, df_pb_strats, df_theory, df_res, df_cal, df_simple_todos, df_vods, df_lineups = load_data()

# Handle navigation triggers
if "trigger_navigation" in st.session_state:
    st.session_state["navigation_radio"] = st.session_state["trigger_navigation"]
    del st.session_state["trigger_navigation"]

with st.sidebar:
    if os.path.exists("logo.png"): st.image("logo.png", width=140)
    st.title("NEXUS")
    
    # --- NAVIGATION LOGIK ---
    all_pages = ["🏠 DASHBOARD", "👥 COACHING", "⚽ SCRIMS", "📝 MATCH ENTRY", "🗺️ MAP ANALYZER", "📘 STRATEGY BOARD", "📚 RESOURCES", "📅 CALENDAR", "📊 PLAYERS", "📹 VOD REVIEW", "💾 DATABASE"]
    
    # Hole die erlaubten Seiten aus dem Session State (gesetzt beim Login)
    # Fallback auf alle Seiten, falls Session State leer ist
    allowed_pages = st.session_state.get('allowed_pages', all_pages)
    
    # Navigation erstellen
    page = st.radio("NAVIGATION", allowed_pages, label_visibility="collapsed", key="navigation_radio")
    
    st.markdown("---")
    
    # --- RELOAD BUTTON (MIT CACHE CLEAR) ---
    if st.button("🔄 Reload Data (if you created/deleted something)"): 
        for k in ["df_matches","df_p","df_scrims","df_availability","df_todos","df_team_pb","df_legacy_pb","df_pb_strats","df_theory","df_res","df_cal","df_simple_todos","df_vods","df_lineups"]:
            if k in st.session_state: del st.session_state[k]
        st.rerun()
    
    # --- USER INFO & LOGOUT ---
    st.markdown("---")
    col1, col2 = st.columns(2)
    with col1:
        st.markdown(f"**User:** {st.session_state.get('username', 'Unknown')}")
    with col2:
        if st.button("🚪 Logout"):
            logout()

# ==============================================================================
# 1. DASHBOARD
# ==============================================================================
if page == "🏠 DASHBOARD":
    # Get current user info
    current_user = st.session_state.get('username', '')
    user_role = st.session_state.get('role', '')

    # --- Data for "What's Next" widget ---
    incomplete_todos = 0
    next_event_str = "No upcoming events."
    next_event_nav = None

    # Player-specific data
    if user_role == 'player' and current_user:
        player_todos = df_todos[(df_todos['Player'] == current_user) & (df_todos['Completed'] == False)]
        incomplete_todos = len(player_todos)
    # Upcoming events (Scrims and Calendar)
    now = pd.Timestamp.now()
    upcoming_events = []
    if not df_scrims.empty and 'DateTimeObj' in df_scrims.columns:
        upcoming_scrims = df_scrims[df_scrims['DateTimeObj'] >= now].copy()
        if not upcoming_scrims.empty:
            upcoming_scrims['Type'] = 'Scrim'
            upcoming_scrims['Title'] = upcoming_scrims['Title']
            upcoming_scrims['DateTime'] = upcoming_scrims['DateTimeObj']
            upcoming_events.append(upcoming_scrims[['DateTime', 'Title', 'Type']])

    if not df_cal.empty and 'Date' in df_cal.columns and 'Time' in df_cal.columns:
        df_cal_copy = df_cal.copy()
        df_cal_copy['DateTime'] = pd.to_datetime(df_cal_copy['Date'] + ' ' + df_cal_copy['Time'], format="%d.%m.%Y %H:%M", errors='coerce')
        upcoming_cal = df_cal_copy[df_cal_copy['DateTime'] >= now]
        if not upcoming_cal.empty:
            upcoming_cal = upcoming_cal.rename(columns={'Event': 'Title'})
            upcoming_events.append(upcoming_cal[['DateTime', 'Title', 'Type']])

    if upcoming_events:
        all_upcoming = pd.concat(upcoming_events).sort_values('DateTime').iloc[0]
        event_time = all_upcoming['DateTime']
        time_delta = event_time - now
        
        if time_delta.days < 1:
            hours, rem = divmod(time_delta.seconds, 3600)
            minutes, _ = divmod(rem, 60)
            if hours > 0:
                next_event_str = f"**{all_upcoming['Title']}** in {hours}h {minutes}m"
            else:
                next_event_str = f"**{all_upcoming['Title']}** in {minutes}m"
        else:
            next_event_str = f"**{all_upcoming['Title']}** in {time_delta.days} day(s)"
        
        next_event_nav = "⚽ SCRIMS" if all_upcoming['Type'] == 'Scrim' else "📅 CALENDAR"

    st.title("PERFORMANCE DASHBOARD")

    # --- Page Layout ---
    main_col, side_col = st.columns([3.5, 1])

    with side_col:
        st.markdown("##### 🔔 WHAT'S NEXT")
        with st.container(border=True):
            if next_event_nav and st.button(f"🗓️ {next_event_str}", use_container_width=True, key="dash_nav_event"):
                st.session_state['trigger_navigation'] = next_event_nav; st.rerun()
            if incomplete_todos > 0 and st.button(f"📝 **{incomplete_todos}** pending task(s)", use_container_width=True, key="dash_nav_todo"):
                st.session_state['trigger_navigation'] = "👥 COACHING"; st.rerun()
            if not next_event_nav and incomplete_todos == 0:
                st.caption("All clear! ✨")

    # Initialize df_filt
    df_filt = pd.DataFrame()

    with main_col:
        # --- Statistiken und Charts ---
        if not df.empty:
            min_date = df['DateObj'].min() if pd.notna(df['DateObj'].min()) else datetime(2024,1,1)
            c1, c2 = st.columns([1,3])
            with c1: start_d = st.date_input("Stats ab:", min_date)
            df_filt = df[df['DateObj'] >= pd.Timestamp(start_d)].copy()
        else:
            st.info("No match data available. Please import matches in the 'Match Entry' tab.")

    if not df.empty:
        if not df_filt.empty:
            tab_overview, tab_team_stats = st.tabs(["📊 OVERVIEW", "📈 TEAM STATS"])

            with tab_overview:
                # --- CONFIDENCE SCALE ---
                st.divider(); st.markdown("### 📊 MAP CONFIDENCE")
                all_maps = sorted(df_filt['Map'].unique())
                sel_maps = st.multiselect("Wähle Maps:", all_maps, default=all_maps)
                conf_list = []
                for m in (sel_maps if sel_maps else all_maps):
                    md = df_filt[df_filt['Map']==m]
                    if md.empty: continue
                    g=len(md); w=len(md[md['Result']=='W']); delta=md['Delta'].sum()
                    wr=w/g*100; score=(wr*1.0)+(g*2.0)+(delta*0.5)
                    # DASHBOARD: nutzt 'list' (Banner)
                    img_p = get_map_img(m, type='list'); 
                    b64 = img_to_b64(img_p)
                    col = "#00ff80" if score>=60 else "#ffeb3b" if score>=40 else "#ff1493"
                    conf_list.append({'M':m, 'S':score, 'WR':wr, 'I':b64, 'C':col})
                
                conf_list.sort(key=lambda x: x['S'], reverse=True)
                html = "<div class='conf-scroll-wrapper'>"
                for c in conf_list:
                    img_tag = f"<img src='data:image/png;base64,{c['I']}'>" if c['I'] else ""
                    html += f"<div class='conf-card' style='border-bottom: 4px solid {c['C']}'><div class='conf-img-container'>{img_tag}</div><div class='conf-body'><div style='font-weight:bold; color:white'>{c['M']}</div><div class='conf-val' style='color:{c['C']}'>{c['S']:.1f}</div><div class='conf-sub'>{c['WR']:.0f}% WR</div></div></div>"
                st.markdown(html+"</div>", unsafe_allow_html=True)
                
                # --- TABLE ---
                # --- POWER RANKING (CUSTOM UI) ---
                st.divider()
                st.markdown("### 🏆 POWER RANKING")
                
                # Daten vorbereiten (wie vorher)
                rank_df = pd.DataFrame(conf_list).rename(columns={'M':'Map','S':'Score','WR':'Winrate'})
                
                if not rank_df.empty:
                    # Max Score finden für die Berechnung der Balkenlänge (100% Breite)
                    max_score = rank_df['Score'].max() if rank_df['Score'].max() > 0 else 1
                    
                    # Header (Optional, aber schön für Übersicht)
                    # st.caption("MAP PERFORMANCE BREAKDOWN")

                    for idx, row in rank_df.iterrows():
                        # 1. Daten holen
                        map_name = row['Map']
                        score = row['Score']
                        winrate = row['Winrate']
                        
                        # 2. Bild holen (als Base64)
                        # Wir nutzen hier 'list' (breit) oder 'icon' (quadratisch) - was du lieber magst
                        img_path = get_map_img(map_name, 'list') 
                        b64_img = img_to_b64(img_path)
                        img_html = f"<img src='data:image/png;base64,{b64_img}' class='rank-img'>" if b64_img else ""
                        
                        # 3. Balkenbreiten berechnen (in %)
                        # Score relativ zum besten Score
                        score_width = min(int((score / max_score) * 100), 100)
                        # Winrate ist einfach der Prozentwert
                        winrate_width = int(winrate)
                        
                        # 4. Farben für Balken definieren
                        # Score: Cyan (#00BFFF)
                        # Winrate: Pink (#FF1493) oder dynamisch (Grün/Rot)
                        wr_color = "#00ff80" if winrate >= 50 else "#ff4655" # Grün wenn >50%, sonst Rot
                        
                        # 5. HTML zusammenbauen
                        html_rank = f"""
                        <div class="rank-row">
                            <div class="rank-img-box">
                                {img_html}
                            </div>
                            
                            <div class="rank-name">{map_name}</div>
                            
                            <div class="rank-stats">
                                <div class="stat-line">
                                    <div class="stat-label">RATING</div>
                                    <div class="prog-bg">
                                        <div class="prog-fill" style="width: {score_width}%; background-color: #00BFFF; box-shadow: 0 0 10px rgba(0, 191, 255, 0.4);"></div>
                                    </div>
                                    <div class="stat-val" style="color: #00BFFF;">{score:.1f}</div>
                                </div>
                                
                                <div class="stat-line">
                                    <div class="stat-label">WIN %</div>
                                    <div class="prog-bg">
                                        <div class="prog-fill" style="width: {winrate_width}%; background-color: {wr_color};"></div>
                                    </div>
                                    <div class="stat-val" style="color: {wr_color};">{winrate:.0f}%</div>
                                </div>
                            </div>
                        </div>
                        """
                        
                        # Rendern ohne Zeilenumbrüche
                        st.markdown(html_rank.replace("\n", " "), unsafe_allow_html=True)

                else:
                    st.info("Noch keine Daten für das Ranking verfügbar.")
                
                # --- RECENT ---
                st.divider()
                st.markdown("### 📜 RECENT MATCHES")
                matches_to_show = df_filt.sort_values('DateObj', ascending=False).head(5)
                
                for idx, row in matches_to_show.iterrows():
                    
                    # 1. DATEN VORBEREITEN
                    res = row['Result']
                    score_us = int(row['Score_Us'])
                    score_en = int(row['Score_Enemy'])
                    
                    # Farben setzen
                    if res == 'W':
                        main_color = "#00ff80" 
                    elif res == 'L':
                        main_color = "#ff4655"
                    else:
                        main_color = "#aaaaaa"

                    # Map Bild laden
                    map_img_path = get_map_img(row['Map'], 'list')
                    b64_map = img_to_b64(map_img_path)
                    map_bg_style = f"background-image: linear-gradient(to right, rgba(0,0,0,0) 0%, rgba(18,18,18,1) 100%), url('data:image/png;base64,{b64_map}');" if b64_map else "background-color: #222;"

                    # 2. STATS EXTRAHIEREN
                    atk_w = int(row.get('Atk_R_W', 0)) if pd.notna(row.get('Atk_R_W')) else 0
                    atk_l = int(row.get('Atk_R_L', 0)) if pd.notna(row.get('Atk_R_L')) else 0
                    def_w = int(row.get('Def_R_W', 0)) if pd.notna(row.get('Def_R_W')) else 0
                    def_l = int(row.get('Def_R_L', 0)) if pd.notna(row.get('Def_R_L')) else 0
                    date_str = row['Date']

                    # 3. AGENTEN ICONS GENERIEREN
                    def get_agent_row_html(comp_prefix):
                        html = ""
                        for i in range(1, 6):
                            an = row.get(f'{comp_prefix}_{i}')
                            if pd.notna(an) and str(an) != "":
                                b64 = img_to_b64(get_agent_img(an))
                                if b64:
                                    html += f"<img src='data:image/png;base64,{b64}' class='val-agent-img' title='{an}'>"
                                else:
                                    html += f"<div class='val-agent-img' style='background:#333' title='?'></div>"
                            else:
                                    html += f"<div class='val-agent-img' style='background:transparent; border:1px dashed #333'></div>"
                        return html

                    my_agents = get_agent_row_html('MyComp')
                    en_agents = get_agent_row_html('EnComp')

                    # 4. VOD LINK
                    vod_link = row.get('VOD_Link')
                    vod_html = ""
                    if pd.notna(vod_link) and str(vod_link).startswith("http"):
                        vod_html = f'<a href="{vod_link}" target="_blank" class="val-vod-link" style="color: {main_color};">WATCH VOD ↗</a>'

                    # 5. HTML ZUSAMMENBAUEN
                    html_card = f"""
                    <div class="val-card">
                        <div class="val-bar" style="background-color: {main_color};"></div>
                        <div class="val-map-section">
                            <div class="val-map-bg" style="{map_bg_style}"></div>
                            <div class="val-map-text">
                                <div class="val-map-name">{row['Map']}</div>
                            </div>
                        </div>
                        <div class="val-comps-section">
                            <div class="val-agent-row">
                                <div class="val-team-label" style="color:#00ff80">US</div>
                                {my_agents}
                            </div>
                            <div class="val-agent-row">
                                <div class="val-team-label" style="color:#ff4655">EN</div>
                                {en_agents}
                            </div>
                        </div>
                        <div class="val-stats-section">
                            <div class="stat-group">
                                <div class="stat-label">ATTACK</div>
                                <div class="stat-value" style="color:#ff99aa">{atk_w} - {atk_l}</div>
                            </div>
                            <div class="stat-group">
                                <div class="stat-label">DEFENSE</div>
                                <div class="stat-value" style="color:#99ffcc">{def_w} - {def_l}</div>
                            </div>
                            <div class="stat-group">
                                <div class="stat-label">DATE</div>
                                <div class="stat-date">{date_str}</div>
                            </div>
                        </div>
                        <div class="val-score-section">
                            <div class="val-score" style="color: {main_color};">{score_us} : {score_en}</div>
                            {vod_html}
                        </div>
                    </div>
                    """
                    
                    st.markdown(html_card.replace("\n", " "), unsafe_allow_html=True)

            with tab_team_stats:
                st.markdown("### 📈 TEAM STATISTICS")
                
                # KPIs
                total_matches = len(df_filt)
                wins = len(df_filt[df_filt['Result'] == 'W'])
                winrate = (wins / total_matches) * 100 if total_matches > 0 else 0
                
                atk_w = df_filt['Atk_R_W'].sum()
                atk_l = df_filt['Atk_R_L'].sum()
                atk_wr = (atk_w / (atk_w + atk_l) * 100) if (atk_w + atk_l) > 0 else 0
                
                def_w = df_filt['Def_R_W'].sum()
                def_l = df_filt['Def_R_L'].sum()
                def_wr = (def_w / (def_w + def_l) * 100) if (def_w + def_l) > 0 else 0
                
                k1, k2, k3, k4 = st.columns(4)
                k1.metric("Matches", total_matches)
                k2.metric("Winrate", f"{winrate:.1f}%")
                k3.metric("Attack Win%", f"{atk_wr:.1f}%")
                k4.metric("Defense Win%", f"{def_wr:.1f}%")
                
                st.divider()
                
                st.markdown("#### 📉 Winrate Trend")
                if not df_filt.empty:
                    df_trend = df_filt.sort_values('DateObj').copy()
                    df_trend['WinNum'] = df_trend['Result'].apply(lambda x: 1 if x == 'W' else 0)
                    df_trend['Cumulative Winrate'] = df_trend['WinNum'].expanding().mean() * 100
                    df_trend['Winrate (Last 5 Matches)'] = df_trend['WinNum'].rolling(window=5, min_periods=1).mean() * 100
                    
                    # Melt for plotly to show both lines
                    df_melt = df_trend.melt(id_vars=['DateObj', 'MatchID'], value_vars=['Cumulative Winrate', 'Winrate (Last 5 Matches)'], 
                                            var_name='Metric', value_name='Winrate')
                    
                    fig_trend = px.line(df_melt, x='DateObj', y='Winrate', color='Metric',
                                        markers=True,
                                        title="Winrate Trend (Cumulative vs Form)",
                                        color_discrete_map={'Cumulative Winrate': '#00BFFF', 'Winrate (Last 5 Matches)': '#FF1493'})
                    
                    fig_trend.update_layout(
                        xaxis_title=None, yaxis_title="Winrate %",
                        yaxis=dict(range=[0, 110]),
                        paper_bgcolor='rgba(0,0,0,0)', plot_bgcolor='rgba(255,255,255,0.05)',
                        font=dict(color='white'),
                        legend=dict(orientation="h", y=1.1, x=0, title=None),
                        margin=dict(l=0, r=0, t=20, b=0),
                        hovermode="x unified"
                    )
                    st.plotly_chart(fig_trend, use_container_width=True)
                
                st.divider()

                # --- Map Stats Calculation (moved up) ---
                map_stats = pd.DataFrame()
                if not df_filt.empty:
                    map_stats = df_filt.groupby('Map').agg({
                        'Result': lambda x: (x == 'W').sum(),
                        'Date': 'count', # Use Date to count matches (safer than MatchID)
                        'Atk_R_W': 'sum', 'Atk_R_L': 'sum',
                        'Def_R_W': 'sum', 'Def_R_L': 'sum'
                    }).rename(columns={'Date': 'Played', 'Result': 'Wins'})
                    
                    map_stats['Win%'] = (map_stats['Wins'] / map_stats['Played']) * 100
                    map_stats['Atk%'] = (map_stats['Atk_R_W'] / (map_stats['Atk_R_W'] + map_stats['Atk_R_L']).replace(0,1)) * 100
                    map_stats['Def%'] = (map_stats['Def_R_W'] / (map_stats['Def_R_W'] + map_stats['Def_R_L']).replace(0,1)) * 100
                    
                    map_stats = map_stats.sort_values('Win%', ascending=False)

                # --- NEW: Map Winrate Chart ---
                st.markdown("#### 🗺️ Map Winrate Comparison")
                if not map_stats.empty:
                    map_stats['Color'] = map_stats['Win%'].apply(lambda wr: '#00ff80' if wr >= 55 else ('#ffeb3b' if wr >= 45 else '#ff4655'))
                    
                    fig_map_wr = px.bar(map_stats, x=map_stats.index, y='Win%',
                                        text_auto='.2s',
                                        hover_data={'Played': True, 'Wins': True})
                    fig_map_wr.update_traces(marker_color=map_stats['Color'], textposition='outside')
                    fig_map_wr.update_layout(
                        yaxis_title="Winrate %", xaxis_title=None,
                        yaxis=dict(range=[0, 110]),
                        paper_bgcolor='rgba(0,0,0,0)', plot_bgcolor='rgba(255,255,255,0.05)',
                        font=dict(color='white')
                    )
                    st.plotly_chart(fig_map_wr, use_container_width=True)
                
                st.divider()
                
                c1, c2 = st.columns([2, 1])
                
                with c1:
                    st.markdown("#### 🗺️ Map Performance Details")
                    if not map_stats.empty:
                        map_stats_table = map_stats.sort_values('Played', ascending=False)

                        # Add Icons & Reset Index
                        # FIX: Nutze '_list.png' (Banner) statt Icon
                        map_stats_table['Icon'] = map_stats_table.index.map(lambda x: f"data:image/png;base64,{img_to_b64(get_map_img(x, 'list'))}")
                        map_stats_table = map_stats_table.reset_index()
                        
                        st.dataframe(
                            map_stats_table[['Icon', 'Map', 'Played', 'Wins', 'Win%', 'Atk%', 'Def%']],
                            use_container_width=True,
                            hide_index=True,
                            column_config={
                                "Icon": st.column_config.ImageColumn("", width="small"),
                                "Win%": st.column_config.ProgressColumn("Win%", format="%.1f%%", min_value=0, max_value=100),
                                "Atk%": st.column_config.ProgressColumn("Atk%", format="%.1f%%", min_value=0, max_value=100),
                                "Def%": st.column_config.ProgressColumn("Def%", format="%.1f%%", min_value=0, max_value=100),
                            }
                        )
                    else:
                        st.info("No data.")

                with c2:
                    st.markdown("#### ♟️ Agent Pick Rates")
                    all_agents = []
                    for i in range(1, 6):
                        if f'MyComp_{i}' in df_filt.columns:
                            all_agents.extend(df_filt[f'MyComp_{i}'].dropna().tolist())
                    
                    all_agents = [a for a in all_agents if a and str(a).strip() != ""]
                    
                    if all_agents:
                        agent_counts = pd.Series(all_agents).value_counts().reset_index()
                        agent_counts.columns = ['Agent', 'Count']
                        agent_counts['Pick%'] = (agent_counts['Count'] / total_matches) * 100
                        
                        # Add Icons
                        agent_counts['Icon'] = agent_counts['Agent'].map(lambda x: f"data:image/png;base64,{img_to_b64(get_agent_img(x))}")
                        
                        st.dataframe(
                            agent_counts[['Icon', 'Agent', 'Count', 'Pick%']],
                            use_container_width=True,
                            hide_index=True,
                            column_config={
                                "Icon": st.column_config.ImageColumn("", width="small"),
                                "Pick%": st.column_config.ProgressColumn("Pick Rate", format="%.0f%%", min_value=0, max_value=100)
                            }
                        )
                    else:
                        st.info("No agent data.")

        else:
            st.info("No matches found for the selected date range.")

# ==============================================================================
# 👥 COACHING
# ==============================================================================
elif page == "👥 COACHING":
    st.title("👥 PLAYER COACHING")

    # Get current user info
    current_user = st.session_state.get('username', '')
    user_role = st.session_state.get('role', '')
    
    # Testing role context switcher
    if user_role == 'testing':
        st.markdown("---")
        col1, col2 = st.columns(2)
        with col1:
            if st.button("🔄 Switch to Coach View", type="primary"):
                st.session_state.testing_context = 'coach'
                st.rerun()
        with col2:
            if st.button("🔄 Switch to Player View", type="secondary"):
                st.session_state.testing_context = 'player'
                st.rerun()
        
        # Show current context
        current_context = st.session_state.get('testing_context', 'coach')
        st.info(f"🧪 **Testing Mode**: Currently viewing as **{current_context.upper()}**")
        user_role = current_context  # Override role for the rest of the page
        st.markdown("---")
    
    if user_role == 'coach':
        # Coach view - manage all players
        tab1, tab3 = st.tabs(["📝 Assign Todos", "📊 Player Overview"])
        
        with tab1:
            st.markdown("### 📝 Assign Tasks to Players")
            
            with st.form("assign_todo"):
                col1, col2 = st.columns(2)
                
                with col1:
                    player = st.selectbox("Select Player", 
                                            ["Luggi","Benni","Andrei","Luca","Sofi","Remus"],
                                            format_func=lambda x: f"🎮 {x}", key="todo_player")
                    title = st.text_input("Task Title", placeholder="e.g., Review Ascent Defense", key="todo_title")
                
                with col2:
                    # Playbook selection
                    all_playbooks = []
                    if not df_legacy_pb.empty:
                        legacy_pbs = [f"Legacy: {pb}" for pb in df_legacy_pb.get('Name', [])]
                        all_playbooks.extend(legacy_pbs)
                    if not df_team_pb.empty:
                        team_pbs = [f"Team: {pb}" for pb in df_team_pb.get('Name', [])]
                        all_playbooks.extend(team_pbs)
                    
                    playbook_link = st.selectbox("Link Playbook (optional)", 
                                                ["None"] + all_playbooks,
                                                help="Select a playbook to link to this task", key="todo_pb_link")
                    
                    youtube_link = st.text_input("YouTube Link (optional)", 
                                                placeholder="https://youtube.com/watch?v=...", key="todo_yt")
                
                description = st.text_area("Task Description", 
                                            placeholder="Detailed instructions for the player...", key="todo_desc")
                
                submitted = st.form_submit_button("📤 Assign Task")
                
                if submitted:
                    if not title or not description:
                        st.error("Please fill in title and description.")
                    else:
                        # Create todo
                        todo_id = str(uuid.uuid4())[:8]
                        
                        # Process playbook link
                        final_pb_link = ""
                        if playbook_link != "None":
                            if playbook_link.startswith("Legacy: "):
                                pb_name = playbook_link.replace("Legacy: ", "")
                                final_pb_link = f"Legacy Playbook: {pb_name}"
                            elif playbook_link.startswith("Team: "):
                                pb_name = playbook_link.replace("Team: ", "")
                                final_pb_link = f"Team Playbook: {pb_name}"
                        
                        new_todo = {
                            'ID': todo_id,
                            'Player': player,
                            'Title': title,
                            'Description': description,
                            'PlaybookLink': final_pb_link,
                            'YoutubeLink': youtube_link if youtube_link else "",
                            'AssignedBy': current_user,
                            'AssignedAt': datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
                            'Completed': False,
                            'CompletedAt': ""
                        }
                        
                        db_insert("player_todos", new_todo, "df_todos")
                        
                        # DISCORD NOTIFICATION
                        send_discord_notification(player, title, description)
                        
                        # Manual Clear
                        for k in ["todo_title", "todo_yt", "todo_desc"]:
                            if k in st.session_state: del st.session_state[k]
                        
                        st.success(f"Task '{title}' assigned to {player}!")
                        st.rerun()
        
        with tab3:
            st.markdown("### 📊 Player Overview")
            
            players =  ["Luggi","Benni","Andrei","Luca","Sofi","Remus"]
            
            for player in players:
                with st.expander(f"🎮 {player}", expanded=False):
                    col1, col2, col3 = st.columns(3)
                    
                    # Task stats
                    player_todos = df_todos[df_todos['Player'] == player]
                    completed_tasks = len(player_todos[player_todos['Completed'] == True])
                    total_tasks = len(player_todos)
                    
                    with col1:
                        st.metric("Tasks Completed", f"{completed_tasks}/{total_tasks}")
                        if total_tasks > 0:
                            completion_rate = int((completed_tasks / total_tasks) * 100)
                            st.progress(completion_rate / 100)
                            st.caption(f"{completion_rate}% completion rate")
                    
                    # Scrim availability
                    player_availability = df_availability[df_availability['Player'] == player]
                    available_count = len(player_availability[player_availability['Available'] == 'Yes'])
                    total_scrims = len(df_scrims) if not df_scrims.empty else 0
                    
                    with col3:
                        if total_scrims > 0:
                            availability_rate = int((available_count / total_scrims) * 100)
                            st.metric("Scrim Availability", f"{availability_rate}%")
                            st.progress(availability_rate / 100)
                        else:
                            st.metric("Scrim Availability", "N/A")
    
    elif user_role == 'player':
        # Player view - see their own tasks and messages
        
        # Check for notifications
        player_todos = df_todos[(df_todos['Player'] == current_user) & (df_todos['Completed'] == False)]
        incomplete_todos = len(player_todos)
        
        # Note: df_legacy_pb and df_team_pb are already loaded globally from GSheets.
        
        with st.container():
            st.markdown("### 📋 My Assigned Tasks")
            
            if df_todos.empty:
                st.info("No tasks assigned yet.")
            else:
                player_todos = df_todos[df_todos['Player'] == current_user]
                
                if player_todos.empty:
                    st.info("No tasks assigned to you yet.")
                else:
                    for _, todo in player_todos.iterrows():
                        completed = todo['Completed']
                        status_icon = "✅" if completed else "⏳"
                        
                        with st.expander(f"{status_icon} {todo['Title']}", expanded=not completed):
                            st.markdown(f"**Assigned by:** {todo['AssignedBy']}")
                            st.markdown(f"**Assigned:** {todo['AssignedAt']}")
                            st.markdown(f"**Description:** {todo['Description']}")
                            
                            if todo['PlaybookLink']:
                                if todo['PlaybookLink'].startswith("Team Playbook: "):
                                    pb_name = todo['PlaybookLink'].replace("Team Playbook: ", "")
                                    if not df_team_pb.empty:
                                        matching_pb = df_team_pb[df_team_pb['Name'] == pb_name]
                                        if not matching_pb.empty:
                                            pb_id = matching_pb.iloc[0]['ID']
                                            if st.button(f"📖 Open Playbook: {pb_name}", key=f"open_pb_{todo['ID']}"):
                                                st.session_state["trigger_navigation"] = "📘 STRATEGY BOARD"
                                                st.session_state['sel_pb_id'] = pb_id
                                                st.rerun()
                                        else:
                                            st.markdown(f"**📖 Playbook:** {todo['PlaybookLink']} (not found)")
                                    else:
                                        st.markdown(f"**📖 Playbook:** {todo['PlaybookLink']}")
                                elif todo['PlaybookLink'].startswith("Legacy Playbook: "):
                                    pb_name = todo['PlaybookLink'].replace("Legacy Playbook: ", "")
                                    if not df_legacy_pb.empty:
                                        matching_pb = df_legacy_pb[df_legacy_pb['Name'] == pb_name]
                                        if not matching_pb.empty:
                                            pb_link = matching_pb.iloc[0].get('Link', '')
                                            if pb_link:
                                                st.markdown(f"**📖 Playbook:** [{pb_name}]({pb_link})")
                                            else:
                                                st.markdown(f"**📖 Playbook:** {pb_name}")
                                        else:
                                            st.markdown(f"**📖 Playbook:** {todo['PlaybookLink']} (not found)")
                                    else:
                                        st.markdown(f"**📖 Playbook:** {todo['PlaybookLink']}")
                                else:
                                    st.markdown(f"**📖 Playbook:** {todo['PlaybookLink']}")
                            
                            if todo['YoutubeLink']:
                                st.markdown(f"**📺 YouTube:** [{todo['YoutubeLink']}]({todo['YoutubeLink']})")
                            
                            if not completed:
                                if st.button("✅ Mark as Completed (Delete)", key=f"complete_{todo['ID']}"):
                                    # Delete task instead of just marking it
                                    df_todos = df_todos[df_todos['ID'] != todo['ID']]
                                    save_player_todos(df_todos)
                                    db_delete("player_todos", todo['ID'])
                                    st.success("Task completed and removed!")
                                    st.rerun()
                            else:
                                st.success(f"✅ Completed on {todo['CompletedAt']}")
                                if st.button("🗑️ Delete", key=f"del_comp_{todo['ID']}"):
                                    df_todos = df_todos[df_todos['ID'] != todo['ID']]
                                    save_player_todos(df_todos)
                                    db_delete("player_todos", todo['ID'])
                                    st.rerun()

# ==============================================================================
# ⚽ SCRIMS
# ==============================================================================
elif page == "⚽ SCRIMS":
    st.title("⚽ SCRIM SCHEDULER")
    
    current_user = st.session_state.get('username', '')
    user_role = st.session_state.get('role', '')
    
    # --- NAVIGATION (Radio instead of Tabs for programmatic redirect) ---
    if 'scrim_nav' not in st.session_state: st.session_state.scrim_nav = "📅 View Scrims"
    
    nav_opts = ["📅 View Scrims"]
    if user_role == 'coach': nav_opts.append("➕ Create Scrim")
    
    # Ensure valid state
    if st.session_state.scrim_nav not in nav_opts: st.session_state.scrim_nav = nav_opts[0]
    
    st.session_state.scrim_nav = st.radio("Scrim Nav", nav_opts, index=nav_opts.index(st.session_state.scrim_nav), horizontal=True, label_visibility="collapsed", key="scrim_nav_radio")
    
    if st.session_state.scrim_nav == "📅 View Scrims":
        st.markdown("### 📅 Upcoming Scrims")
        if df_scrims.empty:
            st.info("No scrims scheduled yet.")
        else:
            df_scrims_sorted = df_scrims[df_scrims['DateTimeObj'] >= pd.Timestamp.now()].sort_values('DateTimeObj', ascending=True)
            if df_scrims_sorted.empty:
                st.info("No upcoming scrims. Check the 'Create Scrim' tab to schedule one.")
            
            for _, scrim in df_scrims_sorted.iterrows():
                scrim_id = scrim['ID']
                with st.container(border=True):
                    c1, c2 = st.columns([1, 2.5])
                    with c1:
                        map_img = get_map_img(scrim.get('Map'), 'list')
                        if map_img: st.image(map_img)
                        else: st.markdown(f"**🗺️ {scrim.get('Map', 'N/A')}**")
                    
                    with c2:
                        st.markdown(f"#### {scrim['Title']}")
                        st.caption(f"🗓️ {scrim['Date']} at {scrim['Time']} | Created by {scrim['CreatedBy']}")
                        if pd.notna(scrim.get('Description')) and scrim['Description']:
                            st.markdown(f"> _{scrim['Description']}_")
                        
                        # Links
                        l1, l2 = st.columns(2)
                        pb_link = scrim.get('PlaybookLink')
                        if pd.notna(pb_link) and pb_link and pb_link != "None":
                            if st.button(f"📖 Open Playbook", key=f"pb_{scrim_id}", use_container_width=True):
                                if pb_link.startswith("Team: "):
                                    pb_name = pb_link.replace("Team: ", "")
                                    pb_id = df_team_pb[df_team_pb['Name'] == pb_name].iloc[0]['ID']
                                    st.session_state['sel_pb_id'] = pb_id
                                    st.session_state['trigger_navigation'] = "📘 STRATEGY BOARD"
                                    st.rerun()
                        
                        vid_link = scrim.get('VideoLink')
                        if pd.notna(vid_link) and vid_link:
                            l2.link_button("📺 Watch Video", vid_link, use_container_width=True)

                        st.markdown("---")
                        
                        # Availability
                        avail_data = df_availability[df_availability['ScrimID'] == scrim_id]
                        available_players = avail_data[avail_data['Available'] == 'Yes']['Player'].tolist()
                        
                        st.markdown(f"**Availability ({len(available_players)}/{len(OUR_TEAM)})**")
                        st.progress(len(available_players) / len(OUR_TEAM) if OUR_TEAM else 0, text=f"{', '.join(available_players) if available_players else 'None'}")

                        # Player actions
                        if user_role != 'coach':
                            st.caption("Your Status:")
                            b1, b2, b3 = st.columns(3)
                            if b1.button("✅ Yes", key=f"yes_{scrim_id}", use_container_width=True): update_availability(scrim_id, current_user, "Yes"); st.rerun()
                            if b2.button("🤔 Maybe", key=f"maybe_{scrim_id}", use_container_width=True): update_availability(scrim_id, current_user, "Maybe"); st.rerun()
                            if b3.button("❌ No", key=f"no_{scrim_id}", use_container_width=True): update_availability(scrim_id, current_user, "No"); st.rerun()
                        
                        # Coach actions
                        if user_role == 'coach':
                            if st.button("🗑️ Delete Scrim", key=f"del_{scrim_id}", use_container_width=True):
                                db_delete("scrims", scrim_id, "df_scrims")
                                # Availability cleanup is handled by DB constraints usually, or we leave orphans for now
                                st.rerun()

    elif st.session_state.scrim_nav == "➕ Create Scrim":
        st.markdown("### ➕ Create New Scrim")
        with st.form("create_scrim"):
            c1, c2 = st.columns(2)
            with c1:
                title = st.text_input("Scrim Title", placeholder="e.g., Weekly Scrim vs Team X", key="scrim_title")
                date = st.date_input("Date", min_value=datetime.today().date(), key="scrim_date")
                time = st.time_input("Time", key="scrim_time")
            with c2:
                map_name = st.selectbox("Map", sorted(df['Map'].unique()) if not df.empty else ["Ascent"], key="scrim_map")
                description = st.text_area("Description", placeholder="e.g., Focus on B-Site retakes", key="scrim_desc")
            
            st.markdown("---")
            st.markdown("##### Optional Links")
            c3, c4 = st.columns(2)
            with c3:
                all_playbooks = ["None"]
                if not df_team_pb.empty: all_playbooks.extend([f"Team: {pb}" for pb in df_team_pb.get('Name', [])])
                if not df_legacy_pb.empty: all_playbooks.extend([f"Legacy: {pb}" for pb in df_legacy_pb.get('Name', [])])
                playbook_link = st.selectbox("Link Playbook", all_playbooks, key="scrim_pb")
            with c4:
                video_link = st.text_input("Video Link (VOD, YouTube, etc.)", key="scrim_vid")

            if st.form_submit_button("Create Scrim", type="primary", use_container_width=True):
                if not title: st.error("Please enter a scrim title.")
                else:
                    scrim_id = str(uuid.uuid4())[:8]
                    new_scrim = {
                        'ID': scrim_id, 'Title': title, 'Date': date.strftime("%Y-%m-%d"),
                        'Time': time.strftime("%H:%M"), 'Map': map_name, 'Description': description,
                        'CreatedBy': current_user, 'CreatedAt': datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
                        'PlaybookLink': playbook_link if playbook_link != "None" else "",
                        'VideoLink': video_link
                    }
                    db_insert("scrims", new_scrim, "df_scrims")
                    
                    # Manual Clear & Redirect
                    for k in ["scrim_title", "scrim_desc", "scrim_vid"]:
                        if k in st.session_state: del st.session_state[k]
                    
                    st.session_state.scrim_nav = "📅 View Scrims"
                    st.success(f"Scrim '{title}' created!"); st.rerun()

# ==============================================================================
# 2. MATCH ENTRY (AUTO PLAYER STATS)
# ==============================================================================
elif page == "📝 MATCH ENTRY":
    st.header("Match Import")
    if 'fd' not in st.session_state: st.session_state['fd'] = {'d':datetime.today(), 'm':'Ascent', 'r':'W', 'us':13, 'en':8, 'mid':'', 'vod':'', 'my':[""]*5, 'en':[""]*5, 'hm':None, 'p_stats':[]}

    with st.expander("📂 JSON Import (Auto-Stats)", expanded=True):
        up = st.file_uploader("Tracker JSON", type=['json'], key="json")
        if up:
            try:
                data = json.load(up); meta = data['data']['metadata']; segs = data['data']['segments']
                mid = meta.get('matchId', f"J_{int(datetime.now().timestamp())}")
                mid = meta.get('matchId') or data['data'].get('attributes', {}).get('id') or f"J_{int(datetime.now().timestamp())}"
                
                my_tid = None
                # Team ID
                for s in segs:
                    if s['type']=='player-summary' and any(t.lower() in s['attributes']['platformUserIdentifier'].lower() for t in OUR_TEAM):
                        my_tid = s['metadata']['teamId']; break
                
                if my_tid:
                    rw=0; rl=0; my_c=[]; en_c=[]; heatmap=[]; p_stats=[]
                    
                    for s in segs:
                        if s['type']=='round-summary':
                            w = s.get('attributes',{}).get('winningTeamId') or s.get('metadata',{}).get('winningTeamId')
                            if w: (rw:=rw+1) if w==my_tid else (rl:=rl+1)
                            w = s.get('attributes',{}).get('winningTeamId') or \
                                s.get('metadata',{}).get('winningTeamId') or \
                                s.get('stats',{}).get('winningTeam',{}).get('value')
                            if w: 
                                if w==my_tid: rw+=1 
                                else: rl+=1
                        
                        if s['type']=='player-summary':
                            ag = s['metadata']['agentName']; name = s['attributes']['platformUserIdentifier']
                            if s['metadata']['teamId']==my_tid: my_c.append(ag)
                            else: en_c.append(ag)
                            
                            # EXTRACT STATS
                            if any(t.lower() in name.lower() for t in OUR_TEAM):
                                sts = s['stats']
                                rounds = rw + rl if (rw+rl) > 0 else 1
                                p_stats.append({
                                    'MatchID': mid, 'Date': datetime.today().strftime("%d.%m.%Y"),
                                    'Map': meta.get('mapName'), 'Player': name.split('#')[0], 'Agent': ag,
                                    'Kills': sts.get('kills', {}).get('value', 0),
                                    'Deaths': sts.get('deaths', {}).get('value', 0),
                                    'Assists': sts.get('assists', {}).get('value', 0),
                                    'Score': sts.get('score', {}).get('value', 0), 'Rounds': rounds,
                                    'HS': sts.get('hsAccuracy', {}).get('value', 0)
                                })

                        if s['type']=='player-round-kills':
                            v=s['attributes']['opponentPlatformUserIdentifier']; k=s['attributes']['platformUserIdentifier']
                            w=s['metadata']['weaponName']; r=s['attributes'].get('round',0)
                            if any(t.lower() in v.lower() for t in OUR_TEAM) and 'opponentLocation' in s['metadata']:
                                heatmap.append({'X':s['metadata']['opponentLocation']['x'], 'Y':s['metadata']['opponentLocation']['y'], 'Player':v, 'Type':'Death', 'Weapon':w, 'Round':r})
                            if any(t.lower() in k.lower() for t in OUR_TEAM) and 'location' in s['metadata']:
                                heatmap.append({'X':s['metadata']['location']['x'], 'Y':s['metadata']['location']['y'], 'Player':k, 'Type':'Kill', 'Weapon':w, 'Round':r})

                    st.session_state['fd'].update({'m':meta.get('mapName','Ascent'), 'us':int(rw), 'en':int(rl), 'r':'W' if rw>rl else 'L' if rl>rw else 'D', 'mid':mid, 'hm':heatmap, 'p_stats':p_stats})
                    while len(my_c)<5: my_c.append(""); 
                    while len(en_c)<5: en_c.append("")
                    st.session_state['fd']['my']=my_c[:5]; st.session_state['fd']['en']=en_c[:5]
                    st.success(f"✅ Loaded! Stats for {len(p_stats)} players ready.")
                else: st.error("Team not found")
            except Exception as e: st.error(str(e))

    d = st.session_state['fd']
    with st.form("e"):
        c1,c2=st.columns(2)
        with c1:
            maps=sorted(["Ascent","Bind","Haven","Split","Lotus","Sunset","Abyss","Pearl","Fracture","Icebox","Breeze","Corrode"])
            try: mx=maps.index(d['m'])
            except: mx=0
            dt=st.date_input("D",d['d']); mp=st.selectbox("M",maps,index=mx)
            ags = sorted([os.path.basename(x).replace(".png","").capitalize() for x in glob.glob(os.path.join(ASSET_DIR, "agents", "*.png"))])
            ac=st.columns(5); my_final=[]
            for i in range(5):
                pr=d['my'][i]; idx=0
                if pr: 
                    try: idx=ags.index(pr.capitalize())+1
                    except: pass
                my_final.append(ac[i].selectbox(f"A{i}",[""]+ags,index=idx))
        with c2:
            ro=["W","L","D"]; ri=ro.index(d['r'])
            re=st.radio("R",ro,index=ri,horizontal=True)
            su=int(d['us']) if isinstance(d['us'],(int,float)) else 0
            se=int(d['en']) if isinstance(d['en'],(int,float)) else 0
            us=st.number_input("U",0,25,su); en=st.number_input("E",0,25,se)
            mi=st.text_input("ID",d['mid']); vo=st.text_input("VOD",d['vod'])
            en_f=d['en']

        with st.expander("Details"):
            rc1,rc2=st.columns(2)
            aw=rc1.number_input("AW",0); al=rc1.number_input("AL",0)
            dw=rc2.number_input("DW",0); dl=rc2.number_input("DL",0)

        if st.form_submit_button("SAVE"):
            if d['hm']: pd.DataFrame(d['hm']).to_csv(os.path.join(BASE_DIR,"data",f"Locs_{mi}.csv"),index=False)
            row={'Date':dt.strftime("%d.%m.%Y"),'Map':mp,'Result':re,'Score_Us':us,'Score_Enemy':en,'MatchID':mi,'Source':'Tracker' if d['hm'] else 'Manual','VOD_Link':vo,'Atk_R_W':aw,'Atk_R_L':al,'Def_R_W':dw,'Def_R_L':dl}
            for i in range(5): row[f'MyComp_{i+1}']=my_final[i]; row[f'EnComp_{i+1}']=en_f[i] if i<len(en_f) else ""
            
            # Save Match (Single Row Insert)
            db_insert("nexus_matches", row, "df_matches")
            
            # Save Player Stats (Bulk Insert for new rows only)
            if d['p_stats']:
                try: supabase.table("Premier - PlayerStats").insert(d['p_stats']).execute()
                except Exception as e: st.error(f"Stats Error: {e}")
                
            # Reset local state to defaults so form doesn't repopulate with old data
            st.session_state['fd'] = {'d':datetime.today(), 'm':'Ascent', 'r':'W', 'us':13, 'en':8, 'mid':'', 'vod':'', 'my':[""]*5, 'en':[""]*5, 'hm':None, 'p_stats':[]}
            st.success("Saved!"); st.rerun()

# ==============================================================================
# 3. MAP ANALYZER
# ==============================================================================
elif page == "🗺️ MAP ANALYZER":
    st.title("TACTICAL BOARD")
    if not df.empty:
        # --- VISUAL MAP SELECTOR ---
        if 'ana_map' not in st.session_state: 
            st.session_state.ana_map = sorted(df['Map'].unique())[0] if not df.empty else "Ascent"
        
        with st.expander("🗺️ SELECT MAP", expanded=True):
            render_visual_selection(sorted(df['Map'].unique()), 'map', 'ana_sel', multi=False, key_state='ana_map')
        
        sel_map = st.session_state.ana_map
        m_df = df[df['Map'] == sel_map]
        head = get_map_img(sel_map, 'list'); 
        if head: st.image(head, use_container_width=True)
        
        wins = len(m_df[m_df['Result']=='W']); g = len(m_df)
        aw = m_df['Atk_R_W'].sum(); al = m_df['Atk_R_L'].sum()
        dw = m_df['Def_R_W'].sum(); dl = m_df['Def_R_L'].sum()
        wr = wins/g*100 if g>0 else 0
        awr = aw/(aw+al)*100 if (aw+al)>0 else 0
        dwr = dw/(dw+dl)*100 if (dw+dl)>0 else 0
        
        def get_col(v):
            if v >= 55: return "rgba(0, 255, 128, 0.15)", "#00ff80"
            elif v >= 45: return "rgba(255, 235, 59, 0.15)", "#ffeb3b"
            else: return "rgba(255, 20, 147, 0.15)", "#ff1493"

        c1, c2, c3 = st.columns(3)
        bg,bo=get_col(wr); c1.markdown(f"<div class='stat-box' style='background:{bg};border-color:{bo}'><div class='stat-val'>{wr:.1f}%</div><div class='stat-lbl'>WINRATE</div></div>", unsafe_allow_html=True)
        bg,bo=get_col(awr); c2.markdown(f"<div class='stat-box' style='background:{bg};border-color:{bo}'><div class='stat-val'>{awr:.1f}%</div><div class='stat-lbl'>ATK WIN%</div></div>", unsafe_allow_html=True)
        bg,bo=get_col(dwr); c3.markdown(f"<div class='stat-box' style='background:{bg};border-color:{bo}'><div class='stat-val'>{dwr:.1f}%</div><div class='stat-lbl'>DEF WIN%</div></div>", unsafe_allow_html=True)
        
        st.divider(); st.subheader("📍 HEATMAP ANALYSIS")
        
        match_lookup = {}
        for idx, row in m_df.iterrows():
            mid = str(row['MatchID'])
            if os.path.exists(os.path.join(BASE_DIR, "data", f"Locs_{mid}.csv")):
                match_lookup[f"{row['Date']} ({row['Result']})"] = mid
        
        if match_lookup:
            c_sel, c_opt, c_side = st.columns([2, 1, 1])
            with c_sel: sels = st.multiselect("Matches:", list(match_lookup.keys()), default=list(match_lookup.keys())[:1])
            with c_opt: metric = st.radio("Metric:", ["Deaths", "Kills"], horizontal=True)
            with c_side: side_filter = st.selectbox("Side:", ["All", "First Half", "Second Half"])
            
            if sels:
                dfs = []
                for l in sels:
                    try: dfs.append(pd.read_csv(os.path.join(BASE_DIR,"data",f"Locs_{match_lookup[l]}.csv")))
                    except: pass
                if dfs:
                    dfh = pd.concat(dfs, ignore_index=True)
                    if 'Type' not in dfh.columns: dfh['Type'] = 'Death'
                    target = 'Kill' if "Kills" in metric else 'Death'
                    dfh = dfh[dfh['Type'] == target]
                    if 'Round' in dfh.columns:
                        if side_filter == "First Half": dfh = dfh[dfh['Round'] <= 12]
                        elif side_filter == "Second Half": dfh = dfh[dfh['Round'] > 12]
                    
                    col_p = 'Player' if 'Player' in dfh.columns else 'Victim'
                    players = sorted(dfh[col_p].astype(str).unique())
                    sel_p = st.multiselect("Filter Player:", players)
                    if sel_p: dfh = dfh[dfh[col_p].isin(sel_p)]
                    else:
                        pat = '|'.join([p.lower() for p in OUR_TEAM])
                        dfh = dfh[dfh[col_p].str.lower().str.contains(pat, na=False)]

                    if not dfh.empty:
                        try:
                            try:
                                r = requests.get("https://valorant-api.com/v1/maps", timeout=3)
                                if r.status_code == 200:
                                    mapi = next((m for m in r.json()['data'] if m['displayName'].lower() == sel_map.lower()), None)
                                else: mapi = None
                            except: mapi = None

                            if mapi:
                                S = 1024
                                dfh['PX'] = (dfh['Y'] * mapi['xMultiplier'] + mapi['xScalarToAdd']) * S
                                dfh['PY'] = (dfh['X'] * mapi['yMultiplier'] + mapi['yScalarToAdd']) * S
                                fig = px.scatter(dfh, x='PX', y='PY', color=col_p, symbol='Weapon', width=700, height=700)
                                icon = get_map_img(sel_map, 'icon')
                                if icon:
                                    img = Image.open(icon)
                                    fig.add_layout_image(dict(source=img, xref="x", yref="y", x=0, y=0, sizex=S, sizey=S, sizing="stretch", layer="below"))
                                fig.update_layout(xaxis=dict(visible=False, range=[0,S]), yaxis=dict(visible=False, range=[S,0], scaleanchor="x"), margin=dict(l=0,r=0,t=0,b=0))
                                st.plotly_chart(fig)
                        except Exception as e: st.error(f"Map Error: {e}")
        else: st.info("No Heatmaps.")
        
        st.divider(); st.subheader("History")
        for idx, row in m_df.sort_values('DateObj', ascending=False).iterrows():
            with st.expander(f"{'✅' if row['Result']=='W' else '❌'} {row['Date']} | {int(row['Score_Us'])}-{int(row['Score_Enemy'])}"):
                c1,c2,c3=st.columns([2,2,1])
                with c1:
                    st.caption("NEXUS"); cols=st.columns(5)
                    for i in range(1,6):
                        ag=row.get(f'MyComp_{i}')
                        if pd.notna(ag): 
                            p=get_agent_img(ag); 
                            if p: cols[i-1].image(p, width=35)
                with c2:
                    st.caption("ENEMY"); cols=st.columns(5)
                    for i in range(1,6):
                        ag=row.get(f'EnComp_{i}')
                        if pd.notna(ag) and str(ag).strip()!="":
                            p=get_agent_img(ag); 
                            if p: cols[i-1].image(p, width=35)
                with c3:
                    st.write(f"ATK: {int(row.get('Atk_R_W',0))}/{int(row.get('Atk_R_L',0))}")
                    st.write(f"DEF: {int(row.get('Def_R_W',0))}/{int(row.get('Def_R_L',0))}")
                    if pd.notna(row.get('VOD_Link')): st.link_button("📺", row['VOD_Link'])

# ==============================================================================
# 4. STRATEGY BOARD
# ==============================================================================
elif page == "📘 STRATEGY BOARD":
    st.title("COMMAND CENTER")
    
    if 'sel_pb_id' not in st.session_state: st.session_state['sel_pb_id'] = None

    # Load Data (Variables are already loaded globally: df_team_pb, df_pb_strats, df_theory)
    
    # TABS (INTEGRATION OF WHITEBOARD)
    tab_playbooks, tab_theory, tab_lineups, tab_links = st.tabs(["🧠 TACTICAL PLAYBOOKS", "📜 MAP THEORY", "🎯 LINEUPS", "🔗 EXTERNAL LINKS"])

    # --------------------------------------------------------------------------
    # TAB 1: TACTICAL PLAYBOOKS
    # --------------------------------------------------------------------------
    with tab_playbooks:
        if st.session_state['sel_pb_id'] is None:
            c1, c2 = st.columns([3, 1])
            with c1: st.subheader("Active Playbooks")
            with c2: 
                with st.popover("➕ New Playbook"):
                    with st.form("create_pb"):
                        pm = st.selectbox("Map", sorted(df['Map'].unique()) if not df.empty else ["Ascent"], key="new_pb_map")
                        pn = st.text_input("Playbook Name (e.g. 'Standard Default')", key="new_pb_name")
                        st.caption("Select Composition:")
                        ac = st.columns(5)
                        ags = sorted([os.path.basename(x).replace(".png","").capitalize() for x in glob.glob(os.path.join(ASSET_DIR, "agents", "*.png"))])
                        sel_ags = [ac[i].selectbox(f"P{i+1}", [""]+ags, key=f"n_pb_{i}") for i in range(5)]
                        
                        if st.form_submit_button("Create System"):
                            new_id = str(uuid.uuid4())
                            new_row = {'ID': new_id, 'Map': pm, 'Name': pn}
                            for i in range(5): new_row[f'Agent_{i+1}'] = sel_ags[i]
                            
                            db_insert("nexus_playbooks", new_row, "df_team_pb")
                            
                            # Manual Clear
                            if "new_pb_name" in st.session_state: del st.session_state["new_pb_name"]
                            
                            st.rerun()

            if not df_team_pb.empty:
                for idx, row in df_team_pb.iterrows():
                    map_img = get_map_img(row['Map'], 'list')
                    b64_map = img_to_b64(map_img)
                    ag_html = ""
                    for i in range(1, 6):
                        a = row.get(f'Agent_{i}')
                        if a and pd.notna(a):
                            ab64 = img_to_b64(get_agent_img(a))
                            ab64 = get_styled_agent_img_b64(a)
                            if ab64: ag_html += f"<img src='data:image/png;base64,{ab64}' style='width:35px; height:35px; border-radius:50%; border:2px solid #111; margin-right:-10px; z-index:{i}'>"
                    
                    st.markdown(f"""<div class='pb-card'><div style="display:flex;align-items:center;gap:20px;"><div style="width:80px;height:50px;border-radius:5px;background-image:url('data:image/png;base64,{b64_map}');background-size:contain;background-position:center;border:1px solid #444;"></div><div><div style="color:#00BFFF;font-weight:bold;font-size:1.1em;text-transform:uppercase;">{row['Name']}</div><div style="color:#666;font-size:0.8em;">{row['Map']}</div></div><div style="margin-left:auto;padding-right:10px;">{ag_html}</div></div></div>""", unsafe_allow_html=True)
                    if st.button(f"OPEN TACTICS >>", key=f"btn_{row['ID']}"):
                        st.session_state['sel_pb_id'] = row['ID']; st.rerun()
            else:
                st.info("No Playbooks defined yet.")

        else:
            # Single Playbook
            # Safety Check: Ensure playbook exists and data is loaded
            if df_team_pb.empty or 'ID' not in df_team_pb.columns or st.session_state['sel_pb_id'] not in df_team_pb['ID'].values:
                st.warning("Playbook not found (or data is reloading). Returning to lobby.")
                st.session_state['sel_pb_id'] = None
                st.rerun()

            pb = df_team_pb[df_team_pb['ID'] == st.session_state['sel_pb_id']].iloc[0]
            # Define my_strats BEFORE using it in the export button
            my_strats = df_pb_strats[df_pb_strats['PB_ID'] == pb['ID']]
            
            st.button("⬅ BACK TO LOBBY", on_click=lambda: st.session_state.update({'sel_pb_id': None}))
            
            header_col1, header_col2 = st.columns([1, 4])

            # --- EXPORT FUNCTIONALITY ---
            with st.expander("🚀 Export Playbook to Presentation"):
                missing_libs = []
                if not HAS_PPTX: missing_libs.append("python-pptx")
                if not HAS_FPDF: missing_libs.append("fpdf")

                if missing_libs:
                    libs_str = " ".join(missing_libs)
                    st.error(f"⚠️ Fehlende Bibliotheken: {', '.join(missing_libs)}\n\n"
                             "Mögliche Lösung:\n"
                             "1. Stoppe den Server (Ctrl+C).\n"
                             "2. Stelle sicher, dass dein venv aktiv ist.\n"
                             f"3. Führe aus: `pip install {libs_str}`\n"
                             "4. Starte neu: `streamlit run app.py`")
                else:
                    st.info("Anleitung:")
                    st.markdown("""
                    1.  **Erstelle eine PowerPoint-Vorlage:**
                        *   Gestalte eine Titelfolie. Du kannst `{{PLAYBOOK_NAME}}` und `{{MAP_NAME}}` als Platzhalter für Text verwenden.
                        *   Gehe zu **Ansicht -> Folienmaster** und erstelle ein neues **Layout**. Benenne dieses Layout exakt in `Strat Layout` um.
                        *   Füge auf dem `Strat Layout` Platzhalter hinzu: einen **Titel**, ein **Inhalt**- oder **Text**-Feld (für Protokolle) und einen **Bild**-Platzhalter.
                        *   *Optional:* Füge auf der Titelfolie einen **Bild-Platzhalter** für die Agenten hinzu. Das Map-Bild wird automatisch als abgedunkelter Hintergrund gesetzt.
                    2.  **Hochladen & Generieren:** Lade deine Vorlage unten hoch und klicke auf Generieren. Die App erstellt für jede Strategie in diesem Playbook eine neue Folie.
                    """)
                    template_file = st.file_uploader("Upload your PowerPoint Template (.pptx)", type="pptx")
                    pdf_bg_file = st.file_uploader("Optional: Upload PDF Background Design (Image A4 Landscape)", type=['png', 'jpg'], help="This image will be used as the background for every page in the PDF.")

                    c_pptx, c_pdf = st.columns(2)

                    if template_file:
                        def generate_presentation(template, playbook_data, strats_data):
                            # Import locally to ensure it works if global import failed but library exists now
                            import pptx
                            from pptx.enum.shapes import PP_PLACEHOLDER_TYPE
                            prs = pptx.Presentation(template)
                            
                            # Find the custom layout
                            strat_layout = next((lo for lo in prs.slide_layouts if lo.name == 'Strat Layout'), None)
                            if not strat_layout:
                                st.error("Template Error: A slide layout named 'Strat Layout' was not found in your template. Please check the instructions.")
                                return None

                            # --- Populate Title Slide (optional) ---
                            title_slide = prs.slides[0]
                            
                            # Text Replacements
                            for shape in title_slide.shapes:
                                if not shape.has_text_frame: continue
                                if '{{PLAYBOOK_NAME}}' in shape.text:
                                    shape.text = shape.text.replace('{{PLAYBOOK_NAME}}', playbook_data['Name'])
                                if '{{MAP_NAME}}' in shape.text:
                                    shape.text = shape.text.replace('{{MAP_NAME}}', playbook_data['Map'])
                            
                            # 1. Map Image as Background (Dimmed)
                            map_path = get_map_img(playbook_data['Map'], 'list')
                            if map_path:
                                try:
                                    with Image.open(map_path) as img:
                                        enhancer = ImageEnhance.Brightness(img)
                                        dimmed_img = enhancer.enhance(0.7) # Reduce brightness by 30%
                                        output = io.BytesIO()
                                        dimmed_img.save(output, format="PNG")
                                        output.seek(0)
                                        title_slide.background.fill.user_picture(output)
                                except Exception as e:
                                    print(f"Error setting background: {e}")

                            # 2. Agents Composite (Insert into first available picture placeholder)
                            # FIX: Allow Object/Content placeholders too, and ensure we find them
                            pic_placeholders = []
                            for s in title_slide.placeholders:
                                if s.placeholder_format.type in [PP_PLACEHOLDER_TYPE.PICTURE, PP_PLACEHOLDER_TYPE.OBJECT]:
                                    pic_placeholders.append(s)
                            
                            if len(pic_placeholders) > 0:
                                agents = [playbook_data.get(f'Agent_{i}') for i in range(1,6)]
                                comp_img = create_team_composite(agents)
                                if comp_img: pic_placeholders[0].insert_picture(comp_img)

                            # --- Create a slide for each strat ---
                            for _, strat in strats_data.iterrows():
                                slide = prs.slides.add_slide(strat_layout)
                                
                                # Set Title
                                if slide.shapes.title: slide.shapes.title.text = strat['Name']
                                
                                # Set Protocols in Body
                                try: protos = json.loads(strat['Protocols'])
                                except: protos = []
                                body_text = "\n".join([f"• IF: {p['trigger']} -> THEN: {p['reaction']}" for p in protos])
                                    
                                # Find Placeholders dynamically (Robust against different template indices)
                                body_ph = None
                                pic_ph = None
                                
                                for shape in slide.placeholders:
                                    # Skip Title
                                    if slide.shapes.title and shape.shape_id == slide.shapes.title.shape_id: continue
                                    
                                    ph_type = shape.placeholder_format.type
                                    if ph_type == PP_PLACEHOLDER_TYPE.PICTURE:
                                        pic_ph = shape
                                    elif ph_type in [PP_PLACEHOLDER_TYPE.BODY, PP_PLACEHOLDER_TYPE.OBJECT]:
                                        if not body_ph: body_ph = shape
                                
                                # Set Text
                                if body_ph: body_ph.text = body_text

                                # Insert Image
                                img_path = os.path.join(STRAT_IMG_DIR, strat['Image'])
                                if pic_ph and os.path.exists(img_path):
                                    pic_ph.insert_picture(img_path)
                                    
                                    # Add Notes
                                    if pd.notna(strat.get('Notes')) and strat['Notes']:
                                        if not slide.has_notes_slide: slide.notes_slide
                                        slide.notes_slide.notes_text_frame.text = strat['Notes']

                            # Save to buffer
                            bio = io.BytesIO()
                            prs.save(bio)
                            return bio.getvalue()

                        with c_pptx:
                            st.download_button(
                                label="⬇️ Download PowerPoint",
                                data=generate_presentation(template_file, pb, my_strats),
                                file_name=f"{pb['Name']}_Playbook.pptx",
                                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                            )
                
                # PDF EXPORT
                if HAS_FPDF:
                    # --- PRESETS LOGIC ---
                    presets = {}
                    if os.path.exists(PRESETS_FILE):
                        try:
                            with open(PRESETS_FILE, 'r') as f: presets = json.load(f)
                        except: pass

                    # --- LAYOUT SETTINGS UI ---
                    layout_cfg = {}
                    with st.expander("🎨 PDF Layout Customization (Strategy Slides)"):
                        st.caption("Adjust positions (mm) and styles. Canvas size: ~339mm x 190mm.")
                        
                        # Preset Selector
                        c_pre, c_save = st.columns([2, 1])
                        sel_preset = c_pre.selectbox("Load Preset", ["Default"] + list(presets.keys()))
                        
                        default_layout = {'img_x':10, 'img_y':35, 'img_w':210, 'img_h':118, 'proto_x':228, 'proto_y':35, 'note_x':10, 'note_y':160, 'note_w': 80, 'title_x':10, 'title_y':10, 'title_size':24, 'title_color':'#FFFFFF', 'proto_size':12, 'if_color':'#00BFFF', 'then_color':'#FFFFFF', 'logo_x':10, 'logo_y':10, 'logo_w':20, 'tag_color': '#00BFFF', 'note_color': '#FFFFFF'}
                        
                        current_defaults = default_layout.copy()
                        if sel_preset != "Default" and sel_preset in presets:
                            current_defaults.update(presets[sel_preset])
                        
                        t1, t2, t3 = st.tabs(["📏 Positions & Sizes", "🎨 Colors & Fonts", "🖼️ Logo & Extras"])
                        
                        with t1:
                            l1, l2 = st.columns(2)
                            with l1:
                                st.markdown("**🖼️ Image Box**")
                                layout_cfg['img_x'] = st.number_input("Img X", value=current_defaults['img_x'])
                                layout_cfg['img_y'] = st.number_input("Img Y", value=current_defaults['img_y'])
                                layout_cfg['img_w'] = st.number_input("Img Width", value=current_defaults['img_w'])
                                layout_cfg['img_h'] = st.number_input("Img Height", value=current_defaults['img_h'])
                                st.markdown("**📝 Notes**")
                                layout_cfg['note_x'] = st.number_input("Note X", value=current_defaults['note_x'])
                                layout_cfg['note_y'] = st.number_input("Note Y", value=current_defaults['note_y'])
                                layout_cfg['note_w'] = st.number_input("Note Width (0=Full)", value=current_defaults.get('note_w', 80))
                            with l2:
                                st.markdown("**⚡ Protocols**")
                                layout_cfg['proto_x'] = st.number_input("Proto X", value=current_defaults['proto_x'])
                                layout_cfg['proto_y'] = st.number_input("Proto Y", value=current_defaults['proto_y'])
                                st.markdown("**🏷️ Title**")
                                layout_cfg['title_x'] = st.number_input("Title X", value=current_defaults['title_x'])
                                layout_cfg['title_y'] = st.number_input("Title Y", value=current_defaults['title_y'])
                        
                        with t2:
                            c1, c2 = st.columns(2)
                            with c1:
                                layout_cfg['title_size'] = st.number_input("Title Font Size", value=current_defaults.get('title_size', 24))
                                layout_cfg['title_color'] = st.color_picker("Title Color", current_defaults.get('title_color', '#FFFFFF'))
                                layout_cfg['tag_color'] = st.color_picker("Tag Color", current_defaults.get('tag_color', '#00BFFF'))
                            with c2:
                                layout_cfg['proto_size'] = st.number_input("Protocol Font Size", value=current_defaults.get('proto_size', 12))
                                layout_cfg['if_color'] = st.color_picker("IF Color (Bold)", current_defaults.get('if_color', '#00BFFF'))
                                layout_cfg['then_color'] = st.color_picker("THEN Color", current_defaults.get('then_color', '#FFFFFF'))
                                layout_cfg['note_color'] = st.color_picker("Notes Color", current_defaults.get('note_color', '#FFFFFF'))
                        
                        with t3:
                            st.markdown("**Logo Settings**")
                            c1, c2, c3 = st.columns(3)
                            layout_cfg['logo_x'] = c1.number_input("Logo X", value=current_defaults.get('logo_x', 10))
                            layout_cfg['logo_y'] = c2.number_input("Logo Y", value=current_defaults.get('logo_y', 10))
                            layout_cfg['logo_w'] = c3.number_input("Logo Width", value=current_defaults.get('logo_w', 20))

                        # Save Preset
                        with c_save:
                            new_preset_name = st.text_input("Save as Preset", placeholder="Name...")
                            if st.button("💾 Save"):
                                if new_preset_name:
                                    presets[new_preset_name] = layout_cfg
                                    with open(PRESETS_FILE, 'w') as f: json.dump(presets, f)
                                    st.success("Saved!")
                                    st.rerun()

                        if st.button("👁️ Preview Layout"):
                            # 16:9 Aspect Ratio (338.7mm x 190.5mm) -> Scale for Preview
                            scale = 3
                            w_px = int(338.7 * scale)
                            h_px = int(190.5 * scale)
                            
                            # Create Canvas
                            img = Image.new("RGB", (w_px, h_px), "#e0e0e0")
                            draw = ImageDraw.Draw(img)
                            
                            def draw_rect(x, y, w, h, color, text):
                                x_px, y_px = int(x*scale), int(y*scale)
                                w_px, h_px = int(w*scale), int(h*scale)
                                draw.rectangle([x_px, y_px, x_px+w_px, y_px+h_px], fill=color, outline="#333", width=2)
                                draw.text((x_px+10, y_px+10), text, fill="#000")

                            # Draw Elements
                            # 1. Title (Approx 200mm width)
                            draw_rect(layout_cfg['title_x'], layout_cfg['title_y'], 200, 15, "#ffcccb", "TITLE")
                            
                            # 2. Image (Height fixed, assume 16:9 ratio for width)
                            draw_rect(layout_cfg['img_x'], layout_cfg['img_y'], layout_cfg['img_w'], layout_cfg['img_h'], "#90ee90", "STRAT IMAGE")
                            
                            # 3. Protocols (Fixed width 100mm in code)
                            draw_rect(layout_cfg['proto_x'], layout_cfg['proto_y'], 100, 120, "#add8e6", "PROTOCOLS")
                            
                            # 4. Notes (Full width remaining)
                            note_w = layout_cfg.get('note_w', 0) if layout_cfg.get('note_w', 0) > 0 else (338.7 - layout_cfg['note_x'] - 10)
                            draw_rect(layout_cfg['note_x'], layout_cfg['note_y'], note_w, 20, "#ffffe0", "NOTES")
                            
                            st.image(img, caption="Layout Preview (Canvas 16:9)", use_container_width=True)

                    def generate_pdf_report(playbook_data, strats_data, bg_image=None, layout=None):
                        from fpdf import FPDF
                        
                        def _create_pdf(enable_custom_fonts):
                            # ⚙️ PDF SETTINGS (RESOLUTION / FORMAT)
                            # 16:9 Aspect Ratio (Standard PPT Widescreen: 338.7mm x 190.5mm)
                            pdf = FPDF(orientation='L', unit='mm', format=(190.5, 338.7))
                            pdf.set_auto_page_break(auto=True, margin=15)
                            
                            # Default Layout if None
                            local_layout = layout if layout is not None else default_layout.copy()

                            def hex_to_rgb(hex_col):
                                hex_col = hex_col.lstrip('#')
                                return tuple(int(hex_col[i:i+2], 16) for i in (0, 2, 4))

                            # --- FONTS ---
                            font_family = "Arial"
                            if enable_custom_fonts:
                                try:
                                    font_dir = os.path.join(ASSET_DIR, "fonts")
                                    regular_font = os.path.join(font_dir, "Nunito-Regular.ttf")
                                    bold_font = os.path.join(font_dir, "Nunito-Bold.ttf")
                                    if os.path.exists(regular_font) and os.path.exists(bold_font):
                                        pdf.add_font("Nunito", "", regular_font, uni=True)
                                        pdf.add_font("Nunito", "B", bold_font, uni=True)
                                        font_family = "Nunito"
                                except Exception as e:
                                    print(f"Font loading error: {e}")

                            # Handle Background Image (Save to temp file if uploaded)
                            bg_path = None
                            if bg_image:
                                bg_path = "temp_pdf_bg.png"
                                with open(bg_path, "wb") as f:
                                    f.write(bg_image.getbuffer())
                            
                            # --- TITLE PAGE ---
                            pdf.add_page()
                            pdf.set_fill_color(10, 10, 20) # Dark BG
                            pdf.rect(0, 0, pdf.w, pdf.h, 'F')
                            
                            # Title Slide ALWAYS uses Map Background (Dimmed) as requested
                            # Check assets/playbook/[map].png first
                            map_clean = str(playbook_data['Map']).lower().strip()
                            playbook_bg = os.path.join(ASSET_DIR, "playbook", f"{map_clean}.png")
                            
                            if os.path.exists(playbook_bg):
                                map_path = playbook_bg
                            else:
                                map_path = get_map_img(playbook_data['Map'], 'list')

                            if map_path:
                                try:
                                    # 🛠️ DIMMING LOGIC (30% reduced brightness)
                                    with Image.open(map_path) as img:
                                        enhancer = ImageEnhance.Brightness(img)
                                        dimmed_img = enhancer.enhance(0.7) # 0.7 = 70% brightness (30% reduced)
                                        
                                        temp_map_bg = "temp_map_bg_pdf.png"
                                        dimmed_img.save(temp_map_bg)
                                        
                                        # Position: x=0, y=0 (Top Left), w=Page Width, h=Page Height
                                        pdf.image(temp_map_bg, x=0, y=0, w=pdf.w, h=pdf.h)
                                except Exception as e:
                                    print(f"PDF BG Error: {e}")
                                    # Fallback if dimming fails
                                    pdf.image(map_path, x=0, y=0, w=pdf.w, h=pdf.h)
                            
                            # Title
                            pdf.set_y(80) # Move title down to 80mm from top
                            pdf.set_font(font_family, "B", 50)
                            pdf.set_text_color(255, 255, 255)
                            pdf.cell(0, 20, txt=playbook_data['Name'], ln=True, align='C')
                            
                            pdf.ln(5) # Add 5mm vertical space
                            pdf.set_font(font_family, "B", 24)
                            pdf.set_text_color(0, 191, 255) # Cyan
                            pdf.cell(0, 10, txt=playbook_data['Map'].upper(), ln=True, align='C')
                            
                            # Agents
                            agents = [playbook_data.get(f'Agent_{i}') for i in range(1,6)]
                            comp_img = create_team_composite(agents)
                            if comp_img:
                                # Save temp to include in PDF
                                with open("temp_comp.png", "wb") as f: f.write(comp_img.getbuffer())
                                pdf.image("temp_comp.png", x=(pdf.w-150)/2, y=140, w=150)
                            
                            # --- STRAT PAGES ---
                            for _, strat in strats_data.iterrows():
                                pdf.add_page()
                                pdf.set_fill_color(240, 240, 240)
                                pdf.rect(0, 0, pdf.w, pdf.h, 'F')
                                
                                # Apply Custom Background to strat pages too if provided
                                if bg_path:
                                    pdf.image(bg_path, x=0, y=0, w=pdf.w, h=pdf.h)
                                
                                # Logo
                                if os.path.exists("logo.png"):
                                    try:
                                        with Image.open("logo.png") as l_img:
                                            enh = ImageEnhance.Brightness(l_img)
                                            l_dim = enh.enhance(0.6) # 40% reduced brightness
                                            l_dim.save("temp_logo_dim.png")
                                        pdf.image("temp_logo_dim.png", x=local_layout.get('logo_x', 10), y=local_layout.get('logo_y', 10), w=local_layout.get('logo_w', 20))
                                    except: pass
                                
                                # Header
                                pdf.set_xy(local_layout['title_x'], local_layout['title_y'])
                                pdf.set_font(font_family, "B", local_layout.get('title_size', 24))
                                
                                # Title Color
                                tr, tg, tb = hex_to_rgb(local_layout.get('title_color', '#FFFFFF'))
                                pdf.set_text_color(tr, tg, tb)
                                
                                pdf.cell(0, 20, txt=strat['Name'], ln=True)
                                
                                # Tag (Below Title)
                                if pd.notna(strat.get('Tag')) and strat['Tag']:
                                    pdf.set_xy(local_layout['title_x'], local_layout['title_y'] + 12)
                                    pdf.set_font(font_family, "I", local_layout.get('title_size', 24) * 0.6) # Smaller font for tag
                                    
                                    # Tag Color
                                    tr, tg, tb = hex_to_rgb(local_layout.get('tag_color', '#00BFFF'))
                                    pdf.set_text_color(tr, tg, tb)
                                    
                                    pdf.cell(0, 10, txt=f"{strat['Tag']}", ln=True)
                                
                                # Image
                                img_path = os.path.join(STRAT_IMG_DIR, strat['Image'])
                                if os.path.exists(img_path):
                                    # Scale preserving aspect ratio (Fit within box)
                                    try:
                                        with Image.open(img_path) as img_pil:
                                            orig_w, orig_h = img_pil.size
                                        
                                        box_w = local_layout.get('img_w', 210)
                                        box_h = local_layout['img_h']
                                        scale = min(box_w/orig_w, box_h/orig_h)
                                        new_w, new_h = orig_w * scale, orig_h * scale
                                        
                                        # Bottom Right Alignment in Box
                                        # X = local_layout['img_x'] + box_w - new_w (Right)
                                        # Y = local_layout['img_y'] + box_h - new_h (Bottom)
                                        pdf.image(img_path, x=local_layout['img_x'] + box_w - new_w, y=local_layout['img_y'] + box_h - new_h, w=new_w, h=new_h)
                                    except:
                                        pdf.image(img_path, x=local_layout['img_x'], y=local_layout['img_y'], w=local_layout.get('img_w', 210), h=local_layout['img_h'])
                                
                                # Protocols
                                pdf.set_xy(local_layout['proto_x'], local_layout['proto_y'])
                                pdf.set_font(font_family, "B", local_layout.get('proto_size', 12) + 4) # Header slightly larger
                                # Header Color (Use Title Color or separate?) Let's use Title Color for consistency or IF color
                                tr, tg, tb = hex_to_rgb(local_layout.get('title_color', '#FFFFFF'))
                                pdf.set_text_color(tr, tg, tb)
                                pdf.cell(100, 10, txt="PROTOCOLS", ln=True)
                                
                                try: protos = json.loads(strat['Protocols'])
                                except: protos = []
                                
                                if_rgb = hex_to_rgb(local_layout.get('if_color', '#00BFFF'))
                                then_rgb = hex_to_rgb(local_layout.get('then_color', '#FFFFFF'))
                                
                                for p in protos:
                                    pdf.set_x(local_layout['proto_x'])
                                    # IF (Bold + Color)
                                    pdf.set_font(font_family, "B", local_layout.get('proto_size', 12))
                                    pdf.set_text_color(*if_rgb)
                                    pdf.multi_cell(100, 6, txt=f"IF: {p['trigger']}", border=0)
                                    
                                    # THEN (Normal + Color)
                                    pdf.set_x(local_layout['proto_x'])
                                    pdf.set_font(font_family, "", local_layout.get('proto_size', 12))
                                    pdf.set_text_color(*then_rgb)
                                    pdf.multi_cell(100, 6, txt=f"-> {p['reaction']}", border=0)
                                    
                                    pdf.ln(3)
                                
                                # Notes
                                if pd.notna(strat.get('Notes')) and strat['Notes']:
                                    pdf.set_xy(local_layout['note_x'], local_layout['note_y'])
                                    
                                    # Note Color
                                    nr, ng, nb = hex_to_rgb(local_layout.get('note_color', '#FFFFFF'))
                                    pdf.set_text_color(nr, ng, nb)
                                    
                                    pdf.set_font(font_family, "", 12)
                                    pdf.multi_cell(local_layout.get('note_w', 0), 6, txt=strat['Notes'])
                            
                            return pdf

                        try:
                            # Try with custom fonts first
                            pdf = _create_pdf(enable_custom_fonts=True)
                            return pdf.output(dest='S').encode('latin-1')
                        except Exception as e:
                            print(f"PDF Generation Error (Custom Fonts): {e}")
                            # Fallback: No custom fonts (Arial)
                            try:
                                pdf = _create_pdf(enable_custom_fonts=False)
                                return pdf.output(dest='S').encode('latin-1')
                            except Exception as e2:
                                return f"Error generating PDF: {str(e2)}".encode('utf-8')
                    
                    with c_pdf:
                        st.download_button(
                            label="⬇️ Download PDF Report",
                            data=generate_pdf_report(pb, my_strats, pdf_bg_file, layout_cfg),
                            file_name=f"{pb['Name']}_Playbook.pdf",
                            mime="application/pdf"
                        )

            with header_col1: st.image(get_map_img(pb['Map'], 'list'), width='stretch')
            with header_col2:
                st.markdown(f"<h1 style='margin:0'>{pb['Name']} <span style='font-size:0.5em; color:#666'>//{pb['Map']}</span></h1>", unsafe_allow_html=True)
                cols = st.columns(10)
                for i in range(1,6):
                    ag = pb.get(f'Agent_{i}'); 
                    if ag: cols[i-1].image(get_agent_img(ag), width=50)
                    if ag: 
                        b64 = get_styled_agent_img_b64(ag)
                        if b64: cols[i-1].image(f"data:image/png;base64,{b64}", width=50)

            st.divider()
            
            with st.expander("➕ ADD NEW STRATEGY / SET PLAY (IMAGE UPLOAD)"):
                # 1. Paste Button (Outside Form for immediate feedback)
                pasted_image = None
                if HAS_CLIPBOARD:
                    pasted_image = paste(label="📋 Click to Paste Image (from Clipboard)", key=f"paste_{pb['ID']}")
                    if pasted_image:
                        st.success("Image captured from clipboard!")
                else:
                    st.warning(f"Clipboard feature unavailable. Error: {CLIPBOARD_ERR}\n\nTry restarting the app if you just installed `st-img-pastebutton`.")

                with st.form("add_pb_strat"):
                    c_name, c_tag = st.columns([3, 1])
                    sn = c_name.text_input("Strategy Name", key="strat_name")
                    s_tag = c_tag.selectbox("Tag", ["Default", "Set Play", "Pistol", "Eco", "Anti-Eco", "Bonus", "Retake"], key="strat_tag")
                    s_note = st.text_area("Notes (Optional)", key="strat_note")
                    
                    si = st.file_uploader("Or Upload a File", type=['png', 'jpg'])

                    if st.form_submit_button("Add"):
                        image_data = None
                        if pasted_image:
                            # Handle pasted image data (usually base64 string)
                            if isinstance(pasted_image, str):
                                try:
                                    # Remove header if present (e.g., "data:image/png;base64,...")
                                    if "," in pasted_image:
                                        pasted_image = pasted_image.split(",")[1]
                                    # Fix padding if necessary
                                    missing_padding = len(pasted_image) % 4
                                    if missing_padding: pasted_image += '=' * (4 - missing_padding)
                                    image_data = base64.b64decode(pasted_image)
                                except: pass
                            else:
                                image_data = pasted_image
                        elif si:
                            image_data = si.getvalue()

                        if sn and image_data:
                            fname = f"PB_{pb['ID'][:8]}_{sn}_{int(datetime.now().timestamp())}.png".replace(" ", "_")
                            max_order = df_pb_strats[df_pb_strats['PB_ID'] == pb['ID']]['Order'].max()
                            new_order = max_order + 1 if pd.notna(max_order) else 0
                            with open(os.path.join(STRAT_IMG_DIR, fname), "wb") as f: f.write(image_data)
                            new_strat = {'PB_ID': pb['ID'], 'Strat_ID': str(uuid.uuid4()), 'Name': sn, 'Image': fname, 'Protocols': '[]', 'Notes': s_note, 'Tag': s_tag, 'Order': new_order}
                            
                            db_insert("nexus_pb_strats", new_strat, "df_pb_strats")
                            
                            # Manual Clear
                            for k in ["strat_name", "strat_note"]:
                                if k in st.session_state: del st.session_state[k]
                            st.rerun()
                        else:
                            st.error("A Strategy Name and an Image (pasted or uploaded) are required.")

            if not my_strats.empty:
                # Sort by Order
                my_strats_sorted = my_strats.sort_values('Order')
                
                for idx, strat in my_strats_sorted.iterrows():
                    with st.container():
                        # Layout: Up/Down | Content
                        c_nav, c_content = st.columns([0.5, 10])
                        
                        with c_nav:
                            if st.button("⬆️", key=f"up_{strat['Strat_ID']}"):
                                current_order = strat['Order']
                                # Find prev strat
                                prev = my_strats_sorted[my_strats_sorted['Order'] < current_order].sort_values('Order', ascending=False).head(1)
                                if not prev.empty:
                                    prev_idx = prev.index[0]
                                    prev_order = prev.iloc[0]['Order']
                                    
                                    # Swap in main DF
                                    db_update("nexus_pb_strats", {'ID': strat['ID'], 'Order': int(prev_order)}, "df_pb_strats")
                                    db_update("nexus_pb_strats", {'ID': prev.iloc[0]['ID'], 'Order': int(current_order)}, "df_pb_strats")
                                    st.rerun()
                                    
                            if st.button("⬇️", key=f"dn_{strat['Strat_ID']}"):
                                current_order = strat['Order']
                                # Find next strat
                                nxt = my_strats_sorted[my_strats_sorted['Order'] > current_order].sort_values('Order', ascending=True).head(1)
                                if not nxt.empty:
                                    nxt_idx = nxt.index[0]
                                    nxt_order = nxt.iloc[0]['Order']
                                    
                                    # Swap
                                    db_update("nexus_pb_strats", {'ID': strat['ID'], 'Order': int(nxt_order)}, "df_pb_strats")
                                    db_update("nexus_pb_strats", {'ID': nxt.iloc[0]['ID'], 'Order': int(current_order)}, "df_pb_strats")
                                    st.rerun()

                        with c_content:
                            c_img, c_proto = st.columns([1.5, 1])
                            with c_img:
                                st.subheader(f"📍 {strat['Name']}  `{strat.get('Tag', 'Default')}`")
                                spath = os.path.join(STRAT_IMG_DIR, strat['Image'])
                                if os.path.exists(spath): st.image(spath, use_container_width=True)
                                if pd.notna(strat.get('Notes')) and strat['Notes']:
                                    st.info(f"📝 **Notes:** {strat['Notes']}")
                            with c_proto:
                                st.markdown("### ⚡ PROTOCOLS")
                                try: protos = json.loads(strat['Protocols'])
                                except: protos = []
                                if protos:
                                    for p in protos: st.markdown(f"""<div class="proto-box"><div class="proto-if">IF: {p['trigger']}</div><div class="proto-then">👉 {p['reaction']}</div></div>""", unsafe_allow_html=True)
                                else: st.caption("No protocols defined.")
                                
                                with st.popover(f"Edit Protocols"):
                                    with st.form(f"pf_{strat['Strat_ID']}"):
                                        trig = st.text_input("IF (Trigger)"); react = st.text_input("THEN (Reaction)")
                                        if st.form_submit_button("Add"):
                                            protos.append({'trigger': trig, 'reaction': react})
                                            db_update("nexus_pb_strats", {'ID': strat['ID'], 'Protocols': json.dumps(protos)}, "df_pb_strats")
                                            st.rerun()
                                    if st.button("Clear Protocols", key=f"clr_{strat['Strat_ID']}"):
                                        db_update("nexus_pb_strats", {'ID': strat['ID'], 'Protocols': '[]'}, "df_pb_strats")
                                        st.rerun()
                            st.divider()

    # --------------------------------------------------------------------------
    # TAB 3: MAP THEORY
    # --------------------------------------------------------------------------
    with tab_theory:
        st.subheader("CONCEPTUAL FRAMEWORK")
        
        # 1. Visual Map Selector
        if 'theory_map' not in st.session_state: 
            st.session_state.theory_map = sorted(df['Map'].unique())[0] if not df.empty else "Ascent"
            
        with st.expander("🗺️ SELECT MAP", expanded=True):
            render_visual_selection(sorted(df['Map'].unique()), 'map', 'theory_sel', multi=False, key_state='theory_map')
            
        theory_map = st.session_state.theory_map
        
        if 'Image' not in df_theory.columns: df_theory['Image'] = None

        def get_theory_data(m, s):
            row = df_theory[(df_theory['Map'] == m) & (df_theory['Section'] == s)]
            if not row.empty: return row.iloc[0]['Content'], row.iloc[0]['Image']
            return "", None

        def save_theory_data_gsheet(m, s, text, new_img_obj, old_img_name):
            img_name = old_img_name
            if new_img_obj:
                img_name = f"THEORY_{m}_{s}_{int(datetime.now().timestamp())}.png".replace(" ", "_")
                with open(os.path.join(STRAT_IMG_DIR, img_name), "wb") as f: f.write(new_img_obj.getbuffer())
            
            # Remove old entry for this section
            new_df = df_theory[~((df_theory['Map'] == m) & (df_theory['Section'] == s))]
            # Add new entry
            new_entry = pd.DataFrame([{'Map': m, 'Section': s, 'Content': text, 'Image': img_name}])
            updated = pd.concat([new_df, new_entry], ignore_index=True)
            
            save_map_theory(updated)
            return img_name

        t_gen, t_atk, t_def = st.tabs(["🌐 GENERAL", "⚔️ ATTACK", "🛡️ DEFENSE"])
        sections = [("General", t_gen), ("Attack", t_atk), ("Defense", t_def)]
        
        for sec_name, sec_tab in sections:
            with sec_tab:
                # Unique key for edit mode state
                edit_key = f"edit_mode_{theory_map}_{sec_name}"
                if edit_key not in st.session_state: st.session_state[edit_key] = False
                
                curr_txt, curr_img = get_theory_data(theory_map, sec_name)
                
                # Header & Edit Toggle
                c_head, c_btn = st.columns([4, 1])
                with c_head: st.markdown(f"### {theory_map} - {sec_name} Theory")
                with c_btn:
                    if st.session_state[edit_key]:
                        if st.button("❌ Cancel", key=f"cn_{theory_map}_{sec_name}"):
                            st.session_state[edit_key] = False
                            st.rerun()
                    else:
                        if st.button("✏️ Edit", key=f"ed_{theory_map}_{sec_name}"):
                            st.session_state[edit_key] = True
                            st.rerun()
                
                st.divider()
                
                if st.session_state[edit_key]:
                    # EDIT MODE
                    # Improved Workflow: Image in Expander for better visibility
                    with st.expander("🖼️ Reference Image (View Large)", expanded=True):
                        if curr_img:
                            p = os.path.join(STRAT_IMG_DIR, curr_img)
                            if os.path.exists(p): st.image(p, use_container_width=True)
                        else:
                            st.info("No image currently uploaded.")

                    with st.form(f"form_{theory_map}_{sec_name}"):
                        new_txt = st.text_area("Content (Markdown supported)", value=curr_txt if curr_txt else "", height=500)
                        new_img = st.file_uploader("Upload New Image (Overwrites current)", type=['png', 'jpg'])
                        
                        if st.form_submit_button("💾 Save Changes", type="primary"):
                            save_theory_data_gsheet(theory_map, sec_name, new_txt, new_img, curr_img)
                            st.session_state[edit_key] = False
                            st.success("Saved!")
                            st.rerun()
                else:
                    # VIEW MODE
                    if not curr_txt and not curr_img:
                        st.info("No theory content added yet. Click 'Edit' to start writing.")
                    else:
                        # Toolbar: Layout Toggle & Export
                        tb1, tb2 = st.columns([2, 1])
                        with tb1:
                            layout_mode = st.radio("View Layout:", ["Side-by-Side", "Stacked (Large Image)"], horizontal=True, key=f"lo_{theory_map}_{sec_name}")
                        with tb2:
                            # Export Markdown
                            md_out = f"# {theory_map} - {sec_name} Theory\n\n"
                            if curr_txt: md_out += curr_txt + "\n\n"
                            if curr_img: md_out += f"![Reference Image]({curr_img})"
                            
                            st.download_button(
                                label="⬇️ Export Markdown",
                                data=md_out,
                                file_name=f"{theory_map}_{sec_name}_Theory.md",
                                mime="text/markdown",
                                key=f"exp_{theory_map}_{sec_name}"
                            )

                        st.divider()

                        if layout_mode == "Stacked (Large Image)":
                            if curr_img:
                                p = os.path.join(STRAT_IMG_DIR, curr_img)
                                if os.path.exists(p): st.image(p, use_container_width=True)
                            if curr_txt: st.markdown(curr_txt)
                        else:
                            c_view_txt, c_view_img = st.columns([1, 1])
                            with c_view_txt:
                                if curr_txt: st.markdown(curr_txt)
                                else: st.caption("_No text content._")
                            with c_view_img:
                                if curr_img:
                                    p = os.path.join(STRAT_IMG_DIR, curr_img)
                                    if os.path.exists(p): 
                                        st.image(p, use_container_width=True)
                                        with open(p, "rb") as f:
                                            st.download_button(label="Download Image", data=f, file_name=curr_img, mime="image/png", key=f"dl_{theory_map}_{sec_name}")
                                else:
                                    st.caption("_No image attached._")

    # --------------------------------------------------------------------------
    # TAB 3: LINEUPS
    # --------------------------------------------------------------------------
    with tab_lineups:
        st.subheader("🎯 LINEUP LIBRARY")
        
        # --- ADD NEW LINEUP ---
        with st.expander("➕ Add New Lineup"):
            # Paste Logic for Lineups
            pasted_lu = None
            if HAS_CLIPBOARD:
                pasted_lu = paste(label="📋 Paste Image", key="paste_lu")
                if pasted_lu: st.success("Image captured!")
            
            with st.form("add_lineup"):
                c1, c2 = st.columns(2)
                l_map = c1.selectbox("Map", sorted(df['Map'].unique()) if not df.empty else ["Ascent"], key="lu_map")
                l_agent = c2.selectbox("Agent", sorted(df_players['Agent'].unique()) if not df_players.empty else ["Sova"], key="lu_agent")
                
                c3, c4 = st.columns(2)
                l_side = c3.selectbox("Side", ["Attack", "Defense"], key="lu_side")
                
                # Dynamic Ability Select
                avail_abils = AGENT_ABILITIES.get(l_agent, {})
                abil_opts = [f"{k}: {v}" for k, v in avail_abils.items()] if avail_abils else ["Default"]
                l_ability = c4.selectbox("Ability", abil_opts, key="lu_abil")
                
                # Extra Tags (formerly Type)
                l_tags = st.multiselect("Tags / Type", ["Recon", "Shock", "Molly", "Flash", "Smoke", "Wall", "Ult", "One-Way", "God Spot", "Retake"], key="lu_tags")
                
                l_title = st.text_input("Title", placeholder="e.g. B Main God Arrow", key="lu_title")
                l_desc = st.text_area("Description / Instructions", key="lu_desc")
                l_vid = st.text_input("Video Link (optional)", key="lu_vid")
                l_file = st.file_uploader("Upload Image (optional)", type=['png', 'jpg'])
                
                if st.form_submit_button("💾 Save Lineup", type="primary"):
                    if not l_title:
                        st.error("Title is required.")
                    else:
                        # Handle Image
                        img_name = ""
                        img_data = None
                        
                        if pasted_lu:
                            if isinstance(pasted_lu, str) and "," in pasted_lu: pasted_lu = pasted_lu.split(",")[1]
                            try: img_data = base64.b64decode(pasted_lu)
                            except: pass
                        elif l_file:
                            img_data = l_file.getvalue()
                            
                        if img_data:
                            img_name = f"LU_{uuid.uuid4().hex[:8]}.png"
                            with open(os.path.join(STRAT_IMG_DIR, img_name), "wb") as f: f.write(img_data)
                        
                        # Clean Ability Name (remove "C: " prefix)
                        final_type = l_ability.split(": ")[1] if ": " in l_ability else l_ability
                        final_tags = ", ".join(l_tags)
                        
                        new_lu = {
                            'ID': str(uuid.uuid4())[:8],
                            'Map': l_map, 'Agent': l_agent, 'Side': l_side, 'Type': final_type,
                            'Title': l_title, 'Image': img_name, 'VideoLink': l_vid,
                            'Description': l_desc, 'Tags': final_tags, 'CreatedBy': st.session_state.get('username', 'Unknown')
                        }
                        db_insert("nexus_lineups", new_lu, "df_lineups")
                        st.success("Lineup saved!")
                        
                        # Manual Clear
                        for k in ["lu_title", "lu_desc", "lu_vid"]:
                            if k in st.session_state: del st.session_state[k]
                        st.rerun()

        st.divider()
        
        # --- FILTER & VIEW ---
        if not df_lineups.empty:
            c_f1, c_f2, c_f3 = st.columns(3)
            f_lu_map = c_f1.multiselect("Map", df_lineups['Map'].unique(), key="f_lu_map")
            f_lu_agent = c_f2.multiselect("Agent", df_lineups['Agent'].unique(), key="f_lu_agent")
            f_lu_type = c_f3.multiselect("Type", df_lineups['Type'].unique(), key="f_lu_type")
            
            view_lu = df_lineups.copy()
            if f_lu_map: view_lu = view_lu[view_lu['Map'].isin(f_lu_map)]
            if f_lu_agent: view_lu = view_lu[view_lu['Agent'].isin(f_lu_agent)]
            if f_lu_type: view_lu = view_lu[view_lu['Type'].isin(f_lu_type)]
            
            cols = st.columns(3)
            for idx, row in view_lu.iterrows():
                with cols[idx % 3]:
                    with st.container(border=True):
                        st.markdown(f"**{row['Title']}**")
                        st.caption(f"{row['Map']} • {row['Agent']} • {row['Side']}")
                        
                        if row['Image']:
                            ip = os.path.join(STRAT_IMG_DIR, row['Image'])
                            if os.path.exists(ip): st.image(ip, use_container_width=True)
                        
                        if row['VideoLink']: st.video(row['VideoLink'])
                        if row['Description']: st.info(row['Description'])
                        
                        if st.button("🗑️", key=f"del_lu_{row['ID']}"):
                            db_delete("nexus_lineups", row['ID'], "df_lineups")
                            st.rerun()
        else:
            st.info("No lineups found.")

    # --------------------------------------------------------------------------
    # TAB 4: EXTERNAL LINKS
    # --------------------------------------------------------------------------
    with tab_links:
        st.subheader("🔗 Playbooks")
        # Use Legacy Playbooks for external links if that's the intent or df_legacy_pb
        # Based on previous code, pb_df seemed to load PLAYBOOKS_FILE which is now df_legacy_pb
        pb_df = df_legacy_pb 
        
        # --- 1. ADD NEW LINK ---
        with st.expander("➕ Add New External Link"):
            with st.form("pb_link_form"):
                c1, c2 = st.columns([1, 2])
                with c1:
                    pm = st.selectbox("Map", sorted(df['Map'].unique()) if not df.empty else ["Ascent"], key="lm")
                    pn = st.text_input("Strategy Name", placeholder="e.g. Pistol Round Strat", key="ln")
                with c2:
                    pl = st.text_input("External Link (VOD/Doc)", placeholder="https://...", key="ll")
                    st.caption("Select Agents (Optional Composition):")
                    ags = sorted([os.path.basename(x).replace(".png","").capitalize() for x in glob.glob(os.path.join(ASSET_DIR, "agents", "*.png"))])
                    cols = st.columns(5)
                    mas = [cols[i].selectbox(f"Agent {i+1}", [""]+ags, key=f"la{i}", label_visibility="collapsed") for i in range(5)]
                
                if st.form_submit_button("💾 Save Link", type="primary"): 
                    if pn and pl:
                        nr = {'Map': pm, 'Name': pn, 'Link': pl}
                        for i in range(5): nr[f'Agent_{i+1}'] = mas[i]
                        
                        db_insert("playbooks", nr, "df_legacy_pb")
                        st.success("Link saved!")
                        
                        # Manual Clear
                        for k in ["ln", "ll"]:
                            if k in st.session_state: del st.session_state[k]
                        st.rerun()
                    else:
                        st.error("Name and Link are required.")

        st.divider()

        # --- 2. FILTER & DISPLAY ---
        if not pb_df.empty:
            # Visual Map Selector
            if 'link_map_sel' not in st.session_state: 
                st.session_state.link_map_sel = sorted(pb_df['Map'].unique())[0] if not pb_df.empty else "Ascent"
            
            # Toggle for "Show All"
            c_filter, c_toggle = st.columns([3, 1])
            with c_toggle:
                show_all = st.checkbox("Show All Maps", value=False)
            
            selected_map = st.session_state.link_map_sel
            
            if not show_all:
                with c_filter:
                    with st.expander("🗺️ SELECT MAP", expanded=True):
                        render_visual_selection(sorted(pb_df['Map'].unique()), 'map', 'link_sel', multi=False, key_state='link_map_sel')
                v_pb = pb_df[pb_df['Map'] == selected_map]
            else:
                v_pb = pb_df

            if v_pb.empty:
                st.info(f"No links found for {selected_map}.")
            else:
                # --- 3. DISPLAY CARDS ---
                for idx, row in v_pb.iterrows():
                    with st.container():
                        # Layout: Image | Content | Action
                        c_img, c_content, c_action = st.columns([1, 3, 1])
                        
                        # Map Image
                        with c_img:
                            map_img = get_map_img(row['Map'], 'list')
                            if map_img: st.image(map_img, use_container_width=True)
                            else: st.markdown(f"**{row['Map']}**")
                        
                        # Content
                        with c_content:
                            st.markdown(f"### {row['Name']}")
                            st.caption(f"Map: {row['Map']}")
                            
                            # Agents
                            ag_cols = st.columns(10)
                            for i in range(1, 6):
                                a = row.get(f'Agent_{i}')
                                if a and pd.notna(a) and str(a).strip() != "":
                                    b64 = get_styled_agent_img_b64(a)
                                    if b64: ag_cols[i-1].image(f"data:image/png;base64,{b64}", width=40)
                        
                        # Action
                        with c_action:
                            st.link_button("🔗 OPEN", row['Link'], use_container_width=True)
                            if st.button("🗑️", key=f"del_link_{idx}"):
                                db_delete("playbooks", row['ID'], "df_legacy_pb")
                                st.rerun()
                        
                        st.markdown("---")
        else:
            st.info("No external links added yet.")

# ==============================================================================
# 5. RESOURCES
# ==============================================================================
elif page == "📚 RESOURCES":
    st.title("KNOWLEDGE BASE")
    # df_res loaded globally
    
    with st.expander("➕ Add"):
        with st.form("ra"):
            rt = st.text_input("Title", key="res_title"); rl = st.text_input("Link", key="res_link"); rc = st.selectbox("Cat", ["Theory", "Lineups", "Setup", "Playbook Theory"], key="res_cat"); rn = st.text_area("Note", key="res_note")
            if st.form_submit_button("Save"):
                db_insert("resources", {'Title': rt, 'Link': rl, 'Category': rc, 'Note': rn}, "df_res")
                # Manual Clear
                for k in ["res_title", "res_link", "res_note"]:
                    if k in st.session_state: del st.session_state[k]
                st.rerun()
    
    if not df_res.empty:
        cats = st.multiselect("Filter:", df_res['Category'].unique(), default=df_res['Category'].unique())
        view = df_res[df_res['Category'].isin(cats)]
        cols = st.columns(4)
        for i, (idx, row) in enumerate(view.iterrows()):
            with cols[i%4]:
                thumb = get_yt_thumbnail(row['Link'])
                img = f"<img src='{thumb}' class='res-thumb'>" if thumb else "<div style='height:140px; background:#222; display:flex; align-items:center; justify-content:center'>📄</div>"
                st.markdown(f"""<div class="res-tile">{img}<div class="res-info"><div style="color:#00BFFF; font-size:0.8em">{row['Category']}</div><div style="font-weight:bold">{row['Title']}</div><a href="{row['Link']}" target="_blank" style="color:#aaa; font-size:0.8em">OPEN</a></div></div>""", unsafe_allow_html=True)
    
    with st.expander("✏️ Edit"):
        ed = st.data_editor(df_res, num_rows="dynamic")
        if st.button("Save Changes"): 
            save_resources(ed)
            st.success("Saved"); st.rerun()

# ==============================================================================
# 6. CALENDAR
# ==============================================================================
elif page == "📅 CALENDAR":
    st.title("SCHEDULE")
    if 'cy' not in st.session_state: st.session_state['cy'] = datetime.now().year
    if 'cm' not in st.session_state: st.session_state['cm'] = datetime.now().month
    
    # NEW: Get User
    current_user = st.session_state.get('username', '')
    
    c1, c2 = st.columns([3, 1]) # Bigger Calendar (3:1 ratio)
    # df_cal and df_simple_todos loaded globally
    
    with c1:
        # --- FILTER ---
        p_list = ["Luggi","Benni","Andrei","Luca","Sofi","Remus"]
        filter_opts = ["All", "My Events"] + p_list
        def_idx = 1 if st.session_state.get('role') == 'player' else 0
        
        sel_filter = st.selectbox("Filter Events:", filter_opts, index=def_idx)
        
        df_cal_view = df_cal.copy()
        # Robustness: Ensure columns exist to prevent KeyError
        for c in ['Date', 'Time', 'Event', 'Map', 'Type', 'Players']:
            if c not in df_cal_view.columns: df_cal_view[c] = ""

        if sel_filter != "All":
            target_user = current_user if sel_filter == "My Events" else sel_filter
            if target_user:
                df_cal_view = df_cal_view[df_cal_view['Players'].apply(
                    lambda x: pd.isna(x) or str(x).strip() == "" or target_user in str(x)
                )]

        cp, cc, cn = st.columns([1,2,1])
        if cp.button("<"): st.session_state['cm']-=1
        if cn.button(">"): st.session_state['cm']+=1
        if st.session_state['cm']==0: st.session_state['cm']=12; st.session_state['cy']-=1
        if st.session_state['cm']==13: st.session_state['cm']=1; st.session_state['cy']+=1
        
        curr = datetime(st.session_state['cy'], st.session_state['cm'], 1)
        cc.markdown(f"<h3 style='text-align:center'>{curr.strftime('%B %Y')}</h3>", unsafe_allow_html=True)
        
        cal = calendar.monthcalendar(curr.year, curr.month)
        cols = st.columns(7)
        for d in ["Mon","Tue","Wed","Thu","Fri","Sat","Sun"]: cols[["Mon","Tue","Wed","Thu","Fri","Sat","Sun"].index(d)].write(f"**{d}**")
        
        for w in cal:
            cols = st.columns(7)
            for i, d in enumerate(w):
                with cols[i]:
                    if d!=0:
                        d_s = f"{d:02d}.{curr.month:02d}.{curr.year}"
                        evs = df_cal_view[df_cal_view['Date']==d_s]
                        
                        is_today = date(curr.year, curr.month, d) == date.today()
                        extra_cls = " cal-today" if is_today else ""
                        
                        h = f"<div class='cal-day-box{extra_cls}'><div class='cal-date'>{d}</div>"
                        for _, e in evs.iterrows():
                            c = "#FF1493" if e.get('Type')=="Match" else "#00BFFF"
                            p_info = f"&#10;👥 {e['Players']}" if pd.notna(e.get('Players')) and e['Players'] else ""
                            h+=f"<div class='cal-event-pill' style='background:{c}20; border-left-color:{c}' title='{e['Time']} {e['Event']}{p_info}'>{e['Time']} {e['Event']}</div>"
                        st.markdown(h+"</div>", unsafe_allow_html=True)
        
        with st.expander("Add Event"):
            with st.form("ca"):
                c_d, c_t = st.columns(2)
                cd=c_d.date_input("Date", key="cal_date"); ct=c_t.time_input("Time", key="cal_time")
                ce=st.text_input("Event Name", key="cal_name"); 
                c_m, c_ty = st.columns(2)
                cm=c_m.text_input("Map (Optional)", key="cal_map"); cty=c_ty.selectbox("Type",["Match","Scrim","Other"], key="cal_type")
                
                # Player Selector
                cp = st.multiselect("Assign Players", ["Luggi","Benni","Andrei","Luca","Sofi","Remus"], default=[], key="cal_players")
                
                if st.form_submit_button("Add"):
                    p_str = ", ".join(cp)
                    db_insert("calendar", {'Date':cd.strftime("%d.%m.%Y"),'Time':ct.strftime("%H:%M"),'Event':ce,'Map':cm,'Type':cty,'Players':p_str}, "df_cal")
                    # Manual Clear
                    for k in ["cal_name", "cal_map", "cal_players"]:
                        if k in st.session_state: del st.session_state[k]
                    st.rerun()

        with st.expander("Manage Events (Edit Dates / Delete)"):
            st.caption("💡 You can edit dates here to move events.")
            edited_cal = st.data_editor(df_cal, num_rows="dynamic", use_container_width=True)
            if st.button("Save Schedule Changes"):
                save_calendar(edited_cal)
                st.success("Schedule updated!")
                st.rerun()

    with c2:
        st.markdown("### ✅ QUICK TASKS")
        
        # Progress Bar
        if not df_simple_todos.empty:
            done_count = len(df_simple_todos[df_simple_todos['Done']==True])
            total_count = len(df_simple_todos)
            prog = done_count / total_count if total_count > 0 else 0
            st.progress(prog, text=f"{done_count}/{total_count} Completed")
        
        # Quick Add
        with st.form("quick_todo", clear_on_submit=True):
            c_in, c_btn = st.columns([3, 1])
            new_task = c_in.text_input("New Task", label_visibility="collapsed", placeholder="Add task...")
            if c_btn.form_submit_button("➕"):
                if new_task:
                    updated = pd.concat([df_simple_todos, pd.DataFrame([{'Task':new_task, 'Done':False}])], ignore_index=True)
                    save_simple_todos(updated)
                    st.rerun()

        # Visual List
        if not df_simple_todos.empty:
            st.markdown("<div style='height: 10px'></div>", unsafe_allow_html=True)
            # Sort: Pending first
            for i, row in df_simple_todos.sort_values('Done').iterrows():
                done = row['Done']
                label = f"~~{row['Task']}~~" if done else row['Task']
                
                # Custom Checkbox Row
                if st.checkbox(label, value=done, key=f"todo_{i}"):
                    if not done:
                        df_simple_todos.at[i, 'Done'] = True
                        save_simple_todos(df_simple_todos)
                        st.rerun()
                else:
                    if done:
                        df_simple_todos.at[i, 'Done'] = False
                        save_simple_todos(df_simple_todos)
                        st.rerun()
            
            if st.button("🗑️ Clear Completed", use_container_width=True):
                df_simple_todos = df_simple_todos[df_simple_todos['Done'] == False]
                save_simple_todos(df_simple_todos)
                st.rerun()

# ==============================================================================
# 7. PLAYERS (KOMPLETT NEU MIT DEEP DIVE)
# ==============================================================================
elif page == "📊 PLAYERS":
    st.title("PLAYER PERFORMANCE")

    # Styling Functions for Conditional Formatting (moved here to be available for all tabs)
    def style_good_bad(v, good_thresh, bad_thresh, inverse=False):
        if pd.isna(v): return ""
        # Colors: Dark Green / Dark Yellow / Dark Red backgrounds with light text
        c_good = 'background-color: #113321; color: #aaffaa'
        c_avg = 'background-color: #443311; color: #ffffaa'
        c_bad = 'background-color: #441111; color: #ffaaaa'
        
        if inverse: # Lower is better (e.g. Deaths)
            if v <= good_thresh: return c_good
            elif v <= bad_thresh: return c_avg
            else: return c_bad
        else: # Higher is better
            if v >= good_thresh: return c_good
            elif v >= bad_thresh: return c_avg
            else: return c_bad
            
    def style_purple_gradient(s):
        """Custom gradient to avoid matplotlib dependency"""
        try:
            s_num = pd.to_numeric(s, errors='coerce').fillna(0)
            max_val = s_num.max()
            if max_val == 0: return [''] * len(s)
            return [f'background-color: rgba(147, 112, 219, {0.1 + (v/max_val)*0.6})' for v in s_num]
        except:
            return [''] * len(s)

    tab_overview, tab_deep = st.tabs(["📊 TEAM OVERVIEW", "🧬 DEEP DIVE ANALYZER"])
    
    with tab_overview:
        if df_players.empty:
            st.info("No player stats yet. Import JSON matches to see data.")
        else:
            # --- 1. FILTER SECTION ---
            st.markdown("### 🔎 FILTER")
            
            # Get unique values for filters
            available_maps = sorted(df_players['Map'].unique())
            available_agents = sorted(df_players['Agent'].unique())
            available_players = sorted(df_players['Player'].unique())
            
            # Determine default players
            default_players = [p for p in available_players if any(t.lower() in p.lower() for t in OUR_TEAM)]
            if not default_players: default_players = available_players

            # --- VISUAL FILTERS ---
            with st.expander("🗺️ MAPS & ♟️ AGENTS FILTER", expanded=False):
                st.caption("Select Maps:")
                sel_maps = render_visual_selection(available_maps, 'map', 'p_map_flt', default=available_maps, multi=True)
                st.divider()
                st.caption("Select Agents:")
                sel_agents = render_visual_selection(available_agents, 'agent', 'p_ag_flt', default=available_agents, multi=True)

            # Player filter
            sel_players = st.multiselect("Select Players:", available_players, default=default_players)

            # --- 2. DATA PROCESSING ---
            df_filtered = df_players.copy()
            
            if sel_maps:
                df_filtered = df_filtered[df_filtered['Map'].isin(sel_maps)]
            if sel_agents:
                df_filtered = df_filtered[df_filtered['Agent'].isin(sel_agents)]
            if sel_players:
                df_filtered = df_filtered[df_filtered['Player'].isin(sel_players)]
                
            if df_filtered.empty:
                st.warning("No data matches your filters.")
            else:
                # Aggregate Data
                # Ensure advanced stats exist (fallback for old data)
                for c in ['ADR', 'FK', 'FD']:
                    if c not in df_filtered.columns: df_filtered[c] = 0

                # Ensure numeric columns for aggregation to avoid TypeErrors
                num_cols = ['Kills', 'Deaths', 'Assists', 'Score', 'Rounds', 'HS', 'ADR', 'FK', 'FD']
                for col in num_cols:
                    if col in df_filtered.columns:
                        df_filtered[col] = pd.to_numeric(df_filtered[col], errors='coerce').fillna(0)

                p_agg = df_filtered.groupby('Player').agg({
                    'MatchID': 'nunique', 
                    'Kills': 'sum', 
                    'Deaths': 'sum', 
                    'Assists': 'sum', 
                    'Score': 'sum', 
                    'Rounds': 'sum',
                    'HS': 'mean',
                    'ADR': 'mean', 'FK': 'sum', 'FD': 'sum'
                }).reset_index()
                
                # Calculate Metrics
                p_agg['SafeDeaths'] = p_agg['Deaths'].replace(0, 1)
                p_agg['KD'] = p_agg['Kills'] / p_agg['SafeDeaths']
                p_agg['ACS'] = p_agg['Score'] / p_agg['Rounds'].replace(0, 1)
                p_agg['KPR'] = p_agg['Kills'] / p_agg['Rounds'].replace(0, 1)
                p_agg['APR'] = p_agg['Assists'] / p_agg['Rounds'].replace(0, 1)
                p_agg['DPR'] = p_agg['Deaths'] / p_agg['Rounds'].replace(0, 1)
                # Rating Formula: KPR + (ADR/450) + (FK-FD)/Rounds
                p_agg['Rating'] = p_agg['KPR'] + (p_agg['ADR'] / 450) + ((p_agg['FK'] - p_agg['FD']) / p_agg['Rounds'].replace(0, 1))
                
                p_display = p_agg.rename(columns={'MatchID': 'Matches'})
                # Sort by Rating
                p_display = p_display.sort_values('Rating', ascending=False)
                
                # --- 3. TABLE VIEW ---
                st.divider()
                st.markdown("### 📋 STATS OVERVIEW")
                
                # Styling Functions
                def style_good_bad(v, good_thresh, bad_thresh, inverse=False):
                    if pd.isna(v): return ""
                    c_good = 'background-color: #113321; color: #aaffaa'
                    c_avg = 'background-color: #443311; color: #ffffaa'
                    c_bad = 'background-color: #441111; color: #ffaaaa'
                    
                    if inverse: 
                        if v <= good_thresh: return c_good
                        elif v <= bad_thresh: return c_avg
                        else: return c_bad
                    else: 
                        if v >= good_thresh: return c_good
                        elif v >= bad_thresh: return c_avg
                        else: return c_bad

                # Apply Pandas Styling
                styler = p_display[['Player', 'Matches', 'Rating', 'KD', 'ACS', 'HS', 'KPR', 'APR', 'DPR']].style\
                    .format({"Rating": "{:.2f}", "KD": "{:.2f}", "ACS": "{:.0f}", "HS": "{:.1f}%", "KPR": "{:.2f}", "APR": "{:.2f}", "DPR": "{:.2f}"})\
                    .map(lambda v: style_good_bad(v, 1.2, 0.85), subset=['Rating'])\
                    .map(lambda v: style_good_bad(v, 1.1, 0.9), subset=['KD'])\
                    .map(lambda v: style_good_bad(v, 230, 180), subset=['ACS'])\
                    .map(lambda v: style_good_bad(v, 25, 15), subset=['HS'])\
                    .map(lambda v: style_good_bad(v, 0.8, 0.6), subset=['KPR'])\
                    .map(lambda v: style_good_bad(v, 0.35, 0.15), subset=['APR'])\
                    .map(lambda v: style_good_bad(v, 0.65, 0.80, inverse=True), subset=['DPR'])

                st.dataframe(
                    styler,
                    column_config={
                        "Player": st.column_config.TextColumn("Player", width="medium"),
                        "Matches": st.column_config.NumberColumn("Matches", format="%d"),
                        "Rating": st.column_config.NumberColumn("Rating", format="%.2f"),
                        "KD": st.column_config.NumberColumn("K/D"),
                    },
                    use_container_width=True,
                    hide_index=True
                )
                
                # --- 4. SPIDER CHART ---
                st.divider()
                st.markdown("### 🕸️ TRAIT ANALYSIS")
                
                c_chart, c_metrics = st.columns([1.5, 1])
                player_list = p_agg['Player'].tolist()
                
                with c_metrics:
                    st.markdown("#### Compare Players")
                    p1_name = st.selectbox("Player 1 (Blue)", player_list, index=0 if len(player_list)>0 else None)
                    p2_name = st.selectbox("Player 2 (Pink)", ["None"] + player_list, index=1 if len(player_list)>1 else 0)
                    
                    if p1_name:
                        row1 = p_agg[p_agg['Player'] == p1_name].iloc[0]
                        
                        def show_metric_comp(label, val1, val2=None, fmt="{:.2f}"):
                            delta = None
                            if val2 is not None: delta = val1 - val2
                            st.metric(label, fmt.format(val1), f"{delta:.2f}" if delta is not None else None)

                        if p2_name and p2_name != "None":
                            row2 = p_agg[p_agg['Player'] == p2_name].iloc[0]
                            st.markdown("---")
                            c_m1, c_m2 = st.columns(2)
                            with c_m1: show_metric_comp("K/D", row1['KD'], row2['KD'])
                            with c_m2: show_metric_comp("ACS", row1['ACS'], row2['ACS'], "{:.0f}")
                            c_m3, c_m4 = st.columns(2)
                            with c_m3: show_metric_comp("KPR", row1['KPR'], row2['KPR'])
                            with c_m4: show_metric_comp("HS%", row1['HS'], row2['HS'], "{:.1f}%")
                        else:
                            st.markdown("---")
                            c_m1, c_m2 = st.columns(2)
                            with c_m1: st.metric("K/D", f"{row1['KD']:.2f}")
                            with c_m2: st.metric("ACS", f"{row1['ACS']:.0f}")
                            c_m3, c_m4 = st.columns(2)
                            with c_m3: st.metric("KPR", f"{row1['KPR']:.2f}")
                            with c_m4: st.metric("HS%", f"{row1['HS']:.1f}%")

                with c_chart:
                    if p1_name:
                        categories = ['K/D', 'ACS', 'HS%', 'KPR', 'APR']
                        max_values = {'K/D': 2.0, 'ACS': 300, 'HS%': 40, 'KPR': 1.0, 'APR': 0.5}
                        radar_data = []
                        
                        def add_player_data(p_row, p_label):
                            vals = {'K/D': p_row['KD'], 'ACS': p_row['ACS'], 'HS%': p_row['HS'], 'KPR': p_row['KPR'], 'APR': p_row['APR']}
                            for cat in categories:
                                norm = min(vals[cat] / max_values[cat], 1.0)
                                radar_data.append({'Player': p_label, 'Metric': cat, 'Value': norm, 'Display': vals[cat]})
                        
                        add_player_data(row1, p1_name)
                        if p2_name and p2_name != "None": add_player_data(row2, p2_name)
                        
                        fig = px.line_polar(pd.DataFrame(radar_data), r='Value', theta='Metric', color='Player', line_close=True,
                                            color_discrete_sequence=['#00BFFF', '#FF1493'], hover_data={'Value': False, 'Display': True})
                        fig.update_traces(fill='toself')
                        fig.update_layout(polar=dict(radialaxis=dict(visible=False, range=[0, 1]), bgcolor='rgba(255,255,255,0.05)'),
                                          paper_bgcolor='rgba(0,0,0,0)', plot_bgcolor='rgba(0,0,0,0)', font=dict(color='white'),
                                          margin=dict(t=20, b=20, l=20, r=20), legend=dict(orientation="h", y=1.02))
                        st.plotly_chart(fig, use_container_width=True)

    with tab_deep:
        # --- ROLE MAPPING ---
        AGENT_ROLES = {
            "Duelist": ["Jett", "Phoenix", "Reyna", "Raze", "Yoru", "Neon", "Iso"],
            "Initiator": ["Sova", "Breach", "Skye", "KAY/O", "Fade", "Gekko"],
            "Controller": ["Brimstone", "Viper", "Omen", "Astra", "Harbor", "Clove"],
            "Sentinel": ["Sage", "Cypher", "Killjoy", "Chamber", "Deadlock", "Vyse"]
        }

        st.markdown("### 🧬 INDIVIDUAL PLAYER DEEP DIVE (Ranked Data)")
        deep_player = st.selectbox("Select Player", ["Andrei", "Benni", "Luca", "Luggi", "Remus", "Sofi"])
        
        # Pfad laden
        file_path = os.path.join(BASE_DIR, "data", "players", f"{deep_player.lower()}_data.json")
        df_deep = pd.DataFrame()
        
        if os.path.exists(file_path):
            df_deep = parse_tracker_json(file_path)
        else:
            st.info(f"No local data file found for {deep_player}. Please upload a Tracker.gg JSON export.")
            up = st.file_uploader("Upload Player JSON Export", type=['json'], key=f"deep_dive_upload_{deep_player}")
            if up: df_deep = parse_tracker_json(up); st.rerun()

        # --- HIER BEGINNT DIE ANALYSE ---
        if not df_deep.empty:
            
            if df_deep.empty:
                st.warning(f"⚠️ Datei geladen, aber keine Daten für Spieler **{deep_player}** gefunden.")
                st.stop()

            # SICHERHEITS-CHECK: Sind die neuen Daten da?
            if 'Cast_Ult' not in df_deep.columns:
                st.error("⚠️ FEHLER: Dein Parser ist veraltet. Bitte führe SCHRITT 1 aus meiner Nachricht aus (Funktion parse_tracker_json aktualisieren).")
                st.stop()

            # 1. TOP STATS ROW
            k1, k2, k3, k4, k5 = st.columns(5)
            # Fix K/D: Use Total Kills / Total Deaths instead of Mean of Ratios
            dd_k = df_deep['Kills'].sum()
            dd_d = df_deep['Deaths'].sum()
            dd_kd = dd_k / dd_d if dd_d > 0 else dd_k
            k1.metric("K/D RATIO", f"{dd_kd:.2f}")
            k2.metric("HEADSHOT %", f"{df_deep['HS%'].mean():.1f}%")
            k3.metric("ADR", f"{df_deep['ADR'].mean():.0f}")
            upr = df_deep['Total_Util'].sum() / df_deep['Rounds'].sum() if df_deep['Rounds'].sum() > 0 else 0
            k4.metric("UTIL / ROUND", f"{upr:.2f}")
            # VLR Rating Approximation (Simple)
            # NEW FORMULA: KPR + (ADR/450) + (FK-FD)/Rounds -> Avg ~ 1.0, High ~ 1.5
            vlr_approx = (df_deep['Kills'].sum()/df_deep['Rounds'].sum()) + (df_deep['ADR'].mean()/450) + ((df_deep['FK'].sum() - df_deep['FD'].sum())/df_deep['Rounds'].sum())
            k5.metric("VLR RATING", f"{vlr_approx:.2f}")
            
            st.divider()

            # 2. AGENT BREAKDOWN (VLR.gg Style Table)
            st.subheader("♟️ AGENT PERFORMANCE")
            
            # Daten aggregieren
            ag_stats = df_deep.groupby('Agent').agg({
                'MatchesPlayed': 'sum',
                'Kills': 'sum', 'Deaths': 'sum',
                'Wins': 'sum',
                'Rounds': 'sum',
                'HS%': 'mean',
                'Cast_Grenade':'sum', 'Cast_Abil1':'sum', 'Cast_Abil2':'sum', 'Cast_Ult':'sum',
                'FK': 'sum', 'FD': 'sum', 'ADR': 'mean'
            }).reset_index()
            ag_stats['Win%'] = ag_stats['Wins'] / ag_stats['MatchesPlayed'] * 100
            ag_stats['KD'] = ag_stats['Kills'] / ag_stats['Deaths'].replace(0,1)
            
            # Util pro Runde berechnen
            safe_rounds = ag_stats['Rounds'].replace(0, 1)
            ag_stats['C_per_round'] = ag_stats['Cast_Grenade'] / safe_rounds
            ag_stats['Q_per_round'] = ag_stats['Cast_Abil1'] / safe_rounds
            ag_stats['E_per_round'] = ag_stats['Cast_Abil2'] / safe_rounds
            ag_stats['X_per_round'] = ag_stats['Cast_Ult'] / safe_rounds
            
            # VLR Rating Calculation per Agent
            # Formula: KPR + (APR/3) + (ADR/300) + (FK-FD)/Rounds + 0.2 (Offset)
            ag_stats['KPR'] = ag_stats['Kills'] / safe_rounds
            ag_stats['APR'] = (ag_stats['Kills'] * 0.3) / safe_rounds # Approximation if Assists not in agg, but we need Assists
            # Fix: Assists were missing in agg above, let's add them if possible or use Kills proxy. 
            # Better: Add 'Assists' to agg.

            # Sortieren nach meisten Matches
            ag_stats = ag_stats.sort_values('MatchesPlayed', ascending=False)
            
            # Tabelle vorbereiten
            table_data = []
            for _, row in ag_stats.iterrows():
                agent_name = row['Agent']
                img_path = get_agent_img(agent_name)
                img_b64 = f"data:image/png;base64,{img_to_b64(img_path)}" if img_path else ""
                
                table_data.append({
                    "Icon": img_b64, "Agent": agent_name,
                    "Matches": int(row['MatchesPlayed']),
                    # NEW FORMULA: KPR + (ADR/450) + (FK-FD)/Rounds
                    "Rating": (row['Kills']/row['Rounds'] if row['Rounds']>0 else 0) + (row['ADR']/450) + ((row['FK']-row['FD'])/row['Rounds'] if row['Rounds']>0 else 0),
                    "Win%": row['Win%'],
                    "K/D": row['KD'], "HS%": row['HS%'],
                    "C": row['C_per_round'], "Q": row['Q_per_round'],
                    "E": row['E_per_round'], "X": row['X_per_round'],
                })
            
            df_table = pd.DataFrame(table_data)
            
                    # HIER WERTE ÄNDERN: (Wert, Gut_Grenze, Schlecht_Grenze)
            styler = df_table.style\
                .format({
                    "Win%": "{:.0f}%", "K/D": "{:.2f}", "HS%": "{:.1f}%", "Rating": "{:.2f}",
                    "C": "{:.1f}", "Q": "{:.1f}", "E": "{:.1f}", "X": "{:.1f}"
                })\
                .map(lambda v: style_good_bad(v, 1.1, 0.9), subset=['Rating'])\
                .map(lambda v: style_good_bad(v, 55, 45), subset=['Win%'])\
                .map(lambda v: style_good_bad(v, 1.1, 0.9), subset=['K/D'])\
                .map(lambda v: style_good_bad(v, 25, 15), subset=['HS%'])\
                .apply(style_purple_gradient, subset=['C', 'Q', 'E', 'X'])


            st.dataframe(
                styler,
                column_config={
                    "Icon": st.column_config.ImageColumn("Icon", width="small"),
                    "Agent": st.column_config.TextColumn("Agent", width="small"),
                    "Matches": st.column_config.NumberColumn("#", format="%d", width="small"),
                    "Rating": st.column_config.NumberColumn("Rating", format="%.2f"),
                    "Win%": st.column_config.TextColumn("Win%"),
                    "K/D": st.column_config.TextColumn("K/D"),
                    "HS%": st.column_config.TextColumn("HS%"),
                    "C": st.column_config.NumberColumn("C", help="Avg Casts per Round"),
                    "Q": st.column_config.NumberColumn("Q", help="Avg Casts per Round"),
                    "E": st.column_config.NumberColumn("E", help="Avg Casts per Round"),
                    "X": st.column_config.NumberColumn("X", help="Avg Casts per Round"),
                },
                use_container_width=True,
                hide_index=True,
                height=400
            )
            
            # EXPORT BUTTON (HTML Report)
            export_df = df_table.drop(columns=['Icon'])
            html_report = f"""
            <html>
            <head><title>{deep_player} Report</title>
            <style>body{{font-family:sans-serif;}} table{{border-collapse:collapse;width:100%;}} th,td{{border:1px solid #444;padding:8px;}} th{{background:#eee;}}</style>
            </head>
            <body>
            <h2>NEXUS Report: {deep_player}</h2>
            <p><strong>K/D:</strong> {df_deep['KD'].mean():.2f} | <strong>HS%:</strong> {df_deep['HS%'].mean():.1f}% | <strong>Rating:</strong> {vlr_approx:.2f}</p>
            {export_df.to_html(index=False)}
            </body></html>
            """
            st.download_button("📄 Export Report (HTML)", html_report, file_name=f"{deep_player}_report.html", mime="text/html", help="Download as HTML, then Print to PDF (Ctrl+P)")

            with st.expander("Ability Key"):
                for agent in df_table['Agent']:
                    abils = AGENT_ABILITIES.get(agent, {})
                    if abils:
                        st.markdown(f"**{agent}:** C: *{abils.get('C')}*, Q: *{abils.get('Q')}*, E: *{abils.get('E')}*, X: *{abils.get('X')}*")

            st.divider()

            # 3. MAP UTILITY ANALYSIS
            st.subheader("🗺️ MAP UTILITY DEEP DIVE")
            
            # Initialize session state for this selector
            if 'deep_agent' not in st.session_state:
                st.session_state.deep_agent = ag_stats['Agent'].unique()[0] if not ag_stats.empty else None

            with st.expander("♟️ SELECT AGENT", expanded=True):
                render_visual_selection(ag_stats['Agent'].unique(), 'agent', 'deep_a_sel', multi=False, key_state='deep_agent')
            
            sel_agent = st.session_state.deep_agent
            
            if sel_agent:
                # Filter auf Agent
                df_ag = df_deep[df_deep['Agent'] == sel_agent]
                # Durchschnittliche Nutzung pro Map berechnen
                map_util = df_ag.groupby('Map')[['Cast_Grenade', 'Cast_Abil1', 'Cast_Abil2', 'Cast_Ult']].mean().reset_index()
                
                # Umformen für Stacked Bar Chart
                map_melt = map_util.melt(id_vars='Map', var_name='Ability', value_name='Avg Casts')
                
                fig2 = px.bar(map_melt, x='Map', y='Avg Casts', color='Ability',
                              title=f"Average Utility Usage for {sel_agent} per Map",
                              color_discrete_map={
                                  'Cast_Ult': '#FF1493', 
                                  'Cast_Grenade': '#88FFFF',
                                  'Cast_Abil1': '#00BFFF',
                                  'Cast_Abil2': '#0055AA'
                              })
                fig2.update_layout(paper_bgcolor='rgba(0,0,0,0)', plot_bgcolor='rgba(255,255,255,0.05)', font=dict(color='white'))
                st.plotly_chart(fig2, use_container_width=True)
            
            st.divider()
            
            # 4. ROLE PERFORMANCE SPIDER CHART
            st.subheader("🕸️ ROLE PERFORMANCE ANALYZER")
            
            # Assign Roles to df_deep
            def get_role(agent_name):
                for role, agents in AGENT_ROLES.items():
                    if agent_name in agents: return role
                return "Flex"
            
            df_deep['Role'] = df_deep['Agent'].apply(get_role)
            
            # Group by Role
            role_stats = df_deep.groupby('Role').agg({
                'Kills': 'sum', 'Deaths': 'sum', 'Assists': 'sum',
                'FK': 'sum', 'FD': 'sum', 'Clutches': 'sum',
                'Rounds': 'sum', 'ADR': 'mean', 'KAST': 'mean'
            }).reset_index()
            
            available_roles = role_stats['Role'].unique()
            
            # Layout Columns
            c_chart, c_stats = st.columns([2, 1])
            
            if len(available_roles) > 0:
                with c_stats:
                    sel_role = st.selectbox("Select Role to Analyze:", available_roles)
                    comp_player = st.selectbox("Compare with:", ["None"] + [p for p in ["Andrei", "Benni", "Luca", "Luggi", "Remus", "Sofi"] if p != deep_player])
                
                r_data = role_stats[role_stats['Role'] == sel_role].iloc[0]
                rounds = r_data['Rounds'] if r_data['Rounds'] > 0 else 1
                
                # Metrics Definitions based on Role
                metrics = {}
                
                if sel_role == "Duelist":
                    metrics = {
                        "FK/FD Ratio": r_data['FK'] / r_data['FD'] if r_data['FD'] > 0 else r_data['FK'],
                        "FB per Round": r_data['FK'] / rounds,
                        "KPR": r_data['Kills'] / rounds,
                        "ADR": r_data['ADR'],
                        "Survival%": (rounds - r_data['Deaths']) / rounds
                    }
                    max_vals = {"FK/FD Ratio": 2.0, "FB per Round": 0.25, "KPR": 1.2, "ADR": 200, "Survival%": 0.6}
                    
                elif sel_role == "Initiator":
                    metrics = {
                        "APR": r_data['Assists'] / rounds,
                        "KAS%": r_data['KAST'], # Fallback KAS calc
                        "Survival%": (rounds - r_data['Deaths']) / rounds,
                        "Low FD/R": 1 - (r_data['FD'] / rounds), # Inverted for chart (Higher is better)
                        "ADR": r_data['ADR']
                    }
                    max_vals = {"APR": 0.6, "KAS%": 100, "Survival%": 0.6, "Low FD/R": 1.0, "ADR": 180}
                    
                elif sel_role == "Controller": # Smokes
                    metrics = {
                        "APR": r_data['Assists'] / rounds,
                        "KAS%": r_data['KAST'],
                        "Survival%": (rounds - r_data['Deaths']) / rounds,
                        "K/D": r_data['Kills'] / r_data['Deaths'] if r_data['Deaths'] > 0 else r_data['Kills'], # Replaced Clutches
                        "ADR": r_data['ADR']
                    }
                    max_vals = {"APR": 0.6, "KAS%": 100, "Survival%": 0.6, "K/D": 1.5, "ADR": 170}
                    
                elif sel_role == "Sentinel":
                    metrics = {
                        "K/D": r_data['Kills'] / r_data['Deaths'] if r_data['Deaths'] > 0 else r_data['Kills'],
                        "Survival%": (rounds - r_data['Deaths']) / rounds,
                        "ADR": r_data['ADR'],
                        "FK/FD Ratio": r_data['FK'] / r_data['FD'] if r_data['FD'] > 0 else r_data['FK'],
                        "KPR": r_data['Kills'] / rounds
                    }
                    max_vals = {"K/D": 1.5, "Survival%": 0.7, "ADR": 180, "FK/FD Ratio": 1.5, "KPR": 1.0}
                else:
                    st.info("Generic stats for Flex role.")
                    metrics = {"KPR": r_data['Kills']/rounds, "ADR": r_data['ADR']}
                    max_vals = {"KPR": 1.0, "ADR": 150}

                # Prepare Chart Data (Primary Player)
                radar_data = [
                    {'Player': deep_player, 'Metric': k, 'Value': min(v / max_vals.get(k, 1), 1.0), 'Display': v} 
                    for k, v in metrics.items()
                ]
                
                # Comparison Logic
                comp_metrics = {}
                if comp_player != "None":
                    comp_path = os.path.join(BASE_DIR, "data", "players", f"{comp_player.lower()}_data.json")
                    if os.path.exists(comp_path):
                        df_comp = parse_tracker_json(comp_path)
                        if not df_comp.empty:
                            df_comp['Role'] = df_comp['Agent'].apply(get_role)
                            comp_role_stats = df_comp.groupby('Role').agg({
                                'Kills': 'sum', 'Deaths': 'sum', 'Assists': 'sum',
                                'FK': 'sum', 'FD': 'sum', 'Clutches': 'sum',
                                'Rounds': 'sum', 'ADR': 'mean', 'KAST': 'mean'
                            }).reset_index()
                            
                            # Filter for the SAME role as primary player
                            cr_data = comp_role_stats[comp_role_stats['Role'] == sel_role]
                            if not cr_data.empty:
                                cr_data = cr_data.iloc[0]
                                c_rounds = cr_data['Rounds'] if cr_data['Rounds'] > 0 else 1
                                
                                # Calculate metrics using same logic
                                if sel_role == "Duelist":
                                    comp_metrics = {"FK/FD Ratio": cr_data['FK']/cr_data['FD'] if cr_data['FD']>0 else cr_data['FK'], "FB per Round": cr_data['FK']/c_rounds, "KPR": cr_data['Kills']/c_rounds, "ADR": cr_data['ADR'], "Survival%": (c_rounds-cr_data['Deaths'])/c_rounds}
                                elif sel_role == "Initiator":
                                    comp_metrics = {"APR": cr_data['Assists']/c_rounds, "KAS%": cr_data['KAST'], "Survival%": (c_rounds-cr_data['Deaths'])/c_rounds, "Low FD/R": 1-(cr_data['FD']/c_rounds), "ADR": cr_data['ADR']}
                                elif sel_role == "Controller":
                                    comp_metrics = {"APR": cr_data['Assists']/c_rounds, "KAS%": cr_data['KAST'], "Survival%": (c_rounds-cr_data['Deaths'])/c_rounds, "K/D": cr_data['Kills']/cr_data['Deaths'] if cr_data['Deaths']>0 else cr_data['Kills'], "ADR": cr_data['ADR']}
                                elif sel_role == "Sentinel":
                                    comp_metrics = {"K/D": cr_data['Kills']/cr_data['Deaths'] if cr_data['Deaths']>0 else cr_data['Kills'], "Survival%": (c_rounds-cr_data['Deaths'])/c_rounds, "ADR": cr_data['ADR'], "FK/FD Ratio": cr_data['FK']/cr_data['FD'] if cr_data['FD']>0 else cr_data['FK'], "KPR": cr_data['Kills']/c_rounds}
                                
                                for k, v in comp_metrics.items():
                                    radar_data.append({'Player': comp_player, 'Metric': k, 'Value': min(v / max_vals.get(k, 1), 1.0), 'Display': v})

                radar_df = pd.DataFrame(radar_data)
                
                with c_chart:
                    fig_spider = px.line_polar(radar_df, r='Value', theta='Metric', color='Player', line_close=True, 
                                              color_discrete_sequence=['#00BFFF', '#FF1493'], # Cyan for P1, Pink for P2
                                              title=f"{sel_role} Performance Profile")
                    fig_spider.update_traces(fill='toself')
                    fig_spider.update_layout(
                        polar=dict(radialaxis=dict(visible=True, showticklabels=False, range=[0, 1]), bgcolor='rgba(255,255,255,0.05)'), 
                        paper_bgcolor='rgba(0,0,0,0)', font=dict(color='white'),
                        legend=dict(orientation="h", y=-0.1)
                    )
                    st.plotly_chart(fig_spider, use_container_width=True)
                
                with c_stats:
                    st.markdown(f"**{deep_player} Stats**")
                    for k, v in metrics.items():
                        delta = None
                        if k in comp_metrics:
                            delta = v - comp_metrics[k]
                        st.metric(k, f"{v:.2f}", delta=f"{delta:.2f}" if delta is not None else None)
        
        else:
            st.info("👆 Bitte lade eine JSON-Datei hoch (Tracker.gg Match-Export oder Profil-Export), um die Analyse zu starten.")

# ==============================================================================
# 📹 VOD REVIEW
# ==============================================================================
elif page == "📹 VOD REVIEW":
    st.title("📹 VOD REVIEW & ANALYSIS")
    
    # Current User
    current_user = st.session_state.get('username', 'Unknown')

    # --- STATE MANAGEMENT FOR VOD WORKSPACE ---
    if 'active_vod_id' not in st.session_state: st.session_state.active_vod_id = None
    
    # --- VIEW 1: LIBRARY (GRID) ---
    if st.session_state.active_vod_id is None:
        c_head, c_act = st.columns([3, 1])
        c_head.markdown("### 📂 VOD LIBRARY")
        
        with c_act:
            if st.button("➕ NEW REVIEW", use_container_width=True, type="primary"):
                st.session_state.active_vod_id = "NEW"
                st.rerun()
            # Repair Button if loading fails
            if df_vods.empty:
                if st.button("🔧 Initialize Database", help="Click this if loading fails or list is empty"):
                    save_vod_reviews(pd.DataFrame(columns=['ID', 'Title', 'Type', 'VideoLink', 'Map', 'Agent', 'Player', 'Notes', 'Tags', 'CreatedBy', 'CreatedAt', 'Rounds']))
                    st.success("Database initialized! Reloading...")
                    time.sleep(1)
                    st.rerun()

        if not df_vods.empty:
            # Filters
            c_s, c_f1, c_f2, c_f3 = st.columns([2, 1, 1, 1])
            search_q = c_s.text_input("🔍 Search", placeholder="Title, Notes, Tags...")
            f_map = c_f1.multiselect("Map", df_vods['Map'].unique())
            f_type = c_f2.multiselect("Type", df_vods['Type'].unique())
            f_player = c_f3.multiselect("Player", df_vods['Player'].unique())
            
            view_df = df_vods.copy()
            if search_q:
                q = search_q.lower()
                view_df = view_df[view_df.apply(lambda x: q in str(x['Title']).lower() or q in str(x['Notes']).lower() or q in str(x['Tags']).lower(), axis=1)]
            if f_map: view_df = view_df[view_df['Map'].isin(f_map)]
            if f_type: view_df = view_df[view_df['Type'].isin(f_type)]
            if f_player: view_df = view_df[view_df['Player'].isin(f_player)]
            
            # Grid Layout
            cols = st.columns(3)
            for idx, row in view_df.sort_values('CreatedAt', ascending=False).iterrows():
                with cols[idx % 3]:
                    with st.container(border=True):
                        # Thumbnail / Map Image
                        map_img = get_map_img(row['Map'], 'list')
                        if map_img: st.image(map_img, use_container_width=True)
                        
                        st.markdown(f"**{row['Title']}**")
                        st.caption(f"{row['Type']} • {row['Map']} • {row['CreatedAt'][:10]}")
                        if row['Tags']: st.caption(f"🏷️ {row['Tags']}")
                        
                        if st.button("▶️ OPEN REVIEW", key=f"open_{row['ID']}", use_container_width=True):
                            st.session_state.active_vod_id = row['ID']
                            st.rerun()
        else:
            st.info("No reviews found. Click 'NEW REVIEW' to start.")

    # --- VIEW 2: WORKSPACE (VALOLENS STYLE) ---
    else:
        # Load Active Review Data
        is_new = st.session_state.active_vod_id == "NEW"
        
        # FIX: Check if ID exists to prevent crash
        if not is_new and st.session_state.active_vod_id not in df_vods['ID'].values:
            st.warning("Review not found or deleted. Returning to library.")
            st.session_state.active_vod_id = None
            st.rerun()

        if is_new:
            row = {'ID': 'NEW', 'Title': '', 'Type': 'Own Gameplay', 'VideoLink': '', 'Map': 'Ascent', 'Player': '', 'Agent': '', 'Notes': '', 'Tags': '', 'Rounds': '[]'}
        else:
            row = df_vods[df_vods['ID'] == st.session_state.active_vod_id].iloc[0]

        # Header
        c_back, c_title, c_save = st.columns([1, 4, 1])
        if c_back.button("⬅ BACK"):
            st.session_state.active_vod_id = None
            st.rerun()
        
        with c_title:
            if is_new: st.markdown("### 🆕 Creating New Review")
            else: st.markdown(f"### 📹 {row['Title']}")

        # --- LAYOUT ---
        col_video, col_tools = st.columns([2, 1])

        # --- LEFT: VIDEO PLAYER ---
        with col_video:
            # Video Link Input (Always visible for editing)
            v_link = st.text_input("Video Link", value=row['VideoLink'], key="wk_vlink", placeholder="Paste YouTube/Twitch link...")
            
            # Video Player
            # Check if we need to jump to a timestamp (from Round or Timestamp click)
            start_time = 0
            if f"seek_{row['ID']}" in st.session_state:
                start_time = st.session_state[f"seek_{row['ID']}"]
                del st.session_state[f"seek_{row['ID']}"] # Consume event
            
            if v_link:
                # If it's a round jump, we might want to auto-play. 
                # Streamlit video updates when start_time changes.
                st.video(v_link, start_time=start_time)
            else:
                st.info("📺 Video will appear here.")
            
            # Metadata Inputs
            with st.expander("📝 Metadata & Settings", expanded=is_new):
                c1, c2 = st.columns(2)
                new_title = c1.text_input("Title", value=row['Title'], key="wk_title")
                new_map = c2.selectbox("Map", sorted(df['Map'].unique()), index=sorted(df['Map'].unique()).index(row['Map']) if row['Map'] in df['Map'].unique() else 0, key="wk_map")
                c3, c4 = st.columns(2)
                new_type = c3.selectbox("Type", ["Own Gameplay", "Pro Match", "Scrim Review"], index=["Own Gameplay", "Pro Match", "Scrim Review"].index(row['Type']) if row['Type'] in ["Own Gameplay", "Pro Match", "Scrim Review"] else 0, key="wk_type")
                new_tags = c4.multiselect("Tags", ["Macro", "Micro", "Comms", "Utility", "Clutch", "IGL"], default=row['Tags'].split(", ") if row['Tags'] else [], key="wk_tags")

        # --- RIGHT: TOOLS & LOGGING ---
        with col_tools:
            # --- MODE SWITCHER ---
            mode = st.radio("Mode", ["📝 General Notes", "🔢 Round by Round"], horizontal=True, label_visibility="collapsed")
            
            # Initialize Rounds Data
            try: rounds_data = json.loads(row['Rounds']) if pd.notna(row.get('Rounds')) and row['Rounds'] else []
            except: rounds_data = []

            # 2. NOTES EDITOR
            st.markdown("#### 📝 Analysis")
            st.divider()
            
            # Initialize temp notes in session state if not present
            if 'wk_notes_temp' not in st.session_state:
                st.session_state.wk_notes_temp = row['Notes']

            if mode == "📝 General Notes":
                if 'wk_notes_temp' not in st.session_state: st.session_state.wk_notes_temp = row['Notes']
                
                notes_val = st.text_area("General Notes", value=st.session_state.wk_notes_temp, height=400, key="wk_notes_area")
                st.session_state.wk_notes_temp = notes_val
            
            else: # ROUND BY ROUND
                notes_val = st.session_state.wk_notes_temp
                if 'wk_rounds_temp' not in st.session_state: st.session_state.wk_rounds_temp = rounds_data
                curr_rounds = st.session_state.wk_rounds_temp
                
                # Round Navigation
                if not curr_rounds:
                    if st.button("➕ Add Round 1"):
                        curr_rounds.append({'Round': 1, 'Time': '00:00', 'Result': '?', 'Notes': ''})
                        st.session_state.wk_rounds_temp = curr_rounds
                        st.rerun()
                    st.info("No rounds yet.")
                else:
                    if 'curr_round_idx' not in st.session_state: st.session_state.curr_round_idx = 0
                    idx = st.session_state.curr_round_idx
                    
                    # Nav Bar
                    c_prev, c_sel, c_next, c_add = st.columns([1, 3, 1, 1])
                    if c_prev.button("◀", disabled=(idx==0)): st.session_state.curr_round_idx -= 1; st.rerun()
                    if c_next.button("▶", disabled=(idx==len(curr_rounds)-1)): st.session_state.curr_round_idx += 1; st.rerun()
                    if c_add.button("➕"): 
                        new_r = len(curr_rounds) + 1
                        curr_rounds.append({'Round': new_r, 'Time': '00:00', 'Result': '?', 'Notes': ''})
                        st.session_state.curr_round_idx = len(curr_rounds) - 1
                        st.rerun()
                    
                    with c_sel:
                        # Dropdown to jump
                        r_opts = [f"R{r['Round']} ({r['Time']})" for r in curr_rounds]
                        sel_r = st.selectbox("Select Round", r_opts, index=idx, label_visibility="collapsed", key="rnd_sel")
                        # Sync dropdown change
                        new_idx = r_opts.index(sel_r)
                        if new_idx != idx:
                            st.session_state.curr_round_idx = new_idx
                            # Jump Video
                            ts_str = curr_rounds[new_idx]['Time']
                            parts = list(map(int, ts_str.split(':')))
                            sec = parts[0]*60 + parts[1]
                            st.session_state[f"seek_{row['ID']}"] = sec
                            st.rerun()

                    # Current Round Editor
                    cur_r = curr_rounds[st.session_state.curr_round_idx]
                    
                    c_meta1, c_meta2 = st.columns(2)
                    cur_r['Time'] = c_meta1.text_input("Timestamp", cur_r['Time'])
                    cur_r['Result'] = c_meta2.selectbox("Result", ["Win", "Loss", "?"], index=["Win", "Loss", "?"].index(cur_r.get('Result', '?')))
                    
                    cur_r['Notes'] = st.text_area(f"Notes for Round {cur_r['Round']}", value=cur_r['Notes'], height=300, key=f"r_note_{st.session_state.curr_round_idx}")
                    
                    # Update List
                    curr_rounds[st.session_state.curr_round_idx] = cur_r
                    st.session_state.wk_rounds_temp = curr_rounds

            # 3. SAVE BUTTON
            if c_save.button("💾 SAVE & CLOSE", type="primary", use_container_width=True):
                if not new_title:
                    st.error("Title required")
                else:
                    # Construct Data
                    save_data = {
                        'ID': row['ID'] if not is_new else str(uuid.uuid4())[:8],
                        'Title': new_title,
                        'Type': new_type,
                        'VideoLink': v_link,
                        'Map': new_map,
                        'Agent': row['Agent'], # Keep existing or add selector
                        'Player': row['Player'],
                        'Notes': notes_val,
                        'Tags': ", ".join(new_tags),
                        'CreatedBy': row['CreatedBy'] if not is_new else current_user,
                        'CreatedAt': row['CreatedAt'] if not is_new else datetime.now().strftime("%Y-%m-%d %H:%M:%S")
                    }
                    
                    if is_new:
                        db_insert("nexus_vod_reviews", save_data, "df_vods")
                    else:
                        db_update("nexus_vod_reviews", save_data, "df_vods")
                    
                    st.success("Saved!")
                    if is_new:
                        st.session_state.active_vod_id = None # Return to library
                        st.rerun()

            # 5. TELESTRATOR
            with st.expander("🎨 Telestrator"):
                 # ... (Existing Telestrator Code reused) ...
                 # Paste Button
                            if HAS_CLIPBOARD:
                                p_ts = paste(label="📋 Paste Screenshot", key=f"paste_ts_wk")
                                if p_ts:
                                    try:
                                        if isinstance(p_ts, str) and "," in p_ts: p_ts = p_ts.split(",")[1]
                                        ib = base64.b64decode(p_ts)
                                        img = Image.open(io.BytesIO(ib))
                                        # Resize
                                        base_width = 600
                                        w_percent = (base_width / float(img.size[0]))
                                        h_size = int((float(img.size[1]) * float(w_percent)))
                                        img = img.resize((base_width, h_size), Image.Resampling.LANCZOS)
                                        st.session_state['ts_bg_wk'] = img
                                        st.session_state['ts_bg_id'] = st.session_state.get('ts_bg_id', 0) + 1
                                        st.rerun()
                                    except Exception as e: st.error(f"Error: {e}")
                            
                            ts_key_bg = 'ts_bg_wk'
                            if ts_key_bg not in st.session_state: st.session_state[ts_key_bg] = None
                            
                            if st.session_state[ts_key_bg]:
                                # Toolbar
                                c_t1, c_t2, c_t3 = st.columns([1,1,1])
                                mode = c_t1.selectbox("Tool", ["freedraw", "line", "rect", "circle"], key=f"tool_wk")
                                color = c_t2.color_picker("Color", "#00BFFF", key=f"col_wk")
                                stroke = c_t3.slider("Width", 1, 10, 3, key=f"wid_wk")
                                
                                canvas_result = st_canvas(
                                    fill_color="rgba(255, 165, 0, 0.1)", stroke_width=stroke, stroke_color=color,
                                    background_image=st.session_state[ts_key_bg], update_streamlit=True,
                                    height=st.session_state[ts_key_bg].height, width=st.session_state[ts_key_bg].width,
                                    drawing_mode=mode, key=f"canvas_wk_{st.session_state.get('ts_bg_id', 0)}",
                                )
                                
                                if st.button("💾 Save Analysis to Notes", key=f"save_ts_wk"):
                                    if canvas_result.image_data is not None:
                                        bg = st.session_state[ts_key_bg].convert("RGBA")
                                        fg = Image.fromarray(canvas_result.image_data.astype('uint8'), "RGBA")
                                        if fg.size != bg.size: fg = fg.resize(bg.size)
                                        final = Image.alpha_composite(bg, fg)
                                        fn = f"VOD_ANALY_{uuid.uuid4().hex[:8]}.png"
                                        final.save(os.path.join(VOD_IMG_DIR, fn))
                                        
                                        st.session_state.wk_notes_temp += f"\n\n**Analysis:**\n[[img:{fn}]]\n"
                                        st.success("Saved to notes!"); st.rerun()

            # 6. EXPORT TO LINEUPS
            st.divider()
            with st.expander("🎯 Export to Lineups"):
                st.caption("Create a lineup card directly from this VOD.")
                
                l_exp_title = st.text_input("Lineup Title", placeholder="e.g. Sova Dart for B Main", key="lu_exp_title")
                
                c_lu1, c_lu2 = st.columns(2)
                
                # Agent from VOD
                vod_agent = row['Agent'] if row['Agent'] else "Sova"
                
                # Abilities
                avail_abils = AGENT_ABILITIES.get(vod_agent, {})
                abil_opts = [f"{k}: {v}" for k, v in avail_abils.items()] if avail_abils else ["Default"]
                l_exp_ability = c_lu1.selectbox("Ability", abil_opts, key="lu_exp_abil")
                
                l_exp_side = c_lu2.selectbox("Side", ["Attack", "Defense"], key="lu_exp_side")
                
                # Tags
                l_exp_tags = st.multiselect("Tags / Type", ["Recon", "Shock", "Molly", "Flash", "Smoke", "Wall", "Ult", "One-Way", "God Spot", "Retake"], key="lu_exp_tags")
                
                if st.button("📤 Export to Library", key="btn_export_lu", use_container_width=True):
                    if not l_exp_title:
                        st.error("Title required.")
                    else:
                        # Clean Ability Name
                        final_type = l_exp_ability.split(": ")[1] if ": " in l_exp_ability else l_exp_ability
                        final_tags = ", ".join(l_exp_tags)
                        
                        new_lu = {
                            'ID': str(uuid.uuid4())[:8],
                            'Map': row['Map'], 'Agent': vod_agent, 
                            'Side': l_exp_side, 'Type': final_type,
                            'Title': l_exp_title, 'Image': "", 'VideoLink': row['VideoLink'],
                            'Description': f"Exported from VOD: {row['Title']}", 'Tags': final_tags, 
                            'CreatedBy': current_user
                        }
                        db_insert("nexus_lineups", new_lu, "df_lineups")
                        st.success(f"Exported to Lineup Library!")

        # Delete Button (Bottom)
        if not is_new:
            if st.button("🗑️ Delete Review", key="del_wk"):
                db_delete("nexus_vod_reviews", row['ID'], "df_vods")
                st.session_state.active_vod_id = None
                st.rerun()

# ==============================================================================
# 8. DATABASE
# ==============================================================================
elif page == "💾 DATABASE":
    st.header("Database")
    ed = st.data_editor(df, num_rows="dynamic")
    if st.button("Save"): 
        save_matches(ed)
        st.success("Saved to Google Sheets")
                        if final_img_name or final_vid_link:
                            new_lu = {
                                'ID': str(uuid.uuid4())[:8],
                                'Map': row['Map'], 'Agent': row['Agent'] if row['Agent'] else "Sova", 
                                'Side': l_exp_side, 'Type': l_exp_type,
                                'Title': l_exp_title, 'Image': final_img_name, 'VideoLink': final_vid_link,
                                'Description': f"Exported from VOD: {row['Title']}", 'Tags': "VOD Export", 
                                'CreatedBy': current_user
                            }
                            db_insert("nexus_lineups", new_lu, "df_lineups")
                            st.success(f"Exported to Lineup Library!")
                        new_lu = {
                            'ID': str(uuid.uuid4())[:8],
                            'Map': row['Map'], 'Agent': vod_agent, 
                            'Side': l_exp_side, 'Type': final_type,
                            'Title': l_exp_title, 'Image': "", 'VideoLink': row['VideoLink'],
                            'Description': f"Exported from VOD: {row['Title']}", 'Tags': final_tags, 
                            'CreatedBy': current_user
                        }
                        db_insert("nexus_lineups", new_lu, "df_lineups")
                        st.success(f"Exported to Lineup Library!")

        # Delete Button (Bottom)
        if not is_new:
            if st.button("🗑️ Delete Review", key="del_wk"):
                db_delete("nexus_vod_reviews", row['ID'], "df_vods")
                st.session_state.active_vod_id = None
                st.rerun()

# ==============================================================================
# 8. DATABASE
# ==============================================================================
elif page == "💾 DATABASE":
    st.header("Database")
    ed = st.data_editor(df, num_rows="dynamic")
    if st.button("Save"): 
        save_matches(ed)
        st.success("Saved to Google Sheets")